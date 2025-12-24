import streamlit as st
import pandas as pd
import json
import os
from datetime import datetime
import uuid
import hashlib
import io

# 보고서 생성을 위한 라이브러리 체크 (GitHub 배포 시 requirements.txt 필용)
try:
    from docx import Document
    from docx.shared import Pt
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptPt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# ==========================================
# 1. 데이터 모델 및 관리 로직
# ==========================================

class AdvancedCustomerManager:
    """
    고객 데이터 관리 클래스
    - 로컬 모드: data/customers/ 폴더에 JSON 저장
    - 확장성: _save_to_file 내부를 API 호출로 바꾸면 즉시 클라우드 DB 연동 가능
    """
    def __init__(self, storage_path="data/customers"):
        # Streamlit Cloud/GitHub 배포 시 안정성을 위해 절대 경로로 변환
        if not os.path.isabs(storage_path):
            # 현재 작업 디렉토리 기준으로 상대 경로를 절대 경로로 변환
            # Streamlit Cloud에서는 작업 디렉토리가 앱 루트이므로 이 방식이 안전함
            base_dir = os.getcwd()
            self.storage_path = os.path.join(base_dir, storage_path)
        else:
            self.storage_path = storage_path
        
        # 디렉토리 생성 (부모 디렉토리까지 자동 생성)
        os.makedirs(self.storage_path, exist_ok=True)

    def generate_identity_hash(self, phone, email):
        """연락처와 이메일을 조합해 동일 고객 여부 판별용 고유값 생성"""
        raw_str = f"{phone.strip()}|{email.strip().lower()}"
        return hashlib.md5(raw_str.encode()).hexdigest()

    def create_customer(self, name, phone, email, trait="일반"):
        """신규 고객 마스터 생성"""
        identity_hash = self.generate_identity_hash(phone, email)
        customer_id = f"CUST-{uuid.uuid4().hex[:8].upper()}"
        
        customer_data = {
            "basic_info": {
                "customer_id": customer_id,
                "name": name,
                "phone": phone,
                "email": email,
                "identity_hash": identity_hash,
                "created_at": datetime.now().isoformat(),
                "last_access_at": datetime.now().isoformat(),
            },
            "crm_profile": {
                "trait": trait,
                "total_consultations": 0,
                "avg_satisfaction": 0.0,
                "tags": []
            },
            "consultation_history": []
        }
        self._save_to_file(customer_id, customer_data)
        return customer_id

    def add_consultation(self, customer_id, content, summary, sentiment, evaluation):
        """새로운 상담 이력 추가 및 마스터 정보 갱신"""
        data = self.load_customer(customer_id)
        if not data:
            return False

        consult_id = f"CON-{uuid.uuid4().hex[:6].upper()}"
        new_consult = {
            "consult_id": consult_id,
            "date": datetime.now().isoformat(),
            "content": content,
            "summary": summary,
            "analysis": {
                "sentiment": sentiment,
                "customer_emotion": evaluation.get("emotion", "Normal")
            },
            "survey": {
                "score": evaluation.get("score", 5),
                "feedback": evaluation.get("feedback", "")
            }
        }

        data["consultation_history"].append(new_consult)
        data["basic_info"]["last_access_at"] = datetime.now().isoformat()
        data["crm_profile"]["total_consultations"] = len(data["consultation_history"])
        
        # 만족도 평균 계산
        scores = [c["survey"]["score"] for c in data["consultation_history"]]
        data["crm_profile"]["avg_satisfaction"] = sum(scores) / len(scores)

        self._save_to_file(customer_id, data)
        return True

    def load_customer(self, customer_id):
        path = os.path.join(self.storage_path, f"{customer_id}.json")
        if os.path.exists(path):
            try:
                with open(path, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except (json.JSONDecodeError, IOError, OSError) as e:
                # 에러 로깅 (필요시)
                return None
        return None

    def _save_to_file(self, customer_id, data):
        """데이터 저장부 (추후 클라우드 API 연동 시 이 부분만 수정)"""
        path = os.path.join(self.storage_path, f"{customer_id}.json")
        try:
            # 디렉토리가 없으면 생성
            os.makedirs(os.path.dirname(path), exist_ok=True)
            with open(path, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=4)
        except (IOError, OSError) as e:
            raise Exception(f"파일 저장 실패: {str(e)}")

    def list_all_customers(self):
        """모든 고객 리스트 반환"""
        if not os.path.exists(self.storage_path):
            return []
        try:
            files = [f for f in os.listdir(self.storage_path) if f.endswith('.json')]
        except (OSError, PermissionError):
            return []
        
        all_data = []
        for f in files:
            try:
                with open(os.path.join(self.storage_path, f), 'r', encoding='utf-8') as file:
                    all_data.append(json.load(file))
            except (json.JSONDecodeError, IOError, OSError):
                # 손상된 파일은 건너뛰기
                continue
        return all_data

    def find_customer_by_info(self, name=None, phone=None, email=None):
        """
        고객 정보(이름, 전화번호, 이메일)로 이전 응대 이력이 있는 고객을 검색
        최소 1~2개 정보가 일치하면 반환
        """
        all_customers = self.list_all_customers()
        if not all_customers:
            return None
        
        # 검색 조건 정규화
        search_name = name.strip().lower() if name else None
        search_phone = phone.strip().replace("-", "").replace(" ", "") if phone else None
        search_email = email.strip().lower() if email else None
        
        if not any([search_name, search_phone, search_email]):
            return None
        
        # 매칭 점수 계산
        best_match = None
        best_score = 0
        
        for customer in all_customers:
            basic_info = customer.get("basic_info", {})
            customer_name = basic_info.get("name", "").strip().lower()
            customer_phone = basic_info.get("phone", "").strip().replace("-", "").replace(" ", "")
            customer_email = basic_info.get("email", "").strip().lower()
            
            match_score = 0
            match_count = 0
            
            # 이름 매칭
            if search_name and customer_name:
                if search_name == customer_name:
                    match_score += 2
                    match_count += 1
            
            # 전화번호 매칭
            if search_phone and customer_phone:
                if search_phone == customer_phone:
                    match_score += 3  # 전화번호는 더 높은 가중치
                    match_count += 1
            
            # 이메일 매칭
            if search_email and customer_email:
                if search_email == customer_email:
                    match_score += 3  # 이메일도 더 높은 가중치
                    match_count += 1
            
            # 최소 1개 이상 일치하고, 점수가 더 높으면 업데이트
            if match_count >= 1 and match_score > best_score:
                best_score = match_score
                best_match = customer
        
        return best_match if best_match else None

    def load_customer_data(self, customer_id):
        """load_customer의 별칭 (하위 호환성)"""
        return self.load_customer(customer_id)

    # --- 보고서 생성 로직 ---

    def generate_word_report(self, customer_data):
        """Google Docs 호환 Word 리포트 생성"""
        if not HAS_DOCX:
            raise ImportError("python-docx가 설치되지 않았습니다. requirements.txt에서 설치하세요.")
        doc = Document()
        doc.add_heading(f"상담 기록 보고서: {customer_data['basic_info']['name']} 고객님", 0)
        
        doc.add_heading("1. 고객 마스터 정보", level=1)
        p = doc.add_paragraph()
        p.add_run(f"고객 ID: {customer_data['basic_info']['customer_id']}\n")
        p.add_run(f"연락처: {customer_data['basic_info']['phone']}\n")
        p.add_run(f"이메일: {customer_data['basic_info']['email']}\n")
        p.add_run(f"고객성향: {customer_data['crm_profile']['trait']}")

        doc.add_heading("2. 상담 이력 히스토리", level=1)
        for idx, con in enumerate(reversed(customer_data['consultation_history'])):
            doc.add_heading(f"상담 건 #{len(customer_data['consultation_history'])-idx} ({con['date'][:10]})", level=2)
            doc.add_paragraph(f"요약: {con['summary']}")
            doc.add_paragraph(f"상세내용:\n{con['content']}")
            doc.add_paragraph(f"평가 결과: {con['survey']['score']}점 / 피드백: {con['survey']['feedback']}")
            
        bio = io.BytesIO()
        doc.save(bio)
        return bio.getvalue()

    def generate_pptx_profile(self, customer_data):
        """Google Slides 호환 PPTX 프로필 생성"""
        if not HAS_PPTX:
            raise ImportError("python-pptx가 설치되지 않았습니다. requirements.txt에서 설치하세요.")
        prs = Presentation()
        
        # 슬라이드 1: 메인 타이틀
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"Customer Profile Report"
        slide.placeholders[1].text = f"대상: {customer_data['basic_info']['name']}\n보고서 생성일: {datetime.now().strftime('%Y-%m-%d')}"
        
        # 슬라이드 2: CRM 분석 요약
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "고객 성향 및 응대 분석"
        tf = slide.placeholders[1].text_frame
        tf.text = f"현재 고객 성향: {customer_data['crm_profile']['trait']}"
        tf.add_paragraph().text = f"누적 상담 건수: {customer_data['crm_profile']['total_consultations']}회"
        tf.add_paragraph().text = f"평균 고객 만족도: {customer_data['crm_profile']['avg_satisfaction']:.1f} / 5.0"
        
        bio = io.BytesIO()
        prs.save(bio)
        return bio.getvalue()

# ==========================================
# 2. Streamlit UI (모듈형)
# ==========================================

def render_customer_management_ui():
    manager = AdvancedCustomerManager()

    st.sidebar.divider()
    st.sidebar.subheader("CRM 관리 도구")
    menu = st.sidebar.radio("작업 선택", ["현황 대시보드", "상담 기록 입력", "고객 데이터 조회"], key="crm_menu")

    if menu == "현황 대시보드":
        st.subheader("📊 전체 데이터 요약")
        customers = manager.list_all_customers()
        
        if customers:
            c1, c2, c3 = st.columns(3)
            c1.metric("총 등록 고객", f"{len(customers)} 명")
            c2.metric("누적 상담", f"{sum([c['crm_profile']['total_consultations'] for c in customers])} 건")
            avg_score = sum([c['crm_profile']['avg_satisfaction'] for c in customers])/len(customers)
            c3.metric("평균 만족도", f"{avg_score:.2f} / 5.0")
            
            st.divider()
            traits = [c["crm_profile"]["trait"] for c in customers]
            st.subheader("고객 성향 분포")
            st.bar_chart(pd.Series(traits).value_counts())
        else:
            st.info("아직 저장된 데이터가 없습니다. 먼저 상담을 기록해 보세요.")

    elif menu == "상담 기록 입력":
        st.subheader("📝 상담 내용 저장")
        st.caption("고객과의 상담을 마친 후 즉시 기록을 남기세요. (하루 10~20개 권장)")
        
        with st.form("new_consult_form", clear_on_submit=True):
            col1, col2 = st.columns(2)
            with col1:
                name = st.text_input("고객 이름")
                phone = st.text_input("연락처 (동일인 식별 키)")
                email = st.text_input("이메일 주소")
            with col2:
                trait = st.selectbox("고객 성향 분류", ["일반", "부드러움", "이성적", "급함", "까다로움", "강성/진상"])
                sentiment = st.select_slider("상담 감정 분위기", options=["매우나쁨", "나쁨", "중립", "좋음", "매우좋음"], value="중립")
            
            st.divider()
            content = st.text_area("상담 상세 내용 (상황 및 처리 결과)")
            summary = st.text_input("상담 핵심 요약 (한 줄)")
            
            col_s1, col_s2 = st.columns(2)
            with col_s1:
                score = st.slider("응대 평가 점수 (1~5)", 1, 5, 5)
            with col_s2:
                feedback = st.text_input("고객 주관식 피드백/메모")

            submitted = st.form_submit_button("상담 데이터 세이프티 저장")
            
            if submitted:
                if not name or not phone:
                    st.warning("이름과 연락처는 필수 입력값입니다.")
                else:
                    # 고객 식별
                    all_custs = manager.list_all_customers()
                    target_id = None
                    new_hash = manager.generate_identity_hash(phone, email)
                    for c in all_custs:
                        if c["basic_info"]["identity_hash"] == new_hash:
                            target_id = c["basic_info"]["customer_id"]
                            break
                    
                    if not target_id:
                        target_id = manager.create_customer(name, phone, email, trait)
                        st.toast(f"신규 고객(ID: {target_id})으로 등록되었습니다.")
                    
                    # 상담 추가
                    eval_data = {"score": score, "feedback": feedback, "emotion": sentiment}
                    if manager.add_consultation(target_id, content, summary, sentiment, eval_data):
                        st.success(f"데이터가 안전하게 저장되었습니다. (고객: {name})")
                        # ⭐ 데이터 저장 후 즉시 화면 갱신 (대시보드/조회 탭에 반영)
                        st.session_state.last_saved_customer_id = target_id
                        # st.rerun()  # 주석 처리: 폼 제출 후 Streamlit이 자동 rerun함 (필요시 주석 해제)
                    else:
                        st.error("데이터 저장 중 오류가 발생했습니다.")

    elif menu == "고객 데이터 조회":
        st.subheader("🔍 개별 고객 정밀 조회")
        customers = manager.list_all_customers()
        
        if customers:
            # 검색 UI
            cust_labels = {f"{c['basic_info']['name']} ({c['basic_info']['phone']})": c['basic_info']['customer_id'] for c in customers}
            
            # ⭐ 마지막 저장된 고객 자동 선택
            default_index = 0
            if st.session_state.get("last_saved_customer_id"):
                last_saved_id = st.session_state.last_saved_customer_id
                for idx, (label, cust_id) in enumerate(cust_labels.items()):
                    if cust_id == last_saved_id:
                        default_index = idx
                        break
            
            selected_label = st.selectbox("조회 대상을 선택하세요", list(cust_labels.keys()), index=default_index)
            
            if selected_label:
                target_id = cust_labels[selected_label]
                cust = manager.load_customer(target_id)
                
                # 상단 프로필 요약
                with st.container(border=True):
                    c1, c2, c3 = st.columns(3)
                    c1.write(f"**이름:** {cust['basic_info']['name']}")
                    c1.write(f"**연락처:** {cust['basic_info']['phone']}")
                    c2.write(f"**성향:** `{cust['crm_profile']['trait']}`")
                    c2.write(f"**상담:** {cust['crm_profile']['total_consultations']}회")
                    c3.write(f"**만족도:** {cust['crm_profile']['avg_satisfaction']:.1f} / 5.0")
                
                # 보고서 내보내기 영역
                st.divider()
                st.subheader("📤 외부 연동용 리포트 생성")
                exp_col1, exp_col2 = st.columns(2)
                
                with exp_col1:
                    if HAS_DOCX:
                        word_bin = manager.generate_word_report(cust)
                        st.download_button(
                            "📄 구글 Docs용 보고서 (Word)",
                            data=word_bin,
                            file_name=f"Report_{cust['basic_info']['name']}.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            use_container_width=True
                        )
                    else:
                        st.info("Word 생성을 위해 python-docx가 필요합니다.")

                with exp_col2:
                    if HAS_PPTX:
                        ppt_bin = manager.generate_pptx_profile(cust)
                        st.download_button(
                            "📊 구글 Slides용 프로필 (PPT)",
                            data=ppt_bin,
                            file_name=f"Profile_{cust['basic_info']['name']}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True
                        )
                    else:
                        st.info("PPT 생성을 위해 python-pptx가 필요합니다.")

                # 하단 타임라인
                st.divider()
                st.write("### 📅 상담 타임라인 (최신순)")
                for item in reversed(cust["consultation_history"]):
                    with st.expander(f"📍 {item['date'][:16]} - {item['summary']}"):
                        st.write(f"**내용:** {item['content']}")
                        st.caption(f"감정: {item['analysis']['sentiment']} | 점수: {item['survey']['score']}점")
        else:
            st.info("조회할 데이터가 없습니다.")

if __name__ == "__main__":
    st.set_page_config(page_title="고객 데이터 관리 시스템", layout="wide")
    render_customer_management_ui()
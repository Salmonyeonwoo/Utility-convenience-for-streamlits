"""
파일 파싱 모듈 (기본 파싱 함수)
PDF, Word, PPTX 파일에서 텍스트 추출 및 데이터 파싱
"""
import re

# 파일 파싱을 위한 라이브러리 (선택적 import)
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def extract_data_from_text(text):
    """텍스트에서 고객 정보 추출 (정규식 기반)"""
    extracted = []
    
    # 고객명 추출
    name_pattern = r'(?:고객명|이름|성함|Name)[\s:：]*([가-힣a-zA-Z\s]+)'
    name_match = re.search(name_pattern, text, re.IGNORECASE)
    name = name_match.group(1).strip() if name_match else ""
    
    # 연락처 추출
    phone_pattern = r'(?:연락처|전화|Phone|Tel)[\s:：]*([0-9\-+\s]+)'
    phone_match = re.search(phone_pattern, text, re.IGNORECASE)
    phone = phone_match.group(1).strip() if phone_match else ""
    
    # 이메일 추출
    email_pattern = r'([a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,})'
    email_match = re.search(email_pattern, text)
    email = email_match.group(1) if email_match else ""
    
    # 상담 유형 추출
    consult_types = ["배송문의/지연", "환불/반품요청", "결제/오류문의", 
                    "상품정보/재고", "계정/로그인", "강성/컴플레인", "기타"]
    consult_type = "기타"
    for ct in consult_types:
        if ct in text:
            consult_type = ct
            break
    
    # 상태 추출 (Solved/Pending)
    status = "Pending"
    if "Solved" in text or "해결" in text or "완료" in text:
        status = "Solved"
    
    # CSAT 점수 추출
    score = 5
    score_pattern = r'(?:CSAT|만족도|점수|Score)[\s:：]*([1-5])'
    score_match = re.search(score_pattern, text, re.IGNORECASE)
    if score_match:
        score = int(score_match.group(1))
    
    # 감정 분석 추출
    sentiment = "보통"
    if "매우나쁨" in text or "매우 나쁨" in text:
        sentiment = "매우나쁨"
    elif "나쁨" in text:
        sentiment = "나쁨"
    elif "좋음" in text or "좋은" in text:
        sentiment = "좋음"
    elif "매우좋음" in text or "매우 좋음" in text:
        sentiment = "매우좋음"
    
    # 고객 성향 추출
    trait = "일반"
    traits = ["부드러움", "합리적", "급함", "까다로움", "진상/강성", "강성"]
    for t in traits:
        if t in text:
            trait = t
            break
    
    if name or phone:  # 최소한 이름이나 연락처가 있어야 유효한 데이터
        extracted.append({
            "name": name,
            "phone": phone,
            "email": email,
            "trait": trait,
            "consult_type": consult_type,
            "status": status,
            "content": text[:500],  # 처음 500자만
            "summary": text[:100] if len(text) > 100 else text,
            "analysis": {
                "sentiment": sentiment,
                "score": score
            }
        })
    
    return extracted


def parse_pdf(file_path):
    """PDF 파일에서 데이터 추출"""
    if not PDF_AVAILABLE:
        return None
    try:
        import streamlit as st
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return extract_data_from_text(text)
    except Exception as e:
        import streamlit as st
        st.error(f"PDF 파싱 오류: {str(e)}")
        return None


def parse_docx(file_path):
    """Word 파일에서 데이터 추출"""
    if not DOCX_AVAILABLE:
        return None
    try:
        import streamlit as st
        doc = DocxDocument(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return extract_data_from_text(text)
    except Exception as e:
        import streamlit as st
        st.error(f"Word 파싱 오류: {str(e)}")
        return None


def parse_pptx(file_path):
    """PPTX 파일에서 데이터 추출"""
    if not PPTX_AVAILABLE:
        return None
    try:
        import streamlit as st
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return extract_data_from_text(text)
    except Exception as e:
        import streamlit as st
        st.error(f"PPTX 파싱 오류: {str(e)}")
        return None

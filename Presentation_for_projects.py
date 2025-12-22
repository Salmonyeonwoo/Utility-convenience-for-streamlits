from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_ai_simulator_deck():
    prs = Presentation()

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI 고객 응대 시뮬레이터 개발 프로젝트"
    subtitle.text = "OpenAI & Gemini API 기반의 지능형 CS 훈련 시스템\n\n발표자: [성함] | 2024.05.31"

    # Slide 2: Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "프로젝트 개요 및 핵심 성과"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "🎯 목표: 실제 상담 환경과 유사한 고비용 효율의 훈련 시스템 구축"
    p = tf.add_paragraph()
    p.text = "🏆 핵심 성과:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "1. Video RAG: API 한계를 극복한 영상 동기화"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. 기술 통합: STT/TTS 및 3개국어 실시간 통역"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. AI 협업: Cursor AI 활용으로 개발 속도 500% 향상"
    p.level = 1

    # Slide 3: Tech Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "기술 스택 (Tech Stack)"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "🧠 Core Brain: OpenAI GPT-4o / Google Gemini Pro"
    p = tf.add_paragraph()
    p.text = "🗣️ Voice & Interface: OpenAI Whisper (STT) / TTS / Streamlit"
    p = tf.add_paragraph()
    p.text = "💾 Data & Logic: FAISS (Vector DB) / Video Clip DB"

    # Slide 4: AI Collaboration
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "개발 프로세스 혁신 (with Cursor AI)"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "🤖 AI Pair Programming: 전체 코드의 80% 이상 AI 협업 작성"
    p = tf.add_paragraph()
    p.text = "Success Stories:"
    p = tf.add_paragraph()
    p.text = "✅ Video RAG 로직 구현: 감정-영상 매핑 알고리즘"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ 버그 해결: Streamlit 세션 초기화 문제 해결"
    p.level = 1

    # Save
    prs.save('AI_CS_Simulator_Project.pptx')
    print("PPTX Created Successfully!")

if __name__ == "__main__":
    create_ai_simulator_deck()

# "AI 고객 응대 시뮬레이터" 발표 자료가 준비되었습니다! HTML 파일을 브라우저에서 열어 전체 화면(F11)으로 띄우시면 바로 발표에 사용하실 수 있습니다. 혹시 수정이 필요하시면 언제든 말씀해주세요.
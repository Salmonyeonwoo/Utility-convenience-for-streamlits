# Copyright (c) Streamlit Inc. (2018-2022) Snowflake Inc. (2022-2025)
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""
시뮬레이션 처리 모듈
고객 응대 시뮬레이션, 채팅/전화 대화 생성, 힌트 생성 등의 기능을 제공합니다.
"""

import os
import json
import uuid
import tempfile
import random
import hashlib
from datetime import datetime
from typing import List, Dict, Any, Tuple, Optional
import streamlit as st

from config import SIM_META_FILE, RAG_INDEX_DIR, DATA_DIR
from utils import _load_json, _save_json
from llm_client import get_api_key, run_llm
from lang_pack import LANG

# Langchain imports for RAG functionality
from langchain_core.documents import Document
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import PyPDFLoader
from langchain_openai import OpenAIEmbeddings
from langchain_community.embeddings import HuggingFaceEmbeddings

# 임베딩 모델 사용 가능 여부 확인
try:
    from langchain_google_genai import GoogleGenerativeAIEmbeddings
    IS_GEMINI_EMBEDDING_AVAILABLE = True
except ImportError:
    IS_GEMINI_EMBEDDING_AVAILABLE = False

try:
    from langchain_nvidia_ai_endpoints import NVIDIAEmbeddings
    IS_NVIDIA_EMBEDDING_AVAILABLE = True
except ImportError:
    IS_NVIDIA_EMBEDDING_AVAILABLE = False

# Word, PPTX, PDF 내보내기 라이브러리
try:
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    IS_DOCX_AVAILABLE = True
except ImportError:
    IS_DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    IS_PPTX_AVAILABLE = True
except ImportError:
    IS_PPTX_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import black
    from reportlab.lib.units import inch
    IS_REPORTLAB_AVAILABLE = True
except ImportError:
    IS_REPORTLAB_AVAILABLE = False

def translate_text_with_llm(text_content: str, target_lang_code: str, source_lang_code: str) -> Tuple[str, bool]:
    """
    주어진 텍스트를 LLM을 사용하여 대상 언어로 번역합니다. (안정화된 텍스트 출력)
    **수정 사항:** LLM Fallback 순서를 OpenAI 우선으로 조정하고, 응답이 비어있을 경우 원본 텍스트를 반환
    
    Returns:
        tuple: (translated_text, is_success) - 번역된 텍스트와 성공 여부
    """
    target_lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}.get(target_lang_code, "English")
    source_lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}.get(source_lang_code, "English")

    # 순수한 텍스트 번역 결과만 출력하도록 강제
    system_prompt = (
        f"You are a professional translation AI. Translate the entire following customer support chat history "
        f"from '{source_lang_name}' to '{target_lang_name}'. "
        f"You MUST translate the content to {target_lang_name} ONLY. "
        f"Do not include any mixed languages, the source text, or any introductory/concluding remarks. "
        f"Output ONLY the translated chat history text. "
    )
    prompt = f"Original Chat History:\n\n{text_content}"

    # LLM Fallback 순서: OpenAI -> Gemini -> Claude (OpenAI를 최우선으로 조정)
    llm_attempts = [
        ("openai", get_api_key("openai"), "gpt-4o"),  # 1순위: OpenAI (가장 안정적)
        ("gemini", get_api_key("gemini"), "gemini-2.5-flash"),  # 2순위
        ("claude", get_api_key("claude"), "claude-3-5-sonnet-latest"),  # 3순위
    ]

    last_error = ""

    for provider, key, model_name in llm_attempts:
        if not key: continue

        try:
            translated_text = ""

            if provider == "openai":
                o_client = OpenAI(api_key=key)
                resp = o_client.chat.completions.create(
                    model=model_name,
                    messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": prompt}],
                    temperature=0.1
                )
                translated_text = resp.choices[0].message.content.strip()

            elif provider == "gemini":
                genai.configure(api_key=key)
                g_model = genai.GenerativeModel(model_name)
                resp = g_model.generate_content(
                    contents=prompt,
                    config=genai.types.GenerateContentConfig(system_instruction=system_prompt, temperature=0.1)
                )
                translated_text = resp.text.strip()

            elif provider == "claude":
                from anthropic import Anthropic
                c_client = Anthropic(api_key=key)
                resp = c_client.messages.create(
                    model=model_name,
                    messages=[{"role": "user", "content": prompt}],
                    system=system_prompt
                )
                translated_text = resp.content[0].text.strip()

            # 번역 결과가 유효한지 확인
            if translated_text and len(translated_text.strip()) > 0:
                return translated_text, True  # 번역 성공
            else:
                last_error = f"Translation failed: {provider} returned empty response."
                continue  # 다음 LLM 시도

        except Exception as e:
            last_error = f"Translation API call failed with {provider} ({model_name}): {e}"  # 모델명 추가
            print(last_error)
            continue  # 다음 LLM 시도

    # 모든 시도가 실패했을 때, 원본 텍스트를 반환하여 프로세스가 계속 진행되도록 함
    # (오류 메시지 대신 원본 텍스트를 반환하여 번역 실패해도 다음 단계로 진행 가능)
    print(f"Translation failed: {last_error or 'No active API key found.'}. Returning original text.")
    return text_content, False  # 원본 텍스트 반환, 번역 실패 표시


# ----------------------------------------
# Realtime Hint Generation (요청 2 반영)
# ----------------------------------------

def generate_realtime_hint(current_lang_key: str, is_call: bool = False):
    """현재 대화 맥락을 기반으로 에이전트에게 실시간 응대 힌트(키워드/정책/액션)를 제공"""
    # 언어 키 검증 및 기본값 처리
    if not current_lang_key or current_lang_key not in ["ko", "en", "ja"]:
        current_lang_key = st.session_state.get("language", "ko")
        if current_lang_key not in ["ko", "en", "ja"]:
            current_lang_key = "ko"
    L = LANG.get(current_lang_key, LANG["ko"])
    # 채팅/전화 구분하여 이력 사용
    if is_call:
        # 전화 시뮬레이터에서는 현재 CC 영역에 표시된 텍스트와 초기 문의를 함께 사용
        website_url = st.session_state.get("call_website_url", "").strip()
        website_context = f"\nWebsite URL: {website_url}" if website_url else ""
        history_text = (
            f"Initial Query: {st.session_state.call_initial_query}\n"
            f"Previous Customer Utterance: {st.session_state.current_customer_audio_text}\n"
            f"Previous Agent Utterance: {st.session_state.current_agent_audio_text}{website_context}"
        )
    else:
        history_text = get_chat_history_for_prompt(include_attachment=True)

    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}[current_lang_key]

    hint_prompt = f"""
You are an AI Supervisor providing an **urgent, internal hint** to a human agent whose AHT is being monitored.
Analyze the conversation history, especially the customer's last message, which might be about complex issues like JR Pass, Universal Studio Japan (USJ), or a complex refund policy.

Provide ONE concise, actionable hint for the agent. The purpose is to save AHT time.

Output MUST be a single paragraph/sentence in {lang_name} containing actionable advice.
DO NOT use markdown headers or titles.
Do NOT direct the agent to check the general website.
Provide an actionable fact or the next specific step (e.g., check policy section, confirm coverage).

Examples of good hints (based on the content):
- Check the official JR Pass site for current exchange rates.
- The 'Universal Express Pass' is non-refundable; clearly cite policy section 3.2.
- Ask for the order confirmation number before proceeding with any action.
- The solution lies in the section of the Klook site titled '~'.

Conversation History:
{history_text}

HINT:
"""
    if not st.session_state.is_llm_ready:
        return "(Mock Hint: LLM Key is missing. Ask the customer for the booking number.)"

    with st.spinner(f"💡 {L['button_request_hint']}..."):
        try:
            return run_llm(hint_prompt).strip()
        except Exception as e:
            return f"❌ Hint Generation Error. (Try again or check API Key: {e})"



def generate_agent_response_draft(current_lang_key: str) -> str:
    """고객 응답을 기반으로 AI가 에이전트 응답 초안을 생성 (요청 1 반영)"""
    # 언어 키 검증 및 기본값 처리
    if not current_lang_key or current_lang_key not in ["ko", "en", "ja"]:
        current_lang_key = st.session_state.get("language", "ko")
        if current_lang_key not in ["ko", "en", "ja"]:
            current_lang_key = "ko"
    L = LANG.get(current_lang_key, LANG["ko"])
    history_text = get_chat_history_for_prompt(include_attachment=True)
    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}[current_lang_key]

    # 고객의 최신 문의 내용 추출 및 분석 (강화)
    latest_customer_message = ""
    initial_customer_query = st.session_state.get('customer_query_text_area', '')
    customer_query_analysis = ""
    
    # 모든 고객 메시지 수집
    all_customer_messages = []
    if st.session_state.simulator_messages:
        all_customer_messages = [msg["content"] for msg in st.session_state.simulator_messages 
                                if msg.get("role") in ["customer", "customer_rebuttal", "initial_query"]]
    
    # 초기 문의도 포함
    if initial_customer_query and initial_customer_query not in all_customer_messages:
        all_customer_messages.insert(0, initial_customer_query)
    
    if all_customer_messages:
        latest_customer_message = all_customer_messages[-1]
        
        # 문의 내용에서 핵심 정보 추출 (간단한 키워드 추출)
        inquiry_keywords = []
        inquiry_text = " ".join(all_customer_messages).lower()
        
        # 일반적인 문의 키워드 패턴
        important_patterns = [
            r'\b\d{4,}\b',  # 주문번호, 전화번호 등 숫자
            r'\b(주문|order|注文)\b',
            r'\b(환불|refund|返金)\b',
            r'\b(취소|cancel|キャンセル)\b',
            r'\b(배송|delivery|配送)\b',
            r'\b(변경|change|変更)\b',
            r'\b(문제|problem|issue|問題)\b',
            r'\b(도움|help|助け)\b',
        ]
        
        # 핵심 문의 내용 요약
        inquiry_summary = f"""
**CUSTOMER INQUIRY DETAILS:**

Initial Query: "{initial_customer_query if initial_customer_query else 'Not provided'}"

Latest Customer Message: "{latest_customer_message}"

All Customer Messages Context:
{chr(10).join([f"- {msg[:150]}..." if len(msg) > 150 else f"- {msg}" for msg in all_customer_messages[-3:]])}

**YOUR RESPONSE MUST DIRECTLY ADDRESS:**

1. **SPECIFIC ISSUE IDENTIFICATION**: 
   - What EXACT problem or question did the customer mention?
   - Extract and reference specific details: order numbers, dates, product names, locations, error messages, etc.
   - If multiple issues were mentioned, address EACH one specifically

2. **CONCRETE SOLUTION PROVIDED**:
   - Provide STEP-BY-STEP instructions tailored to their EXACT situation
   - Include specific actions they need to take (e.g., "Go to Settings > Account > Order History and click on order #12345")
   - Reference the exact products/services they mentioned
   - If they mentioned a location, reference it in your solution

3. **PERSONALIZATION**:
   - Use the customer's specific words/phrases when appropriate
   - Reference their exact situation (e.g., "Since you mentioned your eSIM isn't activating in Paris...")
   - Acknowledge their specific concern or frustration point

4. **COMPLETENESS**:
   - Answer ALL questions they asked
   - Address ALL problems they mentioned
   - If they asked "why", explain the specific reason for their situation
   - If they asked "how", provide detailed steps for their exact case

**CRITICAL REQUIREMENTS:**
- DO NOT use generic templates like "Thank you for contacting us" without addressing their specific issue
- DO NOT give vague answers like "Please check your settings" - be SPECIFIC about which settings and what to do
- DO NOT ignore specific details they mentioned (order numbers, dates, locations, etc.)
- Your response must read as if it was written SPECIFICALLY for this customer's exact inquiry
- If the customer mentioned "eSIM activation in Paris", your response MUST specifically address eSIM activation and Paris
- If the customer mentioned an order number, your response MUST reference that order number

**EXAMPLE OF GOOD RESPONSE:**
Bad: "Thank you for contacting us. We understand your concern and will help you resolve this issue."
Good: "I understand you're having trouble activating your eSIM in Paris. Let me help you resolve this step by step. First, please check if your phone's APN settings are configured correctly for the Paris network..."

**NOW GENERATE YOUR RESPONSE** following these requirements:
"""
        
        customer_query_analysis = inquiry_summary

    # 첨부 파일 컨텍스트 추가
    attachment_context = st.session_state.sim_attachment_context_for_llm
    if attachment_context:
        attachment_context = f"\n[고객 첨부 파일 정보: {attachment_context}]\n"
    else:
        attachment_context = ""

    # 고객 검증 상태 확인 (로그인/계정 관련 문의인 경우)
    is_login_inquiry = check_if_login_related_inquiry(initial_customer_query)
    is_customer_verified = st.session_state.get("is_customer_verified", False)
    verification_warning = ""
    
    if is_login_inquiry and not is_customer_verified:
        # 검증되지 않은 고객에게는 정보 힌트 제공 금지
        customer_email = st.session_state.get("customer_email", "")
        customer_phone = st.session_state.get("customer_phone", "")
        customer_name = st.session_state.get("customer_name", "")
        
        # 이메일은 마스킹하여 힌트 제공 가능 (앞/뒤 1-3자리만)
        masked_email = mask_email(customer_email, show_chars=2) if customer_email else ""
        
        verification_warning = f"""
**⚠️ CRITICAL SECURITY REQUIREMENT - CUSTOMER VERIFICATION NOT COMPLETED:**

This is a LOGIN/ACCOUNT related inquiry, but the customer has NOT been verified yet.

**STRICT RULES YOU MUST FOLLOW:**
1. **DO NOT provide ANY customer information hints** (email, phone, name, receipt number, card number) in your response
2. **EXCEPTION**: You MAY provide a masked email hint ONLY if absolutely necessary: "{masked_email}" (only first/last 2-3 characters visible, rest masked)
3. **DO NOT reveal**: Full email addresses, phone numbers, customer names, receipt numbers, card numbers, or any other personal information
4. **You MUST request verification information** from the customer before proceeding with account-related assistance
5. **Your response should ask the customer to provide verification details** (receipt number, card last 4 digits, name, email, phone) WITHOUT revealing what information you already have

**VERIFICATION REQUEST TEMPLATE:**
- "To ensure account security, please provide the following information for verification: [receipt number / card last 4 digits / name / email / phone]"
- DO NOT mention what information you already have
- DO NOT provide hints about the information

**ONLY AFTER VERIFICATION IS COMPLETED** can you provide full information hints and proceed with account assistance.

**CURRENT STATUS**: Customer verification: NOT COMPLETED ❌
"""
    elif is_login_inquiry and is_customer_verified:
        verification_warning = """
**✅ CUSTOMER VERIFICATION COMPLETED:**

The customer has been successfully verified. You may now provide information hints and proceed with account-related assistance.
"""

    # 고객 유형 및 반복 불만 패턴 분석
    customer_type = st.session_state.get('customer_type_sim_select', '일반적인 문의')
    is_difficult_customer = customer_type in ["까다로운 고객", "매우 불만족스러운 고객", "Difficult Customer",
                                              "Highly Dissatisfied Customer", "難しい顧客", "非常に不満な顧客"]

    # 고객 메시지 수 및 감정 분석
    customer_message_count = sum(
        1 for msg in st.session_state.simulator_messages if msg.get("role") in ["customer", "customer_rebuttal"])
    agent_message_count = sum(1 for msg in st.session_state.simulator_messages if msg.get("role") == "agent_response")

    # 이전 에이전트 응답들 추출 (다양성 확보를 위해)
    previous_agent_responses = [msg["content"] for msg in st.session_state.simulator_messages 
                                if msg.get("role") == "agent_response"]
    previous_responses_context = ""
    if previous_agent_responses:
        previous_responses_context = f"\n[이전 에이전트 응답들 (참고용, 동일하게 반복하지 말 것):\n"
        for i, prev_resp in enumerate(previous_agent_responses[-3:], 1):  # 최근 3개만
            prev_resp_preview = prev_resp[:200] + "..." if len(prev_resp) > 200 else prev_resp
            previous_responses_context += f"{i}. {prev_resp_preview}\n"
        previous_responses_context += "]\n"

    # 고객이 계속 따지거나 화내는 패턴 감지 (고객 메시지가 에이전트 메시지보다 많거나, 반복적인 불만 표현)
    is_repeating_complaints = False
    if customer_message_count > agent_message_count and customer_message_count >= 2:
        # 마지막 2개 고객 메시지 분석
        recent_customer_messages = [msg["content"].lower() for msg in st.session_state.simulator_messages if
                                    msg.get("role") in ["customer", "customer_rebuttal"]][-2:]
        complaint_keywords = ["왜", "이유", "설명", "말이 안", "이해가 안", "화나", "짜증", "불만", "왜", "why", "reason", "explain",
                              "angry", "frustrated", "complaint", "なぜ", "理由", "説明", "怒り", "不満"]
        if any(any(keyword in msg for keyword in complaint_keywords) for msg in recent_customer_messages):
            is_repeating_complaints = True

    # 대처법 포메이션 추가 여부 결정
    needs_coping_strategy = is_difficult_customer or (is_repeating_complaints and customer_message_count >= 2)

    # 대처법 가이드라인 생성
    coping_guidance = ""
    if needs_coping_strategy:
        coping_guidance = f"""

[CRITICAL: Handling Difficult Customer Situation]
The customer type is "{customer_type}" and the customer has sent {customer_message_count} messages while the agent has sent {agent_message_count} messages.
The customer may be showing signs of continued frustration or dissatisfaction.

**INCLUDE THE FOLLOWING COPING STRATEGY FORMAT IN YOUR RESPONSE:**

1. **Immediate Acknowledgment** (1-2 sentences):
   - Acknowledge their frustration/specific concern explicitly
   - Show deep empathy and understanding
   - Example formats:
     * "{'죄송합니다. 불편을 드려 정말 죄송합니다. 고객님의 상황을 충분히 이해하고 있습니다.' if current_lang_key == 'ko' else ('I sincerely apologize for the inconvenience. I fully understand your situation and frustration.' if current_lang_key == 'en' else '大変申し訳ございません。お客様の状況とご不便を十分に理解しております。')}"
     * "{'고객님의 소중한 의견을 잘 듣고 있습니다. 정말 답답하셨을 것 같습니다.' if current_lang_key == 'ko' else ('I hear your concerns clearly. This must have been very frustrating for you.' if current_lang_key == 'en' else 'お客様のご意見をしっかりと受け止めています。本当にお困りだったと思います。')}"

2. **Specific Solution Recap** (2-3 sentences):
   - Clearly restate the solution/step provided previously (if any)
   - Offer a NEW concrete action or alternative solution
   - Be specific and actionable
   - Example formats:
     * "{'앞서 안내드린 [구체적 해결책] 외에도, [새로운 대안/추가 조치]를 진행해드릴 수 있습니다.' if current_lang_key == 'ko' else ('In addition to the [specific solution] I mentioned earlier, I can also [new alternative/additional action] for you.' if current_lang_key == 'en' else '先ほどご案内した[具体的解決策]に加えて、[新しい代替案/追加措置]も進めることができます。')}"
     * "{'혹시 [구체적 문제점] 때문에 불편하셨다면, [구체적 해결 방법]을 바로 진행해드리겠습니다.' if current_lang_key == 'ko' else ('If you are experiencing [specific issue], I can immediately proceed with [specific solution].' if current_lang_key == 'en' else 'もし[具体的問題]でご不便でしたら、[具体的解決方法]をすぐに進めさせていただきます。')}"

3. **Escalation or Follow-up Offer** (1-2 sentences):
   - Offer to escalate to supervisor/higher level support
   - Promise immediate follow-up within specific time
   - Example formats:
     * "{'만약 여전히 불만이 해소되지 않으신다면, 즉시 상급 관리자에게 이관하여 더 나은 해결책을 찾아드리겠습니다.' if current_lang_key == 'ko' else ('If your concern is still not resolved, I can immediately escalate this to a supervisor to find a better solution.' if current_lang_key == 'en' else 'もしご不満が解消されない場合は、すぐに上級管理者にエスカレートして、より良い解決策を見つけさせていただきます。')}"
     * "{'24시간 이내에 [구체적 조치/결과]를 확인하여 고객님께 다시 연락드리겠습니다.' if current_lang_key == 'ko' else ('I will follow up with you within 24 hours regarding [specific action/result].' if current_lang_key == 'en' else '24時間以内に[具体的措置/結果]を確認し、お客様に再度ご連絡いたします。')}"

4. **Closing with Assurance** (1 sentence):
   - Reassure that their concern is being taken seriously
   - Example formats:
     * "{'고객님의 모든 문의사항을 최우선으로 처리하겠습니다.' if current_lang_key == 'ko' else ('I will prioritize resolving all of your concerns.' if current_lang_key == 'en' else 'お客様のすべてのご質問を最優先で処理いたします。')}"

**IMPORTANT NOTES:**
- DO NOT repeat the exact same solution that was already provided
- DO NOT sound dismissive or automated
- DO sound genuinely concerned and willing to go the extra mile
- If policy restrictions exist, acknowledge them but still offer alternatives
- Use warm, respectful tone while being firm about what can/cannot be done

**RESPONSE STRUCTURE:**
[Immediate Acknowledgment]
[Specific Solution Recap + New Action]
[Escalation/Follow-up Offer]
[Closing with Assurance]

Now generate the agent's response draft following this structure:
"""

    # 다양성 확보를 위한 추가 지시사항 (더 강화)
    diversity_instruction = ""
    if previous_agent_responses:
        # 이전 응답들의 주요 키워드/구문 추출 (반복 방지)
        previous_keywords = []
        for prev_resp in previous_agent_responses[-3:]:
            # 간단한 키워드 추출 (2-3단어 구문)
            words = prev_resp.split()[:20]  # 처음 20단어만
            for i in range(len(words) - 1):
                if len(words[i]) > 3 and len(words[i+1]) > 3:
                    previous_keywords.append(f"{words[i]} {words[i+1]}")
        
        keywords_warning = ""
        if previous_keywords:
            unique_keywords = list(set(previous_keywords))[:10]  # 최대 10개만
            keywords_warning = f"\n- AVOID using these exact phrases from previous responses: {', '.join(unique_keywords[:5])}"
        
        diversity_instruction = f"""
**CRITICAL DIVERSITY REQUIREMENT - STRICTLY ENFORCED:**
- You MUST generate a COMPLETELY DIFFERENT response from ALL previous agent responses shown above
- Use COMPLETELY DIFFERENT wording, phrasing, sentence structures, and vocabulary
- If similar solutions are needed, present them in a COMPLETELY DIFFERENT way or from a COMPLETELY DIFFERENT angle
- Vary your opening sentences, transition phrases, and closing statements - NO REPETITION
- DO NOT copy, paraphrase, or reuse ANY phrases from previous responses - be CREATIVE and UNIQUE while maintaining professionalism
- Consider COMPLETELY different approaches: 
  * If previous responses were formal, try a warmer, more personal tone (or vice versa)
  * If previous responses focused on one aspect, emphasize a COMPLETELY different aspect this time
  * Use different examples, analogies, or explanations
  * Change the order of information presentation
  * Use different sentence lengths and structures
{keywords_warning}
- IMPORTANT: Read ALL previous responses carefully and ensure your response is DISTINCTLY different in style, tone, structure, and content
- If you find yourself writing something similar to a previous response, STOP and rewrite it completely differently
"""

    # 랜덤 요소 추가를 위한 변형 지시사항
    variation_approaches = [
        "Start with a different greeting or acknowledgment style",
        "Use a different problem-solving approach or framework",
        "Present information in a different order",
        "Use different examples or analogies",
        "Vary the level of formality or warmth",
        "Focus on different aspects of the solution",
        "Use different transition words and phrases",
        "Change the length and complexity of sentences"
    ]
    selected_approaches = random.sample(variation_approaches, min(3, len(variation_approaches)))
    variation_note = "\n".join([f"- {approach}" for approach in selected_approaches])

    draft_prompt = f"""
You are an AI assistant helping a customer support agent write a professional, tailored response.

**PRIMARY OBJECTIVE:**
Generate a response draft that is SPECIFICALLY tailored to the customer's EXACT inquiry, providing concrete, actionable solutions that directly address their specific situation. The response must read as if it was written personally for this customer's unique case.

**CRITICAL REQUIREMENTS (IN ORDER OF PRIORITY):**
1. **MOST CRITICAL**: Address the customer's SPECIFIC inquiry/question with DETAILED, ACTIONABLE solutions
   - Extract and reference specific details from their message (order numbers, dates, product names, locations, error messages, etc.)
   - Provide step-by-step instructions tailored to their EXACT situation
   - Answer ALL questions they asked completely
   - Address ALL problems they mentioned specifically

2. The response MUST be in {lang_name}

3. Be professional, empathetic, and solution-oriented

4. If the customer asked a question, provide a COMPLETE and SPECIFIC answer - do NOT give vague or generic responses
   - Bad: "Please check your settings"
   - Good: "Please go to Settings > Mobile Network > APN Settings and ensure the APN is set to 'internet'"

5. If the customer mentioned a problem, acknowledge it SPECIFICALLY and provide STEP-BY-STEP solutions
   - Reference their exact problem description
   - Provide solutions that directly address their specific issue

6. Reference specific details from their inquiry (order numbers, dates, products, locations, etc.) if mentioned
   - If they mentioned "order #12345", your response MUST include "order #12345"
   - If they mentioned "Paris", your response should reference Paris specifically
   - If they mentioned "eSIM", address eSIM specifically, not just "SIM card"

7. Keep the tone appropriate for the customer type: {customer_type}

8. Do NOT include any markdown formatting, just plain text

9. {f'**FOLLOW THE COPING STRATEGY FORMAT BELOW**' if needs_coping_strategy else 'Use natural, conversational flow'}

10. **CRITICAL**: Generate a COMPLETELY UNIQUE and VARIED response - avoid repeating ANY similar phrases, structures, or approaches from previous responses

11. **CRITICAL**: Your response must be HIGHLY RELEVANT to the customer's specific inquiry - generic template responses are NOT acceptable
    - DO NOT start with generic greetings without immediately addressing their specific issue
    - DO NOT use placeholder text like "[specific solution]" - provide ACTUAL specific solutions
    - Your response should make the customer feel like you read and understood THEIR specific message

**VARIATION TECHNIQUES TO APPLY:**
{variation_note}

{customer_query_analysis}

**FULL CONVERSATION HISTORY:**
{history_text}
{attachment_context}

{verification_warning}

**PREVIOUS RESPONSES TO AVOID REPEATING:**
{previous_responses_context if previous_responses_context else "No previous responses to compare against."}

**DIVERSITY REQUIREMENTS:**
{diversity_instruction if diversity_instruction else "This is the first response, so no previous responses to avoid."}

{coping_guidance if needs_coping_strategy else ''}

**YOUR TASK:**
Generate the agent's response draft NOW. The response must:

1. **FIRST**: Read the customer inquiry analysis above CAREFULLY and identify:
   - What is their EXACT problem or question?
   - What specific details did they mention (order numbers, dates, locations, products)?
   - What do they need help with specifically?

2. **SECOND**: Write a response that:
   - Addresses their EXACT problem/question (not a generic version)
   - References the specific details they mentioned
   - Provides concrete, actionable steps tailored to their situation
   - Answers ALL their questions completely
   - Makes them feel understood

3. **THIRD**: Ensure the response is:
   - COMPLETELY DIFFERENT and UNIQUE from any previous responses
   - Professional, empathetic, and solution-oriented
   - Written in {lang_name}
   - Free of markdown formatting

**BEFORE YOU WRITE**: Ask yourself:
- "Does this response address the customer's SPECIFIC inquiry?"
- "Would a generic template response work here?" (If yes, rewrite it to be more specific)
- "Does this response reference specific details from the customer's message?"
- "Would the customer feel like I read and understood THEIR specific message?"

**NOW GENERATE THE RESPONSE:**
"""

    if not st.session_state.is_llm_ready:
        return ""

    try:
        draft = run_llm(draft_prompt).strip()
        # 마크다운 제거 (``` 등)
        if draft.startswith("```"):
            lines = draft.split("\n")
            draft = "\n".join(lines[1:-1]) if len(lines) > 2 else draft
        return draft
    except Exception as e:
        return f"❌ 응답 초안 생성 오류: {e}"


# ⭐ 새로운 함수: 전화 발신 시뮬레이션 요약 생성

def generate_outbound_call_summary(customer_query: str, current_lang_key: str, target: str) -> str:
    """
    Simulates an outbound call to a local partner or customer and generates a summary of the outcome.
    """
    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}[current_lang_key]

    # Get the current chat history for context
    history_text = get_chat_history_for_prompt(include_attachment=True)
    if not history_text:
        history_text = f"Initial Customer Query: {customer_query}"

    # Policy context (from supervisor) should be included to guide the outcome
    policy_context = st.session_state.supervisor_policy_context or ""

    summary_prompt = f"""
You are an AI simulating a quick, high-stakes phone call placed by the customer support agent to a '{target}' (either a local partner/vendor or the customer).

The purpose of the call is to resolve a complex, policy-restricted issue (like an exceptional refund for a non-refundable item, or urgent confirmation of an airport transfer change).

Analyze the conversation history, the initial query, and any provided supervisor policy.
Generate a concise summary of the OUTCOME of this simulated phone call.
The summary MUST be professional and strictly in {lang_name}.

[CRITICAL RULE]: For non-refundable items (e.g., Universal Studio Express Pass, non-refundable hotel/transfer), the local partner should only grant an exception IF the customer has provided strong, unavoidable proof (like a flight cancellation notice, doctor's note, or natural disaster notice). If no such proof is evident in the chat history, the outcome should usually be a denial or a request for more proof, but keep the tone professional.
If the customer's query is about Airport Transfer change, the outcome should be: 'Confirmation complete. Change is approved/denied based on partner policy.'

Conversation History:
{history_text}

Supervisor Policy Context (If any):
{policy_context}

Target of Call: {target}

Generate the phone call summary (Outcome ONLY):
"""
    if not st.session_state.is_llm_ready:
        return f"❌ LLM Key missing. (Simulated Outcome: The {target} requested the agent to send proof via email.)"

    try:
        summary = run_llm(summary_prompt).strip()
        # 마크다운 제거 (``` 등)
        if summary.startswith("```"):
            lines = summary.split("\n")
            summary = "\n".join(lines[1:-1]) if len(lines) > 2 else summary
        return summary
    except Exception as e:
        return f"❌ Phone call simulation error: {e}"


# ========================================
# 3. Whisper / TTS Helper
# ========================================

def transcribe_bytes_with_whisper(audio_bytes: bytes, mime_type: str = "audio/webm", lang_code: str = None, auto_detect: bool = True) -> str:
    """
    OpenAI Whisper API 또는 Gemini API를 사용하여 오디오 바이트를 텍스트로 전사합니다.
    OpenAI가 실패하면 Gemini로 자동 fallback합니다.
    
    Args:
        audio_bytes: 전사할 오디오 바이트
        mime_type: 오디오 MIME 타입
        lang_code: 언어 코드 (ko, en, ja 등). None이거나 auto_detect=True이면 자동 감지
        auto_detect: True이면 언어를 자동 감지 (lang_code 무시)
    """
    # 언어 키 안전하게 가져오기
    current_lang = st.session_state.get("language", "ko")
    if current_lang not in ["ko", "en", "ja"]:
        current_lang = "ko"
    L = LANG.get(current_lang, LANG["ko"])
    
    # 임시 파일 저장 (API 호환성)
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".wav")
    tmp.write(audio_bytes)
    tmp.flush()
    tmp.close()
    
    # 1️⃣ OpenAI Whisper API 시도
    client = st.session_state.openai_client
    if client is not None:
        try:
            with open(tmp.name, "rb") as f:
                # 언어 자동 감지 또는 지정된 언어 사용
                if auto_detect or lang_code is None:
                    # language 파라미터를 생략하면 Whisper가 자동으로 언어를 감지합니다
                    res = client.audio.transcriptions.create(
                        model="whisper-1",
                        file=f,
                        response_format="text",
                    )
                else:
                    whisper_lang = {"ko": "ko", "en": "en", "ja": "ja"}.get(lang_code, "en")
                    res = client.audio.transcriptions.create(
                        model="whisper-1",
                        file=f,
                        response_format="text",
                        language=whisper_lang,
                    )
            # res.text 속성이 있는지 확인하고 없으면 res 자체를 문자열로 변환
            result = res.text.strip() if hasattr(res, 'text') else str(res).strip()
            if result:
                try:
                    os.remove(tmp.name)
                except OSError:
                    pass
                return result
        except Exception as e:
            # OpenAI 실패 시 로그만 남기고 Gemini로 fallback
            print(f"OpenAI Whisper failed: {e}")
    
    # 2️⃣ Gemini API fallback
    gemini_key = get_api_key("gemini")
    if gemini_key:
        try:
            import base64
            genai.configure(api_key=gemini_key)
            
            # Gemini는 오디오 파일을 base64로 인코딩하여 전송
            with open(tmp.name, "rb") as f:
                audio_data = f.read()
                audio_base64 = base64.b64encode(audio_data).decode('utf-8')
            
            # Gemini 2.0 Flash 모델 사용 (오디오 지원)
            model = genai.GenerativeModel("gemini-2.0-flash-exp")
            
            # 프롬프트 구성
            lang_prompt = ""
            if lang_code:
                lang_map = {"ko": "한국어", "en": "English", "ja": "日本語"}
                lang_prompt = f"이 오디오는 {lang_map.get(lang_code, 'English')}로 말하고 있습니다. "
            
            prompt = f"{lang_prompt}이 오디오를 텍스트로 전사해주세요. 오직 전사된 텍스트만 반환하세요."
            
            # Gemini는 파일 업로드 방식 사용 (Gemini 2.0 Flash는 오디오 지원)
            try:
                audio_file = genai.upload_file(path=tmp.name, mime_type=mime_type)
                
                # 파일 업로드 후 잠시 대기 (업로드 완료 대기)
                import time
                time.sleep(1)
                
                response = model.generate_content([prompt, audio_file])
                result = response.text.strip() if response.text else ""
                
                # 파일 삭제
                try:
                    genai.delete_file(audio_file.name)
                except Exception as del_err:
                    print(f"Failed to delete Gemini file: {del_err}")
            except Exception as upload_err:
                # 파일 업로드 실패 시 다른 방법 시도
                print(f"Gemini file upload failed: {upload_err}")
                # 대안: base64 인코딩된 오디오를 직접 전송 (모델이 지원하는 경우)
                raise upload_err
            
            if result:
                try:
                    os.remove(tmp.name)
                except OSError:
                    pass
                return result
            else:
                raise Exception("Gemini returned empty result")
        except Exception as e:
            print(f"Gemini transcription failed: {e}")
            # Gemini도 실패한 경우 에러 메시지 반환
            try:
                os.remove(tmp.name)
            except OSError:
                pass
            return f"❌ {L.get('whisper_client_error', '전사 실패')}: OpenAI와 Gemini 모두 실패했습니다. ({str(e)[:100]})"
    else:
        # 두 API 모두 사용 불가
        try:
            os.remove(tmp.name)
        except OSError:
            pass
        return f"❌ {L.get('openai_missing', 'OpenAI API Key가 필요합니다.')} 또는 Gemini API Key가 필요합니다."


def transcribe_audio(audio_bytes, filename="audio.wav"):
    client = st.session_state.openai_client

    # 1️⃣ OpenAI Whisper 시도
    if client:
        try:
            import io
            bio = io.BytesIO(audio_bytes)
            bio.name = filename
            resp = client.audio.transcriptions.create(
                model="whisper-1",
                file=bio,
            )
            return resp.text
        except Exception as e:
            print("Whisper OpenAI failed:", e)

    # 2️⃣ Gemini STT fallback
    try:
        genai.configure(api_key=get_api_key("gemini"))
        model = genai.GenerativeModel("gemini-2.5-flash")
        text = model.generate_content("Transcribe this audio:").text
        return text or ""
    except Exception as e:
        print("Gemini STT failed:", e)

    return "❌ STT not available"


# ========================================
# 비디오 동기화 관련 함수
# ========================================

def analyze_text_for_video_selection(text: str, current_lang_key: str, 
                                     agent_last_response: str = None,
                                     conversation_context: List[Dict] = None) -> Dict[str, Any]:
    """
    LLM을 사용하여 텍스트를 분석하고 적절한 감정 상태와 제스처를 판단합니다.
    OpenAI/Gemini API를 활용한 영상 RAG의 핵심 기능입니다.
    
    ⭐ Gemini 제안 적용: 긴급도, 만족도 변화, 에이전트 답변 기반 예측 추가
    
    Args:
        text: 분석할 텍스트 (고객의 질문/응답)
        current_lang_key: 현재 언어 키
        agent_last_response: 에이전트의 마지막 답변 (선택적, 예측 정확도 향상)
        conversation_context: 대화 컨텍스트 (선택적, 만족도 변화 분석용)
    
    Returns:
        {
            "emotion": "NEUTRAL" | "HAPPY" | "ANGRY" | "ASKING" | "SAD",
            "gesture": "NONE" | "HAND_WAVE" | "NOD" | "SHAKE_HEAD" | "POINT",
            "urgency": "LOW" | "MEDIUM" | "HIGH",  # ⭐ 추가: 긴급도
            "satisfaction_delta": -1.0 to 1.0,  # ⭐ 추가: 만족도 변화 (-1: 감소, 0: 유지, 1: 증가)
            "confidence": 0.0-1.0
        }
    """
    if not text or not text.strip():
        return {
            "emotion": "NEUTRAL", 
            "gesture": "NONE", 
            "urgency": "LOW",
            "satisfaction_delta": 0.0,
            "confidence": 0.5
        }
    
    L = LANG.get(current_lang_key, LANG["ko"])
    
    # ⭐ Gemini 제안: 에이전트 답변 기반 예측 컨텍스트 구성
    context_info = ""
    if agent_last_response:
        context_info = f"""
에이전트의 마지막 답변: "{agent_last_response}"

에이전트의 답변을 고려했을 때, 고객이 지금 말하는 내용은 어떤 감정을 수반할 것인지 예측하세요.
예를 들어:
- 에이전트가 솔루션을 제시했다면 → 고객은 HAPPY 또는 ASKING (추가 질문)
- 에이전트가 거절했다면 → 고객은 ANGRY 또는 SAD
- 에이전트가 질문을 했다면 → 고객은 ASKING (답변) 또는 NEUTRAL
"""
    
    # ⭐ Gemini 제안: 만족도 변화 분석 컨텍스트
    satisfaction_context = ""
    if conversation_context and len(conversation_context) > 1:
        # 최근 대화의 감정 변화 추적
        recent_emotions = []
        for msg in conversation_context[-3:]:  # 최근 3개 메시지
            if msg.get("role") == "customer_rebuttal" or msg.get("role") == "customer":
                recent_emotions.append(msg.get("content", ""))
        
        if len(recent_emotions) >= 2:
            satisfaction_context = f"""
최근 대화 흐름:
- 이전 고객 메시지: "{recent_emotions[-2] if len(recent_emotions) >= 2 else ''}"
- 현재 고객 메시지: "{recent_emotions[-1]}"

만족도 변화를 분석하세요:
- 이전보다 더 긍정적이면 satisfaction_delta > 0
- 이전보다 더 부정적이면 satisfaction_delta < 0
- 비슷하면 satisfaction_delta ≈ 0
"""
    
    # ⭐ Gemini 제안: 개선된 LLM 프롬프트 구성
    prompt = f"""다음 고객의 텍스트를 분석하여 적절한 감정 상태, 제스처, 긴급도, 만족도 변화를 판단하세요.

고객 텍스트: "{text}"
{context_info}
{satisfaction_context}

다음 JSON 형식으로만 응답하세요 (다른 설명 없이):
{{
    "emotion": "NEUTRAL" | "HAPPY" | "ANGRY" | "ASKING" | "SAD",
    "gesture": "NONE" | "HAND_WAVE" | "NOD" | "SHAKE_HEAD" | "POINT",
    "urgency": "LOW" | "MEDIUM" | "HIGH",
    "satisfaction_delta": -1.0 to 1.0,
    "confidence": 0.0-1.0
}}

감정 판단 기준 (세분화):
- HAPPY: 긍정적 표현, 감사, 만족, 해결됨 ("감사합니다", "좋아요", "완벽해요", "이제 이해했어요")
- ANGRY: 불만, 화남, 거부, 강한 부정 ("화가 나요", "불가능해요", "거절합니다", "말도 안 돼요")
- ASKING: 질문, 궁금함, 확인 요청, 정보 요구 ("어떻게", "왜", "알려주세요", "주문번호가 뭐예요?")
- SAD: 슬픔, 실망, 좌절 ("슬프네요", "실망했어요", "아쉽습니다", "그렇다면 어쩔 수 없네요")
- NEUTRAL: 중립적 표현, 단순 정보 전달 (기본값)

제스처 판단 기준:
- HAND_WAVE: 인사, 환영 ("안녕하세요", "반갑습니다")
- NOD: 동의, 긍정, 이해 ("네", "맞아요", "그렇습니다", "알겠습니다")
- SHAKE_HEAD: 부정, 거부, 불만족 ("아니요", "안 됩니다", "그건 아니에요")
- POINT: 설명, 지시, 특정 항목 언급 ("여기", "이것", "저것", "주문번호는")
- NONE: 특별한 제스처 없음 (기본값)

긴급도 판단 기준:
- HIGH: 즉시 해결 필요, 긴급한 문제 ("지금 당장", "바로", "긴급", "중요해요")
- MEDIUM: 빠른 해결 선호, 중요하지만 긴급하지 않음
- LOW: 일반적인 문의, 긴급하지 않음 (기본값)

만족도 변화 (satisfaction_delta):
- 1.0: 매우 만족, 문제 해결됨, 감사 표현
- 0.5: 만족, 긍정적 반응
- 0.0: 중립, 변화 없음
- -0.5: 불만족, 부정적 반응
- -1.0: 매우 불만족, 화남, 거부

JSON만 응답하세요:"""

    try:
        # LLM 호출
        if st.session_state.is_llm_ready:
            response_text = run_llm(prompt)
            
            # JSON 파싱 시도
            try:
                # JSON 부분만 추출 (코드 블록 제거)
                import re
                json_match = re.search(r'\{[^{}]*\}', response_text, re.DOTALL)
                if json_match:
                    result = json.loads(json_match.group())
                    # 유효성 검사
                    valid_emotions = ["NEUTRAL", "HAPPY", "ANGRY", "ASKING", "SAD"]
                    valid_gestures = ["NONE", "HAND_WAVE", "NOD", "SHAKE_HEAD", "POINT"]
                    valid_urgencies = ["LOW", "MEDIUM", "HIGH"]
                    
                    emotion = result.get("emotion", "NEUTRAL")
                    gesture = result.get("gesture", "NONE")
                    urgency = result.get("urgency", "LOW")
                    satisfaction_delta = float(result.get("satisfaction_delta", 0.0))
                    confidence = float(result.get("confidence", 0.7))
                    
                    if emotion not in valid_emotions:
                        emotion = "NEUTRAL"
                    if gesture not in valid_gestures:
                        gesture = "NONE"
                    if urgency not in valid_urgencies:
                        urgency = "LOW"
                    
                    # ⭐ Gemini 제안: 상황별 키워드 추출
                    context_keywords = []
                    text_lower_for_context = text.lower()
                    
                    # 주요 상황별 키워드 매핑
                    if any(word in text_lower_for_context for word in ["주문번호", "order number", "주문 번호"]):
                        context_keywords.append("order_number")
                    if any(word in text_lower_for_context for word in ["해결", "완료", "감사", "solution", "resolved"]):
                        if satisfaction_delta > 0.3:
                            context_keywords.append("solution_accepted")
                    if any(word in text_lower_for_context for word in ["거절", "불가", "안 됩니다", "denied", "cannot"]):
                        if emotion == "ANGRY":
                            context_keywords.append("policy_denial")
                    
                    return {
                        "emotion": emotion,
                        "gesture": gesture,
                        "urgency": urgency,
                        "satisfaction_delta": max(-1.0, min(1.0, satisfaction_delta)),
                        "context_keywords": context_keywords,  # ⭐ 추가
                        "confidence": max(0.0, min(1.0, confidence))
                    }
            except json.JSONDecodeError:
                pass
        
        # LLM 호출 실패 시 키워드 기반 간단한 분석
        text_lower = text.lower()
        emotion = "NEUTRAL"
        gesture = "NONE"
        urgency = "LOW"
        satisfaction_delta = 0.0
        
        # 감정 키워드 분석
        if any(word in text_lower for word in ["감사", "좋아", "완벽", "만족", "고마워", "해결"]):
            emotion = "HAPPY"
            satisfaction_delta = 0.5
        elif any(word in text_lower for word in ["화", "불만", "거절", "불가능", "안 됩니다", "말도 안 돼"]):
            emotion = "ANGRY"
            satisfaction_delta = -0.5
        elif any(word in text_lower for word in ["어떻게", "왜", "알려", "질문", "궁금", "주문번호"]):
            emotion = "ASKING"
        elif any(word in text_lower for word in ["슬프", "실망", "아쉽", "그렇다면"]):
            emotion = "SAD"
            satisfaction_delta = -0.3
        
        # 긴급도 키워드 분석
        if any(word in text_lower for word in ["지금 당장", "바로", "긴급", "중요해요", "즉시"]):
            urgency = "HIGH"
        elif any(word in text_lower for word in ["빨리", "가능한 한", "최대한"]):
            urgency = "MEDIUM"
        
        # 제스처 키워드 분석
        if any(word in text_lower for word in ["안녕", "반갑", "인사"]):
            gesture = "HAND_WAVE"
        elif any(word in text_lower for word in ["네", "맞아", "그래", "동의", "알겠습니다"]):
            gesture = "NOD"
            if emotion == "HAPPY":
                satisfaction_delta = 0.3
        elif any(word in text_lower for word in ["아니", "안 됩니다", "거절"]):
            gesture = "SHAKE_HEAD"
            satisfaction_delta = -0.2
        elif any(word in text_lower for word in ["여기", "이것", "저것", "이거", "주문번호"]):
            gesture = "POINT"
        
        # ⭐ Gemini 제안: 상황별 키워드 추출 (키워드 기반 분석)
        context_keywords = []
        if any(word in text_lower for word in ["주문번호", "order number", "주문 번호"]):
            context_keywords.append("order_number")
        if any(word in text_lower for word in ["해결", "완료", "감사", "solution"]):
            if satisfaction_delta > 0.3:
                context_keywords.append("solution_accepted")
        if any(word in text_lower for word in ["거절", "불가", "안 됩니다"]):
            if emotion == "ANGRY":
                context_keywords.append("policy_denial")
        
        return {
            "emotion": emotion,
            "gesture": gesture,
            "urgency": urgency,
            "satisfaction_delta": satisfaction_delta,
            "context_keywords": context_keywords,  # ⭐ 추가
            "confidence": 0.6  # 키워드 기반 분석은 낮은 신뢰도
        }
    
    except Exception as e:
        print(f"텍스트 분석 오류: {e}")
        return {
            "emotion": "NEUTRAL", 
            "gesture": "NONE", 
            "urgency": "LOW",
            "satisfaction_delta": 0.0,
            "context_keywords": [],  # ⭐ 추가
            "confidence": 0.5
        }


def get_video_path_by_avatar(gender: str, emotion: str, is_speaking: bool = False, 
                             gesture: str = "NONE", context_keywords: List[str] = None) -> str:
    """
    고객 아바타 정보(성별, 감정 상태, 제스처, 상황)에 따라 적절한 비디오 경로를 반환합니다.
    OpenAI/Gemini 기반 영상 RAG: LLM이 분석한 감정/제스처에 따라 비디오 클립을 선택합니다.
    
    ⭐ Gemini 제안: 상황별 비디오 클립 패턴 확장 (예: male_asking_order_number.mp4)
    
    Args:
        gender: "male" 또는 "female"
        emotion: "NEUTRAL", "HAPPY", "ANGRY", "ASKING", "SAD", "HOLD"
        is_speaking: 말하는 중인지 여부
        gesture: "NONE", "HAND_WAVE", "NOD", "SHAKE_HEAD", "POINT"
        context_keywords: 상황별 키워드 리스트 (예: ["order_number", "solution_accepted", "policy_denial"])
    
    Returns:
        비디오 파일 경로 (없으면 None)
    """
    # 비디오 디렉토리 경로 (사용자가 설정한 비디오 파일들이 저장된 위치)
    video_base_dir = os.path.join(DATA_DIR, "videos")
    os.makedirs(video_base_dir, exist_ok=True)
    
    # ⭐ Gemini 제안: 우선순위 -1 - 데이터베이스 기반 추천 비디오 (가장 우선)
    if context_keywords:
        db_recommended = get_recommended_video_from_database(emotion, gesture, context_keywords)
        if db_recommended:
            return db_recommended
    else:
        db_recommended = get_recommended_video_from_database(emotion, gesture, [])
        if db_recommended:
            return db_recommended
    
    # ⭐ Gemini 제안: 우선순위 0 - 상황별 비디오 클립 (가장 구체적)
    if context_keywords:
        for keyword in context_keywords:
            # 상황별 파일명 패턴 시도 (예: male_asking_order_number.mp4)
            context_filename = f"{gender}_{emotion.lower()}_{keyword}"
            if is_speaking:
                context_filename += "_speaking"
            context_filename += ".mp4"
            context_path = os.path.join(video_base_dir, context_filename)
            if os.path.exists(context_path):
                return context_path
            
            # 세션 상태에서도 확인
            context_video_key = f"video_{gender}_{emotion.lower()}_{keyword}"
            if context_video_key in st.session_state and st.session_state[context_video_key]:
                video_path = st.session_state[context_video_key]
                if os.path.exists(video_path):
                    return video_path
    
    # 우선순위 1: 제스처가 있는 경우 제스처별 비디오 시도
    if gesture != "NONE" and gesture:
        gesture_video_key = f"video_{gender}_{emotion.lower()}_{gesture.lower()}"
        if gesture_video_key in st.session_state and st.session_state[gesture_video_key]:
            video_path = st.session_state[gesture_video_key]
            if os.path.exists(video_path):
                return video_path
        
        # 제스처별 파일명 패턴 시도
        gesture_filename = f"{gender}_{emotion.lower()}_{gesture.lower()}"
        if is_speaking:
            gesture_filename += "_speaking"
        gesture_filename += ".mp4"
        gesture_path = os.path.join(video_base_dir, gesture_filename)
        if os.path.exists(gesture_path):
            return gesture_path
    
    # 우선순위 2: 감정 상태별 비디오 (제스처 없이)
    video_key = f"video_{gender}_{emotion.lower()}"
    if is_speaking:
        video_key += "_speaking"
    
    # 세션 상태에 저장된 비디오 경로가 있으면 사용
    if video_key in st.session_state and st.session_state[video_key]:
        video_path = st.session_state[video_key]
        if os.path.exists(video_path):
            return video_path
    
    # 기본 비디오 파일명 패턴 시도
    video_filename = f"{gender}_{emotion.lower()}"
    if is_speaking:
        video_filename += "_speaking"
    video_filename += ".mp4"
    
    video_path = os.path.join(video_base_dir, video_filename)
    if os.path.exists(video_path):
        return video_path
    
    # 우선순위 3: 기본 비디오 파일 시도 (중립 상태)
    default_video = os.path.join(video_base_dir, f"{gender}_neutral.mp4")
    if os.path.exists(default_video):
        return default_video
    
    # 우선순위 4: 세션 상태에서 업로드된 비디오 확인
    if "current_customer_video" in st.session_state and st.session_state.current_customer_video:
        return st.session_state.current_customer_video
    
    return None


# ⭐ Gemini 제안: 비디오 매핑 데이터베이스 관리 함수
def load_video_mapping_database() -> Dict[str, Any]:
    """비디오 매핑 데이터베이스를 로드합니다."""
    if os.path.exists(VIDEO_MAPPING_DB_FILE):
        try:
            with open(VIDEO_MAPPING_DB_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            print(f"비디오 매핑 데이터베이스 로드 오류: {e}")
            return {"mappings": [], "feedback_history": []}
    return {"mappings": [], "feedback_history": []}


def save_video_mapping_database(db_data: Dict[str, Any]):
    """비디오 매핑 데이터베이스를 저장합니다."""
    try:
        with open(VIDEO_MAPPING_DB_FILE, 'w', encoding='utf-8') as f:
            json.dump(db_data, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"비디오 매핑 데이터베이스 저장 오류: {e}")


def add_video_mapping_feedback(
    customer_text: str,
    selected_video_path: str,
    emotion: str,
    gesture: str,
    context_keywords: List[str],
    user_rating: int,  # 1-5 점수
    user_comment: str = ""
) -> None:
    """
    ⭐ Gemini 제안: 사용자 피드백을 비디오 매핑 데이터베이스에 추가합니다.
    
    Args:
        customer_text: 고객의 텍스트
        selected_video_path: 선택된 비디오 경로
        emotion: 분석된 감정
        gesture: 분석된 제스처
        context_keywords: 상황별 키워드
        user_rating: 사용자 평가 (1-5)
        user_comment: 사용자 코멘트 (선택적)
    """
    db_data = load_video_mapping_database()
    
    feedback_entry = {
        "timestamp": datetime.now().isoformat(),
        "customer_text": customer_text[:200],  # 최대 200자
        "selected_video": os.path.basename(selected_video_path) if selected_video_path else None,
        "video_path": selected_video_path,
        "emotion": emotion,
        "gesture": gesture,
        "context_keywords": context_keywords,
        "user_rating": user_rating,
        "user_comment": user_comment[:500] if user_comment else "",  # 최대 500자
        "is_natural_match": user_rating >= 4  # 4점 이상이면 자연스러운 매칭으로 간주
    }
    
    db_data["feedback_history"].append(feedback_entry)
    
    # 매핑 규칙 업데이트 (평가가 높은 경우)
    if user_rating >= 4:
        mapping_key = f"{emotion}_{gesture}_{'_'.join(context_keywords) if context_keywords else 'none'}"
        
        # 기존 매핑 찾기
        existing_mapping = None
        for mapping in db_data["mappings"]:
            if mapping.get("key") == mapping_key:
                existing_mapping = mapping
                break
        
        if existing_mapping:
            # 기존 매핑 업데이트 (평균 점수 계산)
            total_rating = existing_mapping.get("total_rating", 0) + user_rating
            count = existing_mapping.get("count", 0) + 1
            existing_mapping["total_rating"] = total_rating
            existing_mapping["count"] = count
            existing_mapping["avg_rating"] = total_rating / count
            existing_mapping["last_updated"] = datetime.now().isoformat()
        else:
            # 새 매핑 추가
            db_data["mappings"].append({
                "key": mapping_key,
                "emotion": emotion,
                "gesture": gesture,
                "context_keywords": context_keywords,
                "recommended_video": os.path.basename(selected_video_path) if selected_video_path else None,
                "video_path": selected_video_path,
                "total_rating": user_rating,
                "count": 1,
                "avg_rating": float(user_rating),
                "created_at": datetime.now().isoformat(),
                "last_updated": datetime.now().isoformat()
            })
    
    save_video_mapping_database(db_data)


def get_recommended_video_from_database(
    emotion: str,
    gesture: str,
    context_keywords: List[str]
) -> str:
    """
    ⭐ Gemini 제안: 데이터베이스에서 추천 비디오 경로를 가져옵니다.
    
    Args:
        emotion: 감정 상태
        gesture: 제스처
        context_keywords: 상황별 키워드
    
    Returns:
        추천 비디오 경로 (없으면 None)
    """
    db_data = load_video_mapping_database()
    
    mapping_key = f"{emotion}_{gesture}_{'_'.join(context_keywords) if context_keywords else 'none'}"
    
    # 정확한 매칭 찾기
    for mapping in db_data["mappings"]:
        if mapping.get("key") == mapping_key and mapping.get("avg_rating", 0) >= 4.0:
            video_path = mapping.get("video_path")
            if video_path and os.path.exists(video_path):
                return video_path
    
    # 부분 매칭 시도 (감정과 제스처만)
    partial_key = f"{emotion}_{gesture}_none"
    for mapping in db_data["mappings"]:
        if mapping.get("key") == partial_key and mapping.get("avg_rating", 0) >= 4.0:
            video_path = mapping.get("video_path")
            if video_path and os.path.exists(video_path):
                return video_path
    
    return None


def render_synchronized_video(text: str, audio_bytes: bytes, gender: str, emotion: str, 
                               role: str = "customer", autoplay: bool = True,
                               gesture: str = "NONE", context_keywords: List[str] = None):
    """
    TTS 오디오와 동기화된 비디오를 렌더링합니다.
    
    ⭐ Gemini 제안: 피드백 평가 기능 추가
    
    Args:
        text: 말하는 텍스트 내용
        audio_bytes: TTS로 생성된 오디오 바이트
        gender: 고객 성별 ("male" 또는 "female")
        emotion: 감정 상태 ("NEUTRAL", "HAPPY", "ANGRY", "ASKING", "SAD", "HOLD")
        role: 역할 ("customer" 또는 "agent")
        autoplay: 자동 재생 여부
        gesture: 제스처 (선택적)
        context_keywords: 상황별 키워드 (선택적)
    """
    if role == "customer":
        is_speaking = True
        if context_keywords is None:
            context_keywords = []
        
        # ⭐ Gemini 제안: 데이터베이스 기반 추천 비디오 우선 사용
        video_path = get_video_path_by_avatar(gender, emotion, is_speaking, gesture, context_keywords)
        
        if video_path and os.path.exists(video_path):
            try:
                with open(video_path, "rb") as f:
                    video_bytes = f.read()
                
                # 비디오와 오디오를 함께 재생
                # Streamlit의 st.video는 오디오 트랙이 있는 비디오를 지원합니다
                # 여기서는 비디오만 표시하고, 오디오는 별도로 재생합니다
                st.video(video_bytes, format="video/mp4", autoplay=autoplay, loop=False, muted=False)
                
                # 오디오도 함께 재생 (동기화)
                if audio_bytes:
                    st.audio(audio_bytes, format="audio/mp3", autoplay=autoplay, loop=False)
                
                # ⭐ Gemini 제안: 사용자 피드백 평가 UI 추가 (채팅/이메일 탭용)
                if not autoplay:  # 자동 재생이 아닌 경우에만 피드백 UI 표시
                    st.markdown("---")
                    st.markdown("**💬 비디오 매칭 평가**")
                    st.caption("이 비디오가 고객의 텍스트와 감정에 자연스럽게 매칭되었습니까?")
                    
                    feedback_key = f"video_feedback_chat_{st.session_state.get('sim_instance_id', 'default')}_{hash(text) % 10000}"
                    
                    col_rating, col_comment = st.columns([2, 3])
                    with col_rating:
                        rating = st.slider(
                            "평가 점수 (1-5점)",
                            min_value=1,
                            max_value=5,
                            value=3,
                            key=f"{feedback_key}_rating",
                            help="1점: 매우 부자연스러움, 5점: 매우 자연스러움"
                        )
                    
                    with col_comment:
                        comment = st.text_input(
                            "의견 (선택사항)",
                            key=f"{feedback_key}_comment",
                            placeholder="예: 비디오가 텍스트와 잘 맞았습니다"
                        )
                    
                    if st.button("피드백 제출", key=f"{feedback_key}_submit"):
                        # 피드백을 데이터베이스에 저장
                        add_video_mapping_feedback(
                            customer_text=text[:200],
                            selected_video_path=video_path,
                            emotion=emotion,
                            gesture=gesture,
                            context_keywords=context_keywords,
                            user_rating=rating,
                            user_comment=comment
                        )
                        st.success(f"✅ 피드백이 저장되었습니다! (점수: {rating}/5)")
                        st.info("💡 이 피드백은 향후 비디오 선택 정확도를 개선하는 데 사용됩니다.")
                
                return True
            except Exception as e:
                st.warning(f"비디오 재생 오류: {e}")
                # 비디오 재생 실패 시 오디오만 재생
                if audio_bytes:
                    st.audio(audio_bytes, format="audio/mp3", autoplay=autoplay, loop=False)
                return False
        else:
            # 비디오가 없으면 오디오만 재생
            if audio_bytes:
                st.audio(audio_bytes, format="audio/mp3", autoplay=autoplay, loop=False)
            return False
    else:
        # 에이전트는 비디오 없이 오디오만 재생
        if audio_bytes:
            st.audio(audio_bytes, format="audio/mp3", autoplay=autoplay, loop=False)
        return False


def generate_virtual_human_video(text: str, audio_bytes: bytes, gender: str, emotion: str, 
                                 provider: str = "hyperclova") -> bytes:
    """
    가상 휴먼 기술을 사용하여 텍스트와 오디오에 맞는 비디오를 생성합니다.
    
    ⚠️ 주의: OpenAI/Gemini API만으로는 입모양 동기화 비디오 생성이 불가능합니다.
    가상 휴먼 비디오 생성은 별도의 가상 휴먼 API (예: Hyperclova)가 필요합니다.
    
    현재는 미리 준비된 비디오 파일을 사용하는 방식을 권장합니다.
    
    Args:
        text: 말하는 텍스트 내용
        audio_bytes: TTS로 생성된 오디오 바이트
        gender: 고객 성별 ("male" 또는 "female")
        emotion: 감정 상태 ("NEUTRAL", "HAPPY", "ANGRY", "ASKING", "SAD", "HOLD")
        provider: 가상 휴먼 제공자 ("hyperclova", "other")
    
    Returns:
        생성된 비디오 바이트 (없으면 None)
    """
    # 가상 휴먼 API 키 확인
    if provider == "hyperclova":
        api_key = get_api_key("hyperclova")
        if not api_key:
            return None
        
        # TODO: Hyperclova API 연동 구현 (별도 API 필요)
        # OpenAI/Gemini API만으로는 불가능하므로, 실제 가상 휴먼 API가 필요합니다.
        # 예시 구조:
        # response = requests.post(
        #     "https://api.hyperclova.com/virtual-human/generate",
        #     headers={"Authorization": f"Bearer {api_key}"},
        #     json={
        #         "text": text,
        #         "audio": base64.b64encode(audio_bytes).decode(),
        #         "gender": gender,
        #         "emotion": emotion
        #     }
        # )
        # return response.content
    
    # 다른 제공자도 여기에 추가 가능
    # elif provider == "other":
    #     ...
    
    return None


def get_virtual_human_config() -> Dict[str, Any]:
    """
    가상 휴먼 설정을 반환합니다.
    
    Returns:
        가상 휴먼 설정 딕셔너리
    """
    return {
        "enabled": st.session_state.get("virtual_human_enabled", False),
        "provider": st.session_state.get("virtual_human_provider", "hyperclova"),
        "api_key": get_api_key("hyperclova") if st.session_state.get("virtual_human_provider", "hyperclova") == "hyperclova" else None
    }


# 역할별 TTS 음성 스타일 설정
TTS_VOICES = {
    "customer_male": {
        "gender": "male",
        "voice": "alloy"  # Male voice
    },
    "customer_female": {
        "gender": "female",
        "voice": "nova"  # Female voice
    },
    "customer": {
        "gender": "male",
        "voice": "alloy"  # Default male voice (fallback)
    },
    "agent": {
        "gender": "female",
        "voice": "shimmer"  # Distinct Female, Professional/Agent
    },
    "supervisor": {
        "gender": "female",
        "voice": "nova"  # Another Distinct Female, Informative/Supervisor
    }
}


def synthesize_tts(text: str, lang_key: str, role: str = "agent"):
    # lang_key 검증 및 기본값 처리
    if not lang_key or lang_key not in ["ko", "en", "ja"]:
        lang_key = st.session_state.get("language", "ko")
        if lang_key not in ["ko", "en", "ja"]:
            lang_key = "ko"  # 최종 기본값
    
    L = LANG.get(lang_key, LANG["ko"])  # 안전한 접근
    client = st.session_state.openai_client
    if client is None:
        return None, L.get("openai_missing", "OpenAI API Key가 필요합니다.")

    # ⭐ 수정: 고객 역할인 경우 성별에 따라 음성 선택
    if role == "customer":
        customer_gender = st.session_state.customer_avatar.get("gender", "male")
        if customer_gender == "female":
            voice_key = "customer_female"
        else:
            voice_key = "customer_male"
        
        if voice_key in TTS_VOICES:
            voice_name = TTS_VOICES[voice_key]["voice"]
        else:
            voice_name = TTS_VOICES["customer"]["voice"]  # Fallback
    elif role in TTS_VOICES:
        voice_name = TTS_VOICES[role]["voice"]
    else:
        voice_name = TTS_VOICES["agent"]["voice"]  # Default fallback

    try:
        # ⭐ 수정: 텍스트 길이 제한을 제거하여 전체 문의가 재생되도록 함
        # OpenAI TTS는 최대 4096자를 지원하지만, 실제로는 더 긴 텍스트도 처리 가능
        # 고객의 문의를 끝까지 다 들어야 원활한 응대가 가능하므로 전체 텍스트를 처리
        # 만약 텍스트가 너무 길면 (예: 10000자 이상) 여러 청크로 나눠서 처리할 수 있지만,
        # 일반적인 고객 문의는 4096자 이내이므로 전체를 처리
        
        # tts-1 모델 사용 (안정성)
        resp = client.audio.speech.create(
            model="tts-1",
            voice=voice_name,
            input=text
            # format="mp3"은 기본값입니다.
        )
        return resp.read(), L["tts_status_success"]

    except Exception as e:
        return None, f"{L['tts_status_error']}: {e}"


# ----------------------------------------
# TTS Helper
# ----------------------------------------

def render_tts_button(text, lang_key, role="customer", prefix="", index: int = -1):
    # lang_key 검증 및 기본값 처리
    if not lang_key or lang_key not in ["ko", "en", "ja"]:
        lang_key = st.session_state.get("language", "ko")
        if lang_key not in ["ko", "en", "ja"]:
            lang_key = "ko"  # 최종 기본값
    
    L = LANG.get(lang_key, LANG["ko"])  # 안전한 접근

    # ⭐ 수정: index=-1인 경우, UUID를 사용하여 safe_key 생성
    if index == -1:
        # 이관 요약처럼 인덱스가 고정되지 않는 경우, 텍스트 해시와 세션 인스턴스 ID를 조합
        content_hash = hashlib.md5(text[:100].encode()).hexdigest()
        session_id_part = st.session_state.get('sim_instance_id', 'default_session')
        # ⭐ 수정: 이관 요약의 경우 안정적인 키를 생성 (time.time_ns() 제거하여 매번 같은 키 생성)
        # 언어 코드도 추가하여 이관 후 언어 변경 시에도 고유성 보장
        lang_code = st.session_state.get('language', lang_key)
        safe_key = f"{prefix}_SUMMARY_{session_id_part}_{lang_code}_{content_hash}"
    else:
        # 대화 로그처럼 인덱스가 존재하는 경우 (기존 로직 유지)
        content_hash = hashlib.md5(text[:100].encode()).hexdigest()
        safe_key = f"{prefix}_{index}_{content_hash}"

    # 재생 버튼을 누를 때만 TTS 요청
    if st.button(L["button_listen_audio"], key=safe_key):
        if not st.session_state.openai_client:
            st.error(L["openai_missing"])
            return  # 키 없으면 종료

        with st.spinner(L["tts_status_generating"]):
            try:
                audio_bytes, msg = synthesize_tts(text, lang_key, role=role)
                if audio_bytes:
                    # ⭐ st.audio 호출 시 성공한 경우에만 재생 시간을 확보
                    # Streamlit 문서: autoplay는 브라우저 정책상 사용자 상호작용 없이는 작동하지 않을 수 있음
                    try:
                        st.audio(audio_bytes, format="audio/mp3", autoplay=True, loop=False)
                        st.success(msg)
                        # ⭐ 수정: 재생이 시작될 충분한 시간을 확보하기 위해 대기 시간을 3초로 늘림
                        time.sleep(3)
                    except Exception as e:
                        st.warning(f"오디오 재생 중 오류: {e}. 오디오 파일은 생성되었지만 자동 재생에 실패했습니다.")
                        st.audio(audio_bytes, format="audio/mp3", autoplay=False)
                        st.success(msg)
                else:
                    st.error(msg)
                    time.sleep(1)  # 에러 발생 시도 잠시 대기
            except Exception as e:
                # TTS API 호출 자체에서 예외 발생 시 (네트워크 등)
                st.error(f"❌ TTS 생성 중 치명적인 오류 발생: {e}")
                time.sleep(1)

            # 버튼 클릭 이벤트 후, 불필요한 재실행을 막기 위해 여기서 함수 종료
            return
        # [중략: TTS Helper 끝]


# ========================================
# 4. 로컬 음성 기록 Helper
# ========================================

def load_voice_records() -> List[Dict[str, Any]]:
    return _load_json(VOICE_META_FILE, [])


def save_voice_records(records: List[Dict[str, Any]]):
    _save_json(VOICE_META_FILE, records)


def save_audio_record_local(
        audio_bytes: bytes,
        filename: str,
        transcript_text: str,
        mime_type: str = "audio/webm",
        meta: Dict[str, Any] = None,
) -> str:
    records = load_voice_records()
    rec_id = str(uuid.uuid4())
    ts = datetime.utcnow().isoformat()

    ext = filename.split(".")[-1] if "." in filename else "webm"
    audio_filename = f"{rec_id}.{ext}"
    audio_path = os.path.join(AUDIO_DIR, audio_filename)
    with open(audio_path, "wb") as f:
        f.write(audio_bytes)

    rec = {
        "id": rec_id,
        "created_at": ts,
        "filename": filename,
        "audio_filename": audio_filename,
        "size": len(audio_bytes),
        "transcript": transcript_text,
        "mime_type": mime_type,
        "language": st.session_state.language,
        "meta": meta or {},
    }
    records.insert(0, rec)
    save_voice_records(records)
    return rec_id


def delete_audio_record_local(rec_id: str) -> bool:
    records = load_voice_records()
    idx = next((i for i, r in enumerate(records) if r.get("id") == rec_id), None)
    if idx is None:
        return False
    rec = records.pop(idx)
    audio_filename = rec.get("audio_filename")
    if audio_filename:
        audio_path = os.path.join(AUDIO_DIR, audio_filename)
        try:
            os.remove(audio_path)
        except FileNotFoundError:
            pass
    save_voice_records(records)
    return True


def get_audio_bytes_local(rec_id: str):
    records = load_voice_records()
    rec = next((r for r in records if r.get("id") == rec_id), None)
    if not rec:
        raise FileNotFoundError("record not found")
    audio_filename = rec["audio_filename"]
    audio_path = os.path.join(AUDIO_DIR, audio_filename)
    with open(audio_path, "rb") as f:
        b = f.read()
    return b, rec


# ========================================
# 5. 로컬 시뮬레이션 이력 Helper (요청 4 반영)
# ========================================


def load_simulation_histories_local(lang_key: str) -> List[Dict[str, Any]]:
    histories = _load_json(SIM_META_FILE, [])
    # 현재 언어와 메시지 리스트가 유효한 이력만 필터링
    return [
        h for h in histories
        if h.get("language_key") == lang_key and (isinstance(h.get("messages"), list) or h.get("summary"))
    ]



def generate_chat_summary(messages: List[Dict[str, Any]], initial_query: str, customer_type: str,
                          current_lang_key: str) -> Dict[str, Any]:
    """채팅 내용을 AI로 요약하여 주요 정보와 점수를 추출 (요청 4)"""
    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}[current_lang_key]

    # 대화 내용 추출
    conversation_text = f"Initial Query: {initial_query}\n\n"
    for msg in messages:
        role = msg.get("role", "")
        content = msg.get("content", "")
        if role in ["customer", "customer_rebuttal", "phone_exchange"]:
            conversation_text += f"Customer: {content}\n"
        elif role == "agent_response" or role == "agent":
            conversation_text += f"Agent: {content}\n"
        # supervisor 메시지는 LLM에 전달하지 않아 역할 혼동 방지

    # 폰 교환 로그는 이미 "Agent: ... | Customer: ..." 형태로 기록되므로,
    # generate_summary_for_call 함수에서 별도로 처리할 필요 없이,
    # 여기서는 범용 채팅 요약 로직을 따르도록 메시지를 정제합니다.

    summary_prompt = f"""
You are an AI analyst summarizing a customer support conversation. Your task is to extract comprehensive customer profile data and score various aspects numerically.

Analyze the conversation and provide a structured summary in JSON format (ONLY JSON, no markdown).

Extract and score:
1. Main inquiry topic (what the customer asked about)
2. Key responses provided by the agent (list of max 3 core actions/solutions)
3. Customer sentiment score (0-100, where 0=very negative, 50=neutral, 100=very positive)
4. Customer satisfaction score (0-100, based on final response)
5. Customer characteristics with detailed scoring:
   - Language preference (detected language code: ko/en/ja)
   - Cultural background hints (score 0-100, where higher = more cultural context detected)
   - Location/region (general region only, anonymize specific addresses)
   - Communication style (formal/casual, brief/detailed) with scores:
     * Formality score (0-100, 0=casual, 100=very formal)
     * Detail level score (0-100, 0=brief, 100=very detailed)
   - Customer personality traits (score each 0-100):
     * Patience level (0-100)
     * Assertiveness (0-100)
     * Politeness level (0-100)
     * Technical proficiency (0-100, if technical inquiry)
6. Privacy-sensitive information (anonymize: names, emails, phone numbers, specific addresses)
   - Extract patterns only (e.g., "email provided", "phone number provided", "resides in Asia region")
7. Customer behavior patterns:
   - Response time pattern (fast/moderate/slow based on message frequency)
   - Question complexity (simple/moderate/complex)
   - Escalation tendency (0-100, likelihood to escalate)

Output format (JSON only):
{{
  "main_inquiry": "brief description of main issue",
  "key_responses": ["response 1", "response 2"],
  "customer_sentiment_score": 75,
  "customer_satisfaction_score": 80,
  "customer_characteristics": {{
    "language": "ko/en/ja or unknown",
    "cultural_hints": "brief description or unknown",
    "cultural_score": 60,
    "region": "general region or unknown",
    "communication_style": "formal/casual/brief/detailed",
    "formality_score": 70,
    "detail_level_score": 65,
    "personality_traits": {{
      "patience_level": 60,
      "assertiveness": 70,
      "politeness_level": 80,
      "technical_proficiency": 50
    }}
  }},
  "privacy_info": {{
    "has_email": true/false,
    "has_phone": true/false,
    "has_address": true/false,
    "region_hint": "general region or unknown"
  }},
  "behavior_patterns": {{
    "response_time": "fast/moderate/slow",
    "question_complexity": "simple/moderate/complex",
    "escalation_tendency": 30
  }},
  "summary": "overall conversation summary in {lang_name}"
}}

Conversation:
{conversation_text}

JSON Output:
"""

    if not st.session_state.is_llm_ready:
        # Fallback summary with enhanced structure
        return {
            "main_inquiry": initial_query[:100],
            "key_responses": [],
            "customer_sentiment_score": 50,
            "customer_satisfaction_score": 50,
            "customer_characteristics": {
                "language": current_lang_key,
                "cultural_hints": "unknown",
                "cultural_score": 0,
                "region": "unknown",
                "communication_style": "unknown",
                "formality_score": 50,
                "detail_level_score": 50,
                "personality_traits": {
                    "patience_level": 50,
                    "assertiveness": 50,
                    "politeness_level": 50,
                    "technical_proficiency": 50
                }
            },
            "privacy_info": {
                "has_email": False,
                "has_phone": False,
                "has_address": False,
                "region_hint": "unknown"
            },
            "behavior_patterns": {
                "response_time": "moderate",
                "question_complexity": "moderate",
                "escalation_tendency": 50
            },
            "summary": f"Customer inquiry about: {initial_query[:100]}"
        }

    try:
        summary_text = run_llm(summary_prompt).strip()
        # JSON 추출 (마크다운 코드 블록 제거)
        if "```json" in summary_text:
            summary_text = summary_text.split("```json")[1].split("```")[0].strip()
        elif "```" in summary_text:
            summary_text = summary_text.split("```")[1].split("```")[0].strip()

        import json
        summary_data = json.loads(summary_text)
        return summary_data
    except Exception as e:
        # JSON 파싱 실패 시 fallback with enhanced structure
        st.warning(f"요약 생성 중 오류 발생: {e}")
        return {
            "main_inquiry": initial_query[:100],
            "key_responses": [],
            "customer_sentiment_score": 50,
            "customer_satisfaction_score": 50,
            "customer_characteristics": {
                "language": current_lang_key,
                "cultural_hints": "unknown",
                "cultural_score": 0,
                "region": "unknown",
                "communication_style": "unknown",
                "formality_score": 50,
                "detail_level_score": 50,
                "personality_traits": {
                    "patience_level": 50,
                    "assertiveness": 50,
                    "politeness_level": 50,
                    "technical_proficiency": 50
                }
            },
            "privacy_info": {
                "has_email": False,
                "has_phone": False,
                "has_address": False,
                "region_hint": "unknown"
            },
            "behavior_patterns": {
                "response_time": "moderate",
                "question_complexity": "moderate",
                "escalation_tendency": 50
            },
            "summary": f"Error generating summary: {str(e)}"
        }



def recommend_guideline_for_customer(new_customer_summary: Dict[str, Any], histories: List[Dict[str, Any]], language: str = "ko") -> Optional[str]:
    """
    신규 고객의 문의사항과 말투 등을 종합하여 고객 성향 점수를 수치화하고,
    저장된 데이터를 바탕으로 최적의 가이드라인을 추천합니다.
    
    Args:
        new_customer_summary: 신규 고객의 요약 데이터 (성향 점수 포함)
        histories: 과거 고객 응대 이력 리스트
        language: 언어 코드
    
    Returns:
        추천 가이드라인 텍스트 또는 None
    """
    if not histories or not get_api_key("gemini") and not get_api_key("openai"):
        return None
    
    try:
        # 유사한 고객 프로필 찾기
        similar_customers = []
        new_scores = {
            "sentiment": new_customer_summary.get("customer_sentiment_score", 50),
            "satisfaction": new_customer_summary.get("customer_satisfaction_score", 50),
            "formality": new_customer_summary.get("customer_characteristics", {}).get("formality_score", 50),
            "patience": new_customer_summary.get("customer_characteristics", {}).get("personality_traits", {}).get("patience_level", 50),
            "assertiveness": new_customer_summary.get("customer_characteristics", {}).get("personality_traits", {}).get("assertiveness", 50),
        }
        
        for h in histories:
            if not h.get("summary") or not isinstance(h.get("summary"), dict):
                continue
            
            summary = h["summary"]
            old_scores = {
                "sentiment": summary.get("customer_sentiment_score", 50),
                "satisfaction": summary.get("customer_satisfaction_score", 50),
                "formality": summary.get("customer_characteristics", {}).get("formality_score", 50),
                "patience": summary.get("customer_characteristics", {}).get("personality_traits", {}).get("patience_level", 50),
                "assertiveness": summary.get("customer_characteristics", {}).get("personality_traits", {}).get("assertiveness", 50),
            }
            
            # 유사도 계산 (점수 차이의 절대값 합)
            similarity = sum(abs(new_scores[k] - old_scores[k]) for k in new_scores.keys())
            
            if similarity < 100:  # 임계값 이하인 경우 유사 고객으로 간주
                similar_customers.append({
                    "history": h,
                    "similarity": similarity,
                    "scores": old_scores
                })
        
        # 유사도 순으로 정렬
        similar_customers.sort(key=lambda x: x["similarity"])
        
        # 상위 5개 유사 고객의 성공 사례 분석
        if similar_customers:
            lang_name = {"ko": "한국어", "en": "English", "ja": "日本語"}.get(language, "한국어")
            
            similar_cases_text = json.dumps([
                {
                    "initial_query": c["history"].get("initial_query", ""),
                    "key_responses": c["history"].get("summary", {}).get("key_responses", []),
                    "scores": c["scores"],
                    "satisfaction": c["history"].get("summary", {}).get("customer_satisfaction_score", 50)
                }
                for c in similar_customers[:5]
            ], ensure_ascii=False, indent=2)
            
            recommendation_prompt = (
                f"당신은 CS 센터 전문가입니다. 신규 고객의 성향 점수를 분석하고, 유사한 과거 고객들의 성공 사례를 바탕으로 최적의 응대 가이드라인을 추천하세요.\n\n"
                f"신규 고객 프로필:\n{json.dumps(new_customer_summary, ensure_ascii=False, indent=2)}\n\n"
                f"유사한 과거 고객 사례 (상위 5개):\n{similar_cases_text}\n\n"
                f"다음 내용을 포함하여 {lang_name}로 가이드라인을 작성하세요:\n"
                f"1. 고객 성향 분석 (점수 기반)\n"
                f"2. 예상되는 고객 반응 패턴\n"
                f"3. 효과적인 응대 전략 (유사 사례 기반)\n"
                f"4. 주의해야 할 사항\n"
                f"5. 권장 응대 톤 및 스타일\n\n"
                f"실용적이고 구체적인 가이드라인을 제공하세요."
            )
            
            recommendation = run_llm(recommendation_prompt)
            return recommendation if recommendation and not recommendation.startswith("❌") else None
        
        return None
        
    except Exception as e:
        print(f"가이드라인 추천 중 오류 발생: {e}")
        return None


def generate_daily_customer_guide(histories: List[Dict[str, Any]], language: str = "ko") -> Optional[str]:
    """
    일일 고객 가이드 생성 함수
    고객 응대 시뮬레이터 데이터를 분석하여 고객 가이드라인을 생성합니다.
    
    Args:
        histories: 고객 응대 이력 리스트
        language: 언어 코드 (ko, en, ja)
    
    Returns:
        생성된 가이드 내용 (문자열) 또는 None (실패 시)
    """
    if not histories or not get_api_key("gemini") and not get_api_key("openai"):
        return None
    
    try:
        # 요약 데이터가 있는 이력만 필터링
        histories_with_summary = [h for h in histories if h.get("summary") and isinstance(h.get("summary"), dict)]
        
        if not histories_with_summary:
            return None
        
        # 최근 50개 이력 사용 (동일 고객 데이터 누적을 위해 전체 이력 사용)
        recent_histories = histories_with_summary[:50]
        
        # 고객별로 데이터 그룹화 (동일 고객의 상황에 따라 데이터 누적)
        customer_data_map = {}
        for h in recent_histories:
            customer_id = h.get("id", "")
            customer_type = h.get("customer_type", "")
            summary = h.get("summary", {})
            
            if customer_id not in customer_data_map:
                customer_data_map[customer_id] = {
                    "customer_type": customer_type,
                    "histories": [],
                    "total_interactions": 0
                }
            
            customer_data_map[customer_id]["histories"].append({
                "initial_query": h.get("initial_query", ""),
                "summary": summary,
                "timestamp": h.get("timestamp", ""),
                "language": h.get("language_key", language)
            })
            customer_data_map[customer_id]["total_interactions"] += 1
        
        # LLM을 사용하여 고객 가이드 생성
        lang_name = {"ko": "한국어", "en": "English", "ja": "日本語"}.get(language, "한국어")
        
        guide_prompt = (
            f"당신은 CS 센터 교육 전문가입니다. 다음 고객 응대 이력 데이터를 분석하여 종합적인 고객 응대 가이드라인을 작성하세요.\n\n"
            f"분석할 이력 데이터 (고객별 누적 데이터 포함):\n{json.dumps(list(customer_data_map.values())[:20], ensure_ascii=False, indent=2)}\n\n"
            f"다음 내용을 포함하여 가이드라인을 {lang_name}로 작성하세요:\n"
            f"1. 고객 유형별 응대 전략 (일반/까다로운/매우 불만족)\n"
            f"2. 문화권별 응대 가이드 (언어, 문화적 배경 고려)\n"
            f"3. 주요 문의 유형별 해결 방법\n"
            f"4. 고객 감정 점수에 따른 응대 전략\n"
            f"5. 개인정보 처리 가이드\n"
            f"6. 효과적인 소통 스타일 권장사항\n"
            f"7. 동일 고객의 반복 문의에 대한 대응 전략\n"
            f"8. 강성 고객 가이드라인 (까다로운 고객, 매우 불만족 고객)\n\n"
            f"가이드라인을 {lang_name}로 작성하세요. 실제 사례를 바탕으로 구체적이고 실용적인 내용으로 작성해주세요."
        )
        
        guide_content = run_llm(guide_prompt)
        
        if not guide_content or guide_content.startswith("❌"):
            return None
        
        # 가이드 내용 포맷팅
        today_str = datetime.now().strftime("%y%m%d")
        formatted_guide = (
            f"고객 응대 가이드라인\n"
            f"생성일: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
            f"분석 이력 수: {len(recent_histories)}\n"
            f"고객 수: {len(customer_data_map)}\n"
            f"=" * 80 + "\n\n"
            f"{guide_content}\n\n"
            f"=" * 80 + "\n"
            f"이 가이드는 AI 고객 응대 시뮬레이터 데이터를 기반으로 자동 생성되었습니다.\n"
            f"고객 데이터가 추가될 때마다 업데이트됩니다."
        )
        
        return formatted_guide
        
    except Exception as e:
        print(f"고객 가이드 생성 중 오류 발생: {e}")
        return None


def save_daily_customer_guide(guide_content: str, language: str = "ko") -> Optional[str]:
    """
    일일 고객 가이드를 파일로 저장합니다.
    
    Args:
        guide_content: 가이드 내용
        language: 언어 코드
    
    Returns:
        저장된 파일 경로 또는 None
    """
    try:
        today_str = datetime.now().strftime("%y%m%d")
        guide_filename = f"{today_str}_고객가이드.TXT"
        guide_filepath = os.path.join(DATA_DIR, guide_filename)
        
        # 기존 파일이 있으면 추가 (동일 고객 데이터 누적)
        if os.path.exists(guide_filepath):
            with open(guide_filepath, "r", encoding="utf-8") as f:
                existing_content = f.read()
            
            # 새 내용을 기존 내용에 추가
            updated_content = (
                f"{existing_content}\n\n"
                f"{'=' * 80}\n"
                f"업데이트: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
                f"{'=' * 80}\n\n"
                f"{guide_content}"
            )
            
            with open(guide_filepath, "w", encoding="utf-8") as f:
                f.write(updated_content)
        else:
            # 새 파일 생성
            with open(guide_filepath, "w", encoding="utf-8") as f:
                f.write(guide_content)
        
        return guide_filepath
        
    except Exception as e:
        print(f"고객 가이드 저장 중 오류 발생: {e}")
        return None


def get_daily_data_statistics(language: str = "ko") -> Dict[str, Any]:
    """
    일일 데이터 수집 통계를 반환합니다.
    하루에 최소 1개 이상씩 최소 5인의 데이터 확보를 추적합니다.
    
    Returns:
        통계 딕셔너리
    """
    histories = _load_json(SIM_META_FILE, [])
    today = datetime.now().date()
    
    # 오늘 날짜의 이력 필터링
    today_histories = []
    for h in histories:
        try:
            ts = datetime.fromisoformat(h.get("timestamp", "")).date()
            if ts == today and h.get("summary") and isinstance(h.get("summary"), dict):
                today_histories.append(h)
        except:
            continue
    
    # 고유 고객 수 계산 (이메일/전화번호 기반 또는 ID 기반)
    unique_customers = set()
    for h in today_histories:
        # 고객 ID 또는 초기 문의를 기반으로 고유 고객 식별
        customer_id = h.get("id", "")
        initial_query = h.get("initial_query", "")
        # 간단한 해시 기반 고유 고객 식별
        customer_hash = hashlib.md5(f"{customer_id}_{initial_query[:50]}".encode()).hexdigest()
        unique_customers.add(customer_hash)
    
    return {
        "date": today.isoformat(),
        "total_cases": len(today_histories),
        "unique_customers": len(unique_customers),
        "target_met": len(unique_customers) >= 5,  # 최소 5인 목표 달성 여부
        "cases_with_summary": len([h for h in today_histories if h.get("summary")])
    }


def save_simulation_history_local(initial_query: str, customer_type: str, messages: List[Dict[str, Any]],
                                  is_chat_ended: bool, attachment_context: str, is_call: bool = False):
    """AI 요약 데이터를 중심으로 이력을 저장 (요청 4 반영)"""
    histories = _load_json(SIM_META_FILE, [])
    doc_id = str(uuid.uuid4())
    ts = datetime.utcnow().isoformat()

    # AI 요약 생성 (채팅 종료 시 또는 충분한 대화가 있을 때)
    summary_data = None
    if is_chat_ended or len(messages) > 4 or is_call:  # 전화 통화는 바로 요약 시도
        summary_data = generate_chat_summary(messages, initial_query, customer_type, st.session_state.language)

    # 요약 데이터가 생성된 경우에만 저장 (요약 중심 저장)
    if summary_data:
        # 요약 데이터에 초기 문의와 핵심 정보 포함
        data = {
            "id": doc_id,
            "initial_query": initial_query,  # 초기 문의는 유지
            "customer_type": customer_type,
            "messages": [],  # 전체 메시지는 저장하지 않음 (요약만 저장)
            "summary": summary_data,  # AI 요약 데이터 (주요 저장 내용)
            "language_key": st.session_state.language,
            "timestamp": ts,
            "is_chat_ended": is_chat_ended,
            "attachment_context": attachment_context if attachment_context else "",  # 첨부 파일 컨텍스트
            "is_call": is_call,  # 전화 여부 플래그
        }
    else:
        # 요약이 아직 생성되지 않은 경우 (진행 중인 대화), 최소한의 정보만 저장
        data = {
            "id": doc_id,
            "initial_query": initial_query,
            "customer_type": customer_type,
            "messages": messages[:10] if len(messages) > 10 else messages,  # 최근 10개만 저장
            "summary": None,  # 요약 없음
            "language_key": st.session_state.language,
            "timestamp": ts,
            "is_chat_ended": is_chat_ended,
            "attachment_context": attachment_context if attachment_context else "",
            "is_call": is_call,
        }

    # 기존 이력에 추가 (최신순)
    histories.insert(0, data)
    # 너무 많은 이력 방지 (예: 100개로 증가 - 요약만 저장하므로 용량 부담 적음)
    _save_json(SIM_META_FILE, histories[:100])
    
    # ⭐ 추가: 고객 데이터가 저장될 때마다 자동으로 일일 고객 가이드 생성
    if summary_data and is_chat_ended:  # 채팅이 종료되고 요약이 생성된 경우에만 가이드 생성
        try:
            # 전체 이력 로드 (동일 고객 데이터 누적을 위해)
            all_histories = _load_json(SIM_META_FILE, [])
            
            # 오늘 날짜의 가이드가 이미 생성되었는지 확인
            today_str = datetime.now().strftime("%y%m%d")
            guide_filename = f"{today_str}_고객가이드.TXT"
            guide_filepath = os.path.join(DATA_DIR, guide_filename)
            
            # 가이드 생성 (기존 파일이 있으면 업데이트, 없으면 새로 생성)
            guide_content = generate_daily_customer_guide(all_histories, st.session_state.language)
            
            if guide_content:
                saved_path = save_daily_customer_guide(guide_content, st.session_state.language)
                if saved_path:
                    print(f"✅ 고객 가이드가 자동 생성/업데이트되었습니다: {saved_path}")
        except Exception as e:
            print(f"고객 가이드 자동 생성 중 오류 발생 (무시됨): {e}")
    
    return doc_id


def delete_all_history_local():
    _save_json(SIM_META_FILE, [])


# ========================================
# DB 저장 기능 (Word/PPTX/PDF)
# ========================================

def export_history_to_word(histories: List[Dict[str, Any]], filename: str = None, lang: str = "ko") -> str:
    """이력을 Word 파일로 저장"""
    if not IS_DOCX_AVAILABLE:
        raise ImportError("Word 저장을 위해 python-docx가 필요합니다: pip install python-docx")
    
    # 언어 설정 확인 및 기본값 설정
    if lang not in ["ko", "en", "ja"]:
        lang = "ko"
    L = LANG.get(lang, LANG["ko"])
    
    if filename is None:
        filename = f"customer_history_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
    filepath = os.path.join(DATA_DIR, filename)
    
    doc = Document()
    
    # 제목 추가
    title = doc.add_heading(L.get("download_history_title", "고객 응대 이력 요약"), 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # 각 이력 추가
    for i, hist in enumerate(histories, 1):
        # 이력 제목
        doc.add_heading(f'{L.get("download_history_number", "이력 #")}{i}', level=1)
        
        # 기본 정보
        doc.add_paragraph(f'ID: {hist.get("id", "N/A")}')
        doc.add_paragraph(f'{L.get("date_label", "날짜")}: {hist.get("timestamp", "N/A")}')
        doc.add_paragraph(f'{L.get("download_initial_inquiry", "초기 문의")}: {hist.get("initial_query", "N/A")}')
        doc.add_paragraph(f'{L.get("customer_type_label", "고객 유형")}: {hist.get("customer_type", "N/A")}')
        doc.add_paragraph(f'{L.get("language_label", "언어")}: {hist.get("language_key", "N/A")}')
        
        summary = hist.get('summary', {})
        if summary:
            # 요약 섹션
            doc.add_heading(L.get("download_summary", "요약"), level=2)
            doc.add_paragraph(f'{L.get("download_main_inquiry", "주요 문의")}: {summary.get("main_inquiry", "N/A")}')
            doc.add_paragraph(f'{L.get("download_key_response", "핵심 응답")}: {", ".join(summary.get("key_responses", []))}')
            doc.add_paragraph(f'{L.get("sentiment_score_label", "고객 감정 점수")}: {summary.get("customer_sentiment_score", "N/A")}/100')
            doc.add_paragraph(f'{L.get("customer_satisfaction_score_label", "고객 만족도 점수")}: {summary.get("customer_satisfaction_score", "N/A")}/100')
            
            # 고객 특성
            characteristics = summary.get('customer_characteristics', {})
            doc.add_heading(L.get("download_customer_characteristics", "고객 특성"), level=2)
            doc.add_paragraph(f'{L.get("language_label", "언어")}: {characteristics.get("language", "N/A")}')
            doc.add_paragraph(f'{L.get("download_cultural_background", "문화적 배경")}: {characteristics.get("cultural_hints", "N/A")}')
            doc.add_paragraph(f'{L.get("region_label", "지역")}: {characteristics.get("region", "N/A")}')
            doc.add_paragraph(f'{L.get("download_communication_style", "소통 스타일")}: {characteristics.get("communication_style", "N/A")}')
            
            # 개인정보 요약
            privacy = summary.get('privacy_info', {})
            doc.add_heading(L.get("download_privacy_summary", "개인정보 요약"), level=2)
            doc.add_paragraph(f'{L.get("email_provided_label", "이메일 제공")}: {L.get("download_yes", "예") if privacy.get("has_email") else L.get("download_no", "아니오")}')
            doc.add_paragraph(f'{L.get("phone_provided_label", "전화번호 제공")}: {L.get("download_yes", "예") if privacy.get("has_phone") else L.get("download_no", "아니오")}')
            doc.add_paragraph(f'{L.get("download_address_provided", "주소 제공")}: {L.get("download_yes", "예") if privacy.get("has_address") else L.get("download_no", "아니오")}')
            doc.add_paragraph(f'{L.get("download_region_hint", "지역 힌트")}: {privacy.get("region_hint", "N/A")}')
            
            # 전체 요약
            doc.add_paragraph(f'{L.get("download_overall_summary", "전체 요약")}: {summary.get("summary", "N/A")}')
        
        # 구분선
        if i < len(histories):
            doc.add_paragraph('-' * 80)
    
    doc.save(filepath)
    return filepath



def export_history_to_pptx(histories: List[Dict[str, Any]], filename: str = None, lang: str = "ko") -> str:
    """이력을 PPTX 파일로 저장"""
    if not IS_PPTX_AVAILABLE:
        raise ImportError("PPTX 저장을 위해 python-pptx가 필요합니다: pip install python-pptx")
    
    # 언어 설정 확인 및 기본값 설정
    if lang not in ["ko", "en", "ja"]:
        lang = "ko"
    L = LANG.get(lang, LANG["ko"])
    
    if filename is None:
        filename = f"customer_history_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    filepath = os.path.join(DATA_DIR, filename)
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 제목 슬라이드
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = L.get("download_history_title", "고객 응대 이력 요약")
    subtitle.text = f"{L.get('download_created_date', '생성일')}: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    # 각 이력에 대해 슬라이드 생성
    for i, hist in enumerate(histories, 1):
        # 제목 및 내용 레이아웃
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = f"{L.get('download_history_number', '이력 #')}{i}"
        
        tf = body_shape.text_frame
        tf.text = f"ID: {hist.get('id', 'N/A')}"
        
        p = tf.add_paragraph()
        p.text = f"{L.get('date_label', '날짜')}: {hist.get('timestamp', 'N/A')}"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = f"{L.get('download_initial_inquiry', '초기 문의')}: {hist.get('initial_query', 'N/A')}"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = f"{L.get('customer_type_label', '고객 유형')}: {hist.get('customer_type', 'N/A')}"
        p.level = 0
        
        summary = hist.get('summary', {})
        if summary:
            p = tf.add_paragraph()
            p.text = f"{L.get('download_main_inquiry', '주요 문의')}: {summary.get('main_inquiry', 'N/A')}"
            p.level = 0
            
            p = tf.add_paragraph()
            p.text = f"{L.get('sentiment_score_label', '고객 감정 점수')}: {summary.get('customer_sentiment_score', 'N/A')}/100"
            p.level = 0
            
            p = tf.add_paragraph()
            p.text = f"{L.get('customer_satisfaction_score_label', '고객 만족도 점수')}: {summary.get('customer_satisfaction_score', 'N/A')}/100"
            p.level = 0
    
    prs.save(filepath)
    return filepath



def export_history_to_pdf(histories: List[Dict[str, Any]], filename: str = None, lang: str = "ko") -> str:
    """이력을 PDF 파일로 저장 (한글/일본어 인코딩 지원 강화)"""
    if not IS_REPORTLAB_AVAILABLE:
        raise ImportError("PDF 저장을 위해 reportlab이 필요합니다: pip install reportlab")
    
    # 언어 설정 확인 및 기본값 설정
    if lang not in ["ko", "en", "ja"]:
        lang = "ko"
    L = LANG.get(lang, LANG["ko"])
    
    if filename is None:
        filename = f"customer_history_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
    filepath = os.path.join(DATA_DIR, filename)
    
    # ⭐ 개선: 한글/일본어 폰트 지원 강화 - 둘 다 등록하여 혼합 사용 가능
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    
    # 한글/일본어 폰트 등록 상태
    korean_font_registered = False
    japanese_font_registered = False
    korean_font_name = 'KoreanFont'
    japanese_font_name = 'JapaneseFont'
    
    def register_font(font_name: str, font_path: str) -> bool:
        """폰트를 등록하는 헬퍼 함수"""
        try:
            if font_path.endswith('.ttf'):
                # TTF 파일 등록
                font = TTFont(font_name, font_path)
                pdfmetrics.registerFont(font)
                if font_name in pdfmetrics.getRegisteredFontNames():
                    return True
            elif font_path.endswith('.ttc'):
                # TTC 파일 처리 (여러 서브폰트 시도)
                for subfont_idx in range(8):  # 서브폰트 인덱스 확대 (0-7)
                    try:
                        font = TTFont(font_name, font_path, subfontIndex=subfont_idx)
                        pdfmetrics.registerFont(font)
                        if font_name in pdfmetrics.getRegisteredFontNames():
                            return True
                    except Exception:
                        continue
            return False
        except Exception as e:
            print(f"⚠️ 폰트 등록 실패 ({font_name}, {font_path}): {e}")
            return False
    
    try:
        # 운영체제별 폰트 경로 설정
        import platform
        system = platform.system()
        
        if system == 'Windows':
            # Windows 기본 한글 폰트 경로 (우선순위 순)
            korean_font_paths = [
                "C:/Windows/Fonts/malgun.ttf",  # 맑은 고딕 (TTF)
                "C:/Windows/Fonts/malgunsl.ttf",  # 맑은 고딕 (TTF, 대체)
                "C:/Windows/Fonts/NanumGothic.ttf",  # 나눔고딕
                "C:/Windows/Fonts/NanumBarunGothic.ttf",  # 나눔바른고딕
                "C:/Windows/Fonts/NanumGothicBold.ttf",  # 나눔고딕 볼드
                "C:/Windows/Fonts/gulim.ttc",  # 굴림 (TTC)
                "C:/Windows/Fonts/batang.ttc",  # 바탕 (TTC)
                "C:/Windows/Fonts/malgun.ttc",  # 맑은 고딕 (TTC)
                "C:/Windows/Fonts/NanumGothic.ttc",  # 나눔고딕 (TTC)
            ]
            
            # Windows 일본어 폰트 경로 (한자 지원 강화)
            japanese_font_paths = [
                "C:/Windows/Fonts/msgothic.ttc",  # MS Gothic (일본어 한자 지원)
                "C:/Windows/Fonts/msmincho.ttc",  # MS Mincho (일본어 한자 지원)
                "C:/Windows/Fonts/meiryo.ttc",  # Meiryo (일본어)
                "C:/Windows/Fonts/yuanti.ttc",  # Microsoft YaHei (중국어/일본어 한자 지원)
                "C:/Windows/Fonts/notosanscjksc-regular.otf",  # Noto Sans CJK (한중일 통합)
            ]
        elif system == 'Darwin':  # macOS
            korean_font_paths = [
                "/System/Library/Fonts/Supplemental/AppleGothic.ttf",
                "/Library/Fonts/AppleGothic.ttf",
                "/System/Library/Fonts/AppleGothic.ttc",
            ]
            japanese_font_paths = [
                "/System/Library/Fonts/Supplemental/AppleGothic.ttf",  # 한중일 통합
                "/System/Library/Fonts/Hiragino Sans GB.ttc",
                "/System/Library/Fonts/STHeiti Light.ttc",
            ]
        else:  # Linux
            korean_font_paths = [
                "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
                "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
                "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
            ]
            japanese_font_paths = [
                "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",  # 한중일 통합
                "/usr/share/fonts/truetype/takao/TakaoGothic.ttf",
            ]
        
        # 한글 폰트 등록 (모든 경로 시도)
        for font_path in korean_font_paths:
            if os.path.exists(font_path):
                if register_font(korean_font_name, font_path):
                    korean_font_registered = True
                    print(f"✅ 한글 폰트 등록 성공: {font_path}")
                    break
        
        # 일본어 폰트 등록 (한글과 독립적으로 등록 - 둘 다 사용 가능)
        for font_path in japanese_font_paths:
            if os.path.exists(font_path):
                if register_font(japanese_font_name, font_path):
                    japanese_font_registered = True
                    print(f"✅ 일본어 폰트 등록 성공: {font_path}")
                    break
        
        # 폰트 등록 실패 시 경고
        if not korean_font_registered and not japanese_font_registered:
            print("⚠️ 경고: 한글/일본어 폰트를 찾을 수 없습니다. PDF에서 한글이 깨질 수 있습니다.")
            print(f"   시스템: {system}")
            print("   등록된 폰트 목록:", pdfmetrics.getRegisteredFontNames())
            if system == 'Windows':
                print("   폰트 경로 확인 필요: C:/Windows/Fonts/")
            elif system == 'Darwin':
                print("   폰트 경로 확인 필요: /System/Library/Fonts/")
            else:
                print("   폰트 경로 확인 필요: /usr/share/fonts/")
            
    except Exception as e:
        error_msg = str(e)
        print(f"⚠️ 폰트 등록 실패: {error_msg}")
        korean_font_registered = False
        japanese_font_registered = False
    
    doc = SimpleDocTemplate(filepath, pagesize=letter)
    story = []
    styles = getSampleStyleSheet()
    
    # ⭐ 개선: 텍스트 내용에 따라 적절한 폰트를 선택하는 스타일 생성 함수
    def get_multilingual_style(base_style_name, default_font=None, **kwargs):
        """다국어 지원 스타일 생성 (한글/일본어/영어)"""
        base_style = styles[base_style_name]
        style_kwargs = {
            'parent': base_style,
            **kwargs
        }
        
        # 기본 폰트 설정 (한글 우선, 없으면 일본어, 없으면 기본)
        registered_fonts = pdfmetrics.getRegisteredFontNames()
        if default_font and default_font in registered_fonts:
            style_kwargs['fontName'] = default_font
        elif korean_font_registered and korean_font_name in registered_fonts:
            style_kwargs['fontName'] = korean_font_name
        elif japanese_font_registered and japanese_font_name in registered_fonts:
            style_kwargs['fontName'] = japanese_font_name
        elif not korean_font_registered and not japanese_font_registered:
            print("⚠️ 경고: 한글/일본어 폰트가 없어 기본 폰트를 사용합니다. 한글이 깨질 수 있습니다.")
        
        return ParagraphStyle(f'Multilingual{base_style_name}', **style_kwargs)
    
    # 제목 스타일 (한글 폰트 우선 사용)
    title_style = get_multilingual_style(
        'Heading1',
        fontSize=24,
        textColor=black,
        spaceAfter=30,
        alignment=1,  # 중앙 정렬
        default_font=korean_font_name if korean_font_registered else japanese_font_name
    )
    
    # 일반 텍스트 스타일
    normal_style = get_multilingual_style('Normal')
    heading1_style = get_multilingual_style('Heading1')
    heading2_style = get_multilingual_style('Heading2')
    
    # ⭐ 개선: 텍스트를 안전하게 처리하고 적절한 폰트를 선택하는 헬퍼 함수
    def safe_text(text, detect_font=True):
        """텍스트를 안전하게 처리하여 PDF에 표시 (한글/일본어/한자 지원 강화)
        
        Args:
            text: 처리할 텍스트
            detect_font: 텍스트 내용에 따라 폰트를 자동 선택할지 여부
        
        Returns:
            (처리된 텍스트, 추천 폰트명) 튜플
        """
        if text is None:
            return ("N/A", None)
        
        # 문자열로 변환 (UTF-8 인코딩 명시적 처리)
        text_str = None
        if isinstance(text, bytes):
            # 바이트 문자열인 경우 UTF-8로 디코딩 시도
            try:
                text_str = text.decode('utf-8', errors='replace')
            except:
                try:
                    # UTF-8 실패 시 다른 인코딩 시도
                    text_str = text.decode('cp949', errors='replace')  # 한국어 Windows 인코딩
                except:
                    try:
                        text_str = text.decode('shift_jis', errors='replace')  # 일본어 인코딩
                    except:
                        try:
                            text_str = text.decode('euc-kr', errors='replace')  # 한국어 EUC-KR
                        except:
                            text_str = text.decode('latin-1', errors='replace')
        else:
            text_str = str(text)
        
        # None 체크
        if text_str is None:
            return ("N/A", None)
        
        # 유니코드 정규화 (NFC 형식으로 통일) - 한글/일본어 문자 정확도 향상
        try:
            import unicodedata
            text_str = unicodedata.normalize('NFC', text_str)
        except:
            pass
        
        # 특수 문자 이스케이프 (HTML 엔티티로 변환) - ReportLab Paragraph는 HTML을 지원
        # 하지만 &는 먼저 처리해야 함
        text_str = text_str.replace('&', '&amp;')
        text_str = text_str.replace('<', '&lt;')
        text_str = text_str.replace('>', '&gt;')
        text_str = text_str.replace('"', '&quot;')
        text_str = text_str.replace("'", '&#39;')
        
        # 폰트 선택 로직 (텍스트 내용 분석)
        recommended_font = None
        if detect_font:
            try:
                # 유니코드 범위 확인
                # 한글: AC00-D7AF (완성형), 1100-11FF (자모)
                # 일본어 히라가나: 3040-309F, 가타카나: 30A0-30FF, 한자: 4E00-9FFF
                has_korean = any('\uAC00' <= char <= '\uD7AF' or '\u1100' <= char <= '\u11FF' for char in text_str)
                has_japanese = any('\u3040' <= char <= '\u309F' or '\u30A0' <= char <= '\u30FF' or '\u4E00' <= char <= '\u9FFF' for char in text_str)
                
                registered_fonts = pdfmetrics.getRegisteredFontNames()
                
                if has_korean and korean_font_registered and korean_font_name in registered_fonts:
                    recommended_font = korean_font_name
                elif has_japanese and japanese_font_registered and japanese_font_name in registered_fonts:
                    recommended_font = japanese_font_name
                elif has_korean or has_japanese:
                    # 한글/일본어 문자가 있지만 적절한 폰트가 없는 경우
                    if korean_font_registered and korean_font_name in registered_fonts:
                        recommended_font = korean_font_name
                    elif japanese_font_registered and japanese_font_name in registered_fonts:
                        recommended_font = japanese_font_name
                    else:
                        print(f"⚠️ 경고: 한글/일본어 문자가 포함되어 있지만 폰트가 등록되지 않았습니다.")
                        print(f"   텍스트 샘플: {text_str[:50]}")
                        print(f"   등록된 폰트: {registered_fonts}")
            except Exception as check_error:
                # 확인 중 오류가 발생해도 계속 진행
                pass
        
        return (text_str, recommended_font)
    
    # Paragraph 생성 헬퍼 함수 (폰트 자동 선택)
    def create_paragraph(text, style, auto_font=True):
        """텍스트와 스타일로 Paragraph 생성 (폰트 자동 선택)"""
        text_str, recommended_font = safe_text(text, detect_font=auto_font)
        
        # 추천 폰트가 있고 스타일에 폰트가 설정되지 않은 경우
        if recommended_font and auto_font:
            # 새로운 스타일 생성 (폰트 포함)
            style_with_font = ParagraphStyle(
                name=f'{style.name}_with_font',
                parent=style,
                fontName=recommended_font
            )
            return Paragraph(text_str, style_with_font)
        
        return Paragraph(text_str, style)
    
    # 제목 추가
    title_text, _ = safe_text(L.get("download_history_title", "고객 응대 이력 요약"))
    story.append(Paragraph(title_text, title_style))
    story.append(Spacer(1, 0.2*inch))
    
    # 각 이력 추가
    for i, hist in enumerate(histories, 1):
        # 이력 제목
        heading_text, _ = safe_text(f'{L.get("download_history_number", "이력 #")}{i}')
        story.append(Paragraph(heading_text, heading1_style))
        story.append(Spacer(1, 0.1*inch))
        
        # 기본 정보 (폰트 자동 선택)
        id_text, _ = safe_text(f'ID: {hist.get("id", "N/A")}')
        story.append(create_paragraph(id_text, normal_style))
        
        timestamp_text, _ = safe_text(f'{L.get("date_label", "날짜")}: {hist.get("timestamp", "N/A")}')
        story.append(create_paragraph(timestamp_text, normal_style))
        
        query_text, _ = safe_text(f'{L.get("download_initial_inquiry", "초기 문의")}: {hist.get("initial_query", "N/A")}')
        story.append(create_paragraph(query_text, normal_style))
        
        customer_type_text, _ = safe_text(f'{L.get("customer_type_label", "고객 유형")}: {hist.get("customer_type", "N/A")}')
        story.append(create_paragraph(customer_type_text, normal_style))
        
        language_text, _ = safe_text(f'{L.get("language_label", "언어")}: {hist.get("language_key", "N/A")}')
        story.append(create_paragraph(language_text, normal_style))
        
        summary = hist.get('summary', {})
        if summary:
            story.append(Spacer(1, 0.1*inch))
            summary_title, _ = safe_text(L.get("download_summary", "요약"))
            story.append(Paragraph(summary_title, heading2_style))
            
            main_inquiry_text, _ = safe_text(f'{L.get("download_main_inquiry", "주요 문의")}: {summary.get("main_inquiry", "N/A")}')
            story.append(create_paragraph(main_inquiry_text, normal_style))
            
            key_responses = summary.get("key_responses", [])
            if isinstance(key_responses, list):
                responses_list = []
                for r in key_responses:
                    r_text, _ = safe_text(r)
                    responses_list.append(r_text)
                responses_text = ", ".join(responses_list)
            else:
                responses_text, _ = safe_text(key_responses)
            responses_para_text, _ = safe_text(f'{L.get("download_key_response", "핵심 응답")}: {responses_text}')
            story.append(create_paragraph(responses_para_text, normal_style))
            
            sentiment_text, _ = safe_text(f'{L.get("sentiment_score_label", "고객 감정 점수")}: {summary.get("customer_sentiment_score", "N/A")}/100')
            story.append(create_paragraph(sentiment_text, normal_style))
            
            satisfaction_text, _ = safe_text(f'{L.get("customer_satisfaction_score_label", "고객 만족도 점수")}: {summary.get("customer_satisfaction_score", "N/A")}/100')
            story.append(create_paragraph(satisfaction_text, normal_style))
            
            characteristics = summary.get('customer_characteristics', {})
            story.append(Spacer(1, 0.1*inch))
            char_title, _ = safe_text(L.get("download_customer_characteristics", "고객 특성"))
            story.append(Paragraph(char_title, heading2_style))
            
            lang_char_text, _ = safe_text(f'{L.get("language_label", "언어")}: {characteristics.get("language", "N/A")}')
            story.append(create_paragraph(lang_char_text, normal_style))
            
            cultural_text, _ = safe_text(f'{L.get("download_cultural_background", "문화적 배경")}: {characteristics.get("cultural_hints", "N/A")}')
            story.append(create_paragraph(cultural_text, normal_style))
            
            region_text, _ = safe_text(f'{L.get("region_label", "지역")}: {characteristics.get("region", "N/A")}')
            story.append(create_paragraph(region_text, normal_style))
            
            comm_style_text, _ = safe_text(f'{L.get("download_communication_style", "소통 스타일")}: {characteristics.get("communication_style", "N/A")}')
            story.append(create_paragraph(comm_style_text, normal_style))
            
            privacy = summary.get('privacy_info', {})
            story.append(Spacer(1, 0.1*inch))
            privacy_title, _ = safe_text(L.get("download_privacy_summary", "개인정보 요약"))
            story.append(Paragraph(privacy_title, heading2_style))
            
            email_text, _ = safe_text(f'{L.get("email_provided_label", "이메일 제공")}: {L.get("download_yes", "예") if privacy.get("has_email") else L.get("download_no", "아니오")}')
            story.append(create_paragraph(email_text, normal_style))
            
            phone_text, _ = safe_text(f'{L.get("phone_provided_label", "전화번호 제공")}: {L.get("download_yes", "예") if privacy.get("has_phone") else L.get("download_no", "아니오")}')
            story.append(create_paragraph(phone_text, normal_style))
            
            address_text, _ = safe_text(f'{L.get("download_address_provided", "주소 제공")}: {L.get("download_yes", "예") if privacy.get("has_address") else L.get("download_no", "아니오")}')
            story.append(create_paragraph(address_text, normal_style))
            
            region_hint_text, _ = safe_text(f'{L.get("download_region_hint", "지역 힌트")}: {privacy.get("region_hint", "N/A")}')
            story.append(create_paragraph(region_hint_text, normal_style))
            
            full_summary_text, _ = safe_text(f'{L.get("download_overall_summary", "전체 요약")}: {summary.get("summary", "N/A")}')
            story.append(create_paragraph(full_summary_text, normal_style))
        
        # 구분선
        if i < len(histories):
            story.append(Spacer(1, 0.2*inch))
            divider_text, _ = safe_text('-' * 80)
            story.append(Paragraph(divider_text, normal_style))
            story.append(Spacer(1, 0.2*inch))
    
    # PDF 빌드 (UTF-8 인코딩 명시, 폰트 서브셋팅 강화)
    try:
        # 폰트 등록 상태 확인 및 경고
        registered_fonts = pdfmetrics.getRegisteredFontNames()
        print(f"📋 등록된 폰트 목록: {registered_fonts}")
        
        if not korean_font_registered and not japanese_font_registered:
            print("⚠️ 경고: 한글/일본어 폰트가 등록되지 않았습니다. PDF에서 한글이 깨질 수 있습니다.")
            print("   가능한 해결 방법:")
            import platform
            system = platform.system()
            if system == 'Windows':
                print("   1. Windows 폰트 폴더(C:/Windows/Fonts/)에 한글 폰트가 설치되어 있는지 확인")
                print("   2. 관리자 권한으로 실행")
                print("   3. 맑은 고딕(malgun.ttf) 또는 나눔고딕(NanumGothic.ttf) 설치 확인")
            elif system == 'Darwin':
                print("   1. macOS 시스템 폰트(/System/Library/Fonts/) 확인")
                print("   2. AppleGothic 폰트 설치 확인")
            else:
                print("   1. Linux 시스템 폰트(/usr/share/fonts/) 확인")
                print("   2. Noto Sans CJK 또는 Nanum 폰트 설치 확인")
        else:
            if korean_font_registered:
                print(f"✅ 한글 폰트 등록 확인: {korean_font_name} in {registered_fonts}")
            if japanese_font_registered:
                print(f"✅ 일본어 폰트 등록 확인: {japanese_font_name} in {registered_fonts}")
            print("✅ 한글/일본어 텍스트가 올바르게 표시됩니다.")
        
        # PDF 빌드 실행 (폰트 서브셋팅 자동 적용)
        doc.build(story)
        print(f"✅ PDF 생성 완료: {filepath}")
        print(f"   파일 크기: {os.path.getsize(filepath) / 1024:.2f} KB")
        
    except Exception as e:
        # 인코딩 오류가 발생하면 에러 메시지와 함께 재시도
        error_msg = str(e)
        print(f"⚠️ PDF 빌드 오류: {error_msg}")
        
        # 폰트 관련 오류인 경우 추가 정보 제공
        if 'font' in error_msg.lower() or 'encoding' in error_msg.lower():
            print("   폰트/인코딩 오류로 보입니다. 폰트 등록 상태를 확인하세요.")
            registered_fonts = pdfmetrics.getRegisteredFontNames()
            print(f"   등록된 폰트: {registered_fonts}")
            if korean_font_registered:
                print(f"   - 한글 폰트: 등록됨 ({korean_font_name})")
            else:
                print(f"   - 한글 폰트: 등록되지 않음")
            if japanese_font_registered:
                print(f"   - 일본어 폰트: 등록됨 ({japanese_font_name})")
            else:
                print(f"   - 일본어 폰트: 등록되지 않음")
        
        # 재시도 (단순 재시도는 위험할 수 있으므로 에러를 다시 발생시킴)
        raise Exception(f"PDF 생성 실패: {error_msg}")
    
    return filepath


# ========================================
# 6. RAG Helper (FAISS)
# ========================================
# RAG 관련 함수는 시뮬레이터와 무관하므로 기존 코드를 유지합니다.

def load_documents(files) -> List[Document]:
    docs: List[Document] = []
    for f in files:
        name = f.name
        lower = name.lower()
        if lower.endswith(".pdf"):
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
            tmp.write(f.read())
            tmp.flush()
            tmp.close()
            loader = PyPDFLoader(tmp.name)
            file_docs = loader.load()
            for d in file_docs:
                d.metadata["source"] = name
            docs.extend(file_docs)
            try:
                os.remove(tmp.name)
            except OSError:
                pass
        elif lower.endswith(".txt"):
            text = f.read().decode("utf-8", errors="ignore")
            docs.append(Document(page_content=text, metadata={"source": name}))
        elif lower.endswith(".html") or lower.endswith(".htm"):
            text = f.read().decode("utf-8", errors="ignore")
            docs.append(Document(page_content=text, metadata={"source": name}))
    return docs


def split_documents(docs: List[Document]) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=150,
        separators=["\n\n", "\n", ".", " ", ""],
    )
    return splitter.split_documents(docs)


def get_embedding_model():
    if get_api_key("openai"):
        try:
            return OpenAIEmbeddings(model="text-embedding-3-small")
        except:
            pass
    if get_api_key("gemini"):
        try:
            return GoogleGenerativeAIEmbeddings(model="models/text-embedding-004")
        except:
            pass
    return None


def get_embedding_function():
    """
    RAG 임베딩에 사용할 임베딩 모델을 결정합니다.
    API 키 유효성 순서: OpenAI (사용자 설정 시) -> Gemini -> NVIDIA -> HuggingFace (fallback)
    API 인증 오류 발생 시 다음 모델로 이동하도록 처리합니다.
    """

    # 1. OpenAI 임베딩 시도 (사용자가 유효한 키를 설정했을 경우)
    openai_key = get_api_key("openai")
    if openai_key:
        try:
            st.info("🔹 RAG: OpenAI Embedding 사용 중")
            return OpenAIEmbeddings(openai_api_key=openai_key)
        except Exception as e:
            st.warning(f"OpenAI 임베딩 실패 → Gemini로 Fallback: {e}")

    # 2. Gemini 임베딩 시도
    gemini_key = get_api_key("gemini")
    if IS_GEMINI_EMBEDDING_AVAILABLE and gemini_key:
        try:
            st.info("🔹 RAG: Gemini Embedding 사용 중")
            # ⭐ 수정: 모델 이름 형식을 'models/model-name'으로 수정
            return GoogleGenerativeAIEmbeddings(google_api_key=gemini_key, model="models/text-embedding-004")
        except Exception as e:
            st.warning(f"Gemini 임베딩 실패 → NVIDIA로 Fallback: {e}")

    # 3. NVIDIA 임베딩 시도
    nvidia_key = get_api_key("nvidia")
    if IS_NVIDIA_EMBEDDING_AVAILABLE and nvidia_key:
        try:
            st.info("🔹 RAG: NVIDIA Embedding 사용 중")
            # NIM 모델 사용 (실제 키가 유효해야 함)
            return NVIDIAEmbeddings(api_key=nvidia_key, model="ai-embed-qa-4")
        except Exception as e:
            st.warning(f"NVIDIA 임베딩 실패 → HuggingFace Fallback: {e}")

    # 4. HuggingFace Embeddings (Local Fallback)
    try:
        st.info("🔹 RAG: Local HuggingFace Embedding 사용 중")
        # 경량 모델 사용
        return HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    except Exception as e:
        st.warning(f"최종 Fallback 임베딩 실패: {e}")

    st.error("❌ RAG 임베딩 실패: 사용 가능한 API Key가 없습니다.")
    return None


def build_rag_index(files):
    # 언어 키 안전하게 가져오기
    current_lang = st.session_state.get("language", "ko")
    if current_lang not in ["ko", "en", "ja"]:
        current_lang = "ko"
    L = LANG.get(current_lang, LANG["ko"])
    if not files: return None, 0

    # 임베딩 함수를 시도하는 과정에서 에러 메시지가 발생할 수 있으므로 try-except로 감쌉니다.
    try:
        embeddings = get_embedding_function()
    except Exception as e:
        st.error(f"RAG 임베딩 함수 초기화 중 치명적인 오류 발생: {e}")
        return None, 0

    if embeddings is None:
        # 어떤 임베딩 모델도 초기화할 수 없음을 알림
        error_msg = L["rag_embed_error_none"]

        # 상세 오류 정보 구성 (실제 사용 가능한 임베딩 모델이 없는 경우)
        if not get_api_key("openai"):
            error_msg += f"\n- {L['rag_embed_error_openai']}"
        if not get_api_key("gemini"):
            error_msg += f"\n- {L['rag_embed_error_gemini']}"
        if not get_api_key("nvidia"):
            error_msg += f"\n- {L['rag_embed_error_nvidia']}"

        st.error(error_msg)
        return None, 0

    # 임베딩 객체 초기화 성공 후, 데이터 로드 및 분할
    docs = load_documents(files)
    if not docs: return None, 0

    chunks = split_documents(docs)
    if not chunks: return None, 0

    try:
        vectorstore = FAISS.from_documents(chunks, embeddings)
        # 저장
        vectorstore.save_local(RAG_INDEX_DIR)
    except Exception as e:
        # API 인증 실패 등 실제 API 호출 오류 처리
        st.error(f"RAG 인덱스 생성 중 오류: {e}")
        return None, 0

    return vectorstore, len(chunks)


def load_rag_index():
    # RAG 인덱스 로드 시에도 유효한 임베딩 함수가 필요합니다.
    try:
        embeddings = get_embedding_function()
    except Exception:
        # get_embedding_function 내에서 에러 메시지를 처리하거나 스킵하므로 여기서는 조용히 처리
        return None

    if embeddings is None:
        return None

    try:
        # allow_dangerous_deserialization=True는 필수
        vs = FAISS.load_local(RAG_INDEX_DIR, embeddings, allow_dangerous_deserialization=True)
        return vs
    except Exception:
        return None


def rag_answer(question: str, vectorstore: FAISS, lang_key: str) -> str:
    # RAG Answer는 LLM 클라이언트 라우팅을 사용하도록 수정
    llm_client, info = get_llm_client()
    if llm_client is None:
        # 언어 키 검증
        if lang_key not in ["ko", "en", "ja"]:
            lang_key = "ko"
        return LANG.get(lang_key, LANG["ko"]).get("simulation_no_key_warning", "API Key가 필요합니다.")

    # Langchain ChatOpenAI 대신 run_llm을 사용하기 위해 prompt를 직접 구성
    retriever = vectorstore.as_retriever(search_kwargs={"k": 4})
    docs = retriever.get_relevant_documents(question)
    context = "\n\n".join(d.page_content[:1500] for d in docs)

    # ⭐ RAG 다국어 인식 오류 해결: 답변 생성 모델에게 질문 언어로 일관되게 답하도록 강력히 지시
    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}.get(lang_key, "English")

    prompt = (
            f"You are a helpful AI tutor. Answer the question using ONLY the provided context.\n"
            f"The answer MUST be STRICTLY in {lang_name}, which is the language of the question.\n"
            f"If you cannot find the answer in the context, say you don't know in {lang_name}.\n"
            f"Note: The context may be in a different language, but you must still answer in {lang_name}.\n\n"
            "Question:\n" + question + "\n\n"
                                       "Context:\n" + context + "\n\n"
                                                                f"Answer (in {lang_name}):"
    )
    return run_llm(prompt)


# ========================================
# 7. LSTM Helper (간단 Mock + 시각화)
# ========================================

def load_or_train_lstm():
    # 실제 LSTM 대신 랜덤 + sin 파형 기반 Mock
    np.random.seed(42)
    n_points = 50
    ts = 60 + 20 * np.sin(np.linspace(0, 4 * np.pi, n_points)) + np.random.normal(0, 5, n_points)
    ts = np.clip(ts, 50, 100).astype(np.float32)
    return ts





# ========================================
# 8. LLM (ChatOpenAI) for Simulator / Content
# (RAG와 동일하게 run_llm으로 통합)
# ========================================

# ConversationChain 대신 run_llm을 사용하여 메모리 기능을 수동으로 구현
# st.session_state.simulator_memory는 유지하여 대화 기록을 관리합니다.


def get_chat_history_for_prompt(include_attachment=False):
    """메모리에서 대화 기록을 추출하여 프롬프트에 사용할 문자열 형태로 반환 (채팅용)"""
    history_str = ""
    for msg in st.session_state.simulator_messages:
        role = msg["role"]
        content = msg["content"]
        if role == "customer" or role == "customer_rebuttal":
            history_str += f"Customer: {content}\n"
        elif role == "agent_response":
            history_str += f"Agent: {content}\n"
        # supervisor 메시지는 LLM에 전달하지 않아 역할 혼동 방지
    return history_str



def generate_customer_reaction(current_lang_key: str, is_call: bool = False) -> str:
    """
    고객의 다음 반응을 생성하는 LLM 호출 (채팅 전용)
    **수정 사항:** 에이전트 정보 요청 시 필수 정보 (주문번호, eSIM, 자녀 만 나이, 취소 사유) 제공 의무를 강화함.
    """
    history_text = get_chat_history_for_prompt()
    # 언어 키 검증
    if current_lang_key not in ["ko", "en", "ja"]:
        current_lang_key = st.session_state.get("language", "ko")
        if current_lang_key not in ["ko", "en", "ja"]:
            current_lang_key = "ko"
    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}[current_lang_key]
    L_local = LANG.get(current_lang_key, LANG["ko"])

    # 첨부 파일 컨텍스트 추가
    attachment_context = st.session_state.sim_attachment_context_for_llm
    if attachment_context:
        # LLM에게 첨부 파일 컨텍스트를 제공하되, 에이전트에게 반복하지 않도록 주의
        attachment_context = f"[INITIAL ATTACHMENT CONTEXT (for customer reference only, do not repeat to agent)]\n{attachment_context}\n\n"
    else:
        attachment_context = ""

    next_prompt = f"""
{attachment_context}
You are now ROLEPLAYING as the CUSTOMER.

Read the following conversation and respond naturally in {lang_name}.

Conversation so far:
{history_text}

RULES:
1. You are only the customer. Do not write as the agent.
2. **[CRITICAL: Mandatory Information Submission for Problem Resolution]** If the agent requested any of the following critical information, you MUST provide it:
    - Order/Booking Number (e.g., ABC123, 123456)
    - eSIM related details (e.g., Host device compatibility, local status/location, time of activation)
    - Child-related product details (e.g., Child's Date of Birth or Current Age)
    - Exception/Refund Reason (e.g., flight cancellation/delay, illness, local natural disaster)
    - **If you are a difficult customer and the agent requests this information, you MUST still provide it, but you may express frustration or impatience while doing so.**
3. **[Crucial Rule for Repetition/New Inquiry]** After the agent has provided an attempt at a solution or answer:
    - If you are still confused or the problem is not fully solved, you MUST state the remaining confusion/problem clearly and briefly. DO NOT REPEAT THE INITIAL QUERY. Focus only on the unresolved aspect or the new inquiry.
4. **[CRITICAL: Solution Acknowledgment]** If the agent provided a clear and accurate solution/confirmation:
    - You MUST respond with appreciation and satisfaction, like "{L_local['customer_positive_response']}" or similar positive acknowledgment. This applies even if you are a difficult customer.
5. If the agent's LAST message was the closing confirmation: "{L_local['customer_closing_confirm']}"
    - If you have NO additional questions: You MUST reply with "{L_local['customer_no_more_inquiries']}".
    - If you DO have additional questions: You MUST reply with "{L_local['customer_has_additional_inquiries']}" AND MUST FOLLOW UP WITH THE NEW INQUIRY DETAILS IMMEDIATELY. DO NOT just repeat that you have an additional question.
6. Do NOT repeat your initial message or previous responses unless necessary.
7. Output ONLY the customer's next message.
"""
    try:
        reaction = run_llm(next_prompt)

        # ⭐ LLM이 응답했지만 내용이 너무 짧거나 비어있을 경우, 긍정 종료 문구를 반환
        if not reaction or len(reaction.strip()) < 5:
            print("LLM returned insufficient response. Using positive closing fallback.")
            return L_local['customer_positive_response']

        return reaction.strip()
    except Exception as e:
        # ⭐ LLM 호출 자체에서 오류 발생 시 (API 키, 할당량) 긍정 종료 문구를 강제 반환
        print(f"LLM Customer Reaction generation failed: {e}. Falling back to positive closing.")
        return L_local['customer_positive_response']  # 강제 안전장치



def summarize_history_with_ai(current_lang_key: str) -> str:
    """전화 통화 로그를 정리하여 LLM에 전달하고 요약 텍스트를 받는 함수."""
    # 전화 로그는 'phone_exchange' 역할을 가지거나, 'initial_query'에 포함되어 있음

    # 1. 로그 추출
    conversation_text = ""
    initial_query = st.session_state.get("call_initial_query", "N/A")
    website_url = st.session_state.get("call_website_url", "").strip()
    if initial_query and initial_query != "N/A":
        conversation_text += f"Initial Query: {initial_query}\n"
    if website_url:
        conversation_text += f"Website URL: {website_url}\n"

    for msg in st.session_state.simulator_messages:
        role = msg.get("role", "")
        content = msg.get("content", "")
        if role == "phone_exchange":
            # phone_exchange는 "Agent: ... | Customer: ..." 형태로 이미 정리되어 있음
            conversation_text += f"{content}\n"

    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}[current_lang_key]

    summary_prompt = f"""
You are an AI Analyst specialized in summarizing customer phone calls. 
Analyze the full conversation log below, identify the main issue, the steps taken by the agent, and the customer's sentiment.

Provide a concise, easy-to-read summary of the key exchange STRICTLY in {lang_name}.

--- Conversation Log ---
{conversation_text}
---

Summary:
"""
    if not st.session_state.is_llm_ready:
        return "LLM Key가 없어 요약 생성이 불가합니다."

    try:
        summary = run_llm(summary_prompt)
        return summary.strip()
    except Exception as e:
        return f"❌ AI 요약 생성 오류: {e}"



def generate_customer_reaction_for_call(current_lang_key: str, last_agent_response: str) -> str:
    """전화 시뮬레이터 전용 고객 반응 생성 (마지막 에이전트 응답 중심)"""
    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}[current_lang_key]
    L_local = LANG[current_lang_key]
    
    # ⭐ 추가: 고객 성별 및 감정 상태 가져오기
    customer_gender = st.session_state.customer_avatar.get("gender", "male")
    customer_emotion = st.session_state.customer_avatar.get("state", "NEUTRAL")
    
    # 감정 상태에 따른 톤 설정
    emotion_tone_map = {
        "HAPPY": "friendly, positive, and satisfied",
        "ASKING": "slightly frustrated, questioning, and seeking clarification",
        "ANGRY": "angry, frustrated, and demanding",
        "SAD": "sad, depressed, and disappointed",
        "NEUTRAL": "neutral, calm, and polite"
    }
    emotion_tone = emotion_tone_map.get(customer_emotion, "neutral, calm, and polite")
    
    gender_pronoun = "she" if customer_gender == "female" else "he"
    
    # ⭐ 추가: 에이전트가 종료 확인 질문을 했는지 확인
    closing_msg = L_local['customer_closing_confirm']
    is_closing_question = closing_msg in last_agent_response or any(
        phrase in last_agent_response.lower() 
        for phrase in ["다른 문의", "추가 문의", "다른 도움", "anything else", "other questions"]
    )
    
    # ⭐ 수정: 초기 문의를 완전히 제거하고 마지막 에이전트 응답에만 집중
    # 최근 대화 이력만 추출 (최대 3-4개 교환만)
    recent_exchanges = []
    for msg in reversed(st.session_state.simulator_messages):  # 역순으로 최근 것부터
        role = msg.get("role", "")
        content = msg.get("content", "")
        
        if role == "phone_exchange":
            recent_exchanges.insert(0, content)  # 앞에 삽입하여 순서 유지
            if len(recent_exchanges) >= 3:  # 최근 3개만
                break
        elif role == "agent" or role == "agent_response":
            # agent와 agent_response 역할 모두 처리
            recent_exchanges.insert(0, f"Agent: {content}")
            if len(recent_exchanges) >= 3:
                break
    
    # 최근 대화 이력 (있는 경우만)
    recent_history = "\n".join(recent_exchanges) if recent_exchanges else "(No previous exchanges)"
    
    website_url = st.session_state.get("call_website_url", "").strip()
    website_context = f"\nWebsite URL: {website_url}" if website_url else ""
    
    # ⭐ 수정: 마지막 에이전트 응답만 강조 (초기 문의 완전 제거)
    last_agent_text = last_agent_response.strip() if last_agent_response else "None"
    
    history_text = f"""[Recent Conversation Context - For Reference Only]
{recent_history}{website_context}

═══════════════════════════════════════════════════════════════════
🎯 YOUR TASK: Respond ONLY to the Agent's message below
═══════════════════════════════════════════════════════════════════

Agent just said: "{last_agent_text}"

═══════════════════════════════════════════════════════════════════
IMPORTANT: 
- Respond DIRECTLY to what the agent JUST SAID above
- DO NOT repeat your initial query
- DO NOT refer to old conversation unless agent asks
- Keep your response short and conversational
- Your emotional state: {customer_emotion} - respond with {emotion_tone} tone
═══════════════════════════════════════════════════════════════════"""

    # ⭐ 추가: 종료 확인 질문에 대한 특별 처리
    if is_closing_question:
        call_prompt = f"""
You are a CUSTOMER in a phone call. You are a {customer_gender} customer. Respond naturally in {lang_name}.

Your current emotional state: {customer_emotion}
Your response tone should be: {emotion_tone}

{history_text}

The agent just asked: "{last_agent_text}"

CRITICAL RULES FOR CLOSING CONFIRMATION:
1. If you have NO additional questions and the conversation is resolved:
   - You MUST reply with: "{L_local['customer_no_more_inquiries']}"
2. If you DO have additional questions or the issue is NOT fully resolved:
   - You MUST reply with: "{L_local['customer_has_additional_inquiries']}" AND immediately state your additional question
3. Your response MUST be ONLY one of the two options above, in {lang_name}.
4. Output ONLY the customer's response (must be one of the two rule options).

Your response (respond to the closing confirmation question):
"""
    else:
        call_prompt = f"""
You are a CUSTOMER in a phone call. You are a {customer_gender} customer. Respond naturally in {lang_name}.

Your current emotional state: {customer_emotion}
Your response tone should be: {emotion_tone}

{history_text}

RULES:
1. **CRITICAL**: Respond DIRECTLY and ACCURATELY to what the agent JUST SAID: "{last_agent_text}"
2. **If agent asked a question** → Answer it SPECIFICALLY and DIRECTLY. Do not give vague or unrelated answers.
3. **If agent requested information** → Provide the EXACT information requested. Be precise and relevant.
4. **If agent gave a solution or instruction** → Acknowledge it clearly and indicate if you understand or need clarification, based on your emotional state ({customer_emotion})
5. **If agent asked for confirmation** → Confirm or clarify based on what was asked
6. Keep your response short (1-2 sentences max) and focused ONLY on what the agent just said
7. DO NOT repeat your initial query unless the agent specifically asks about it
8. DO NOT mention old conversation unless the agent refers to it
9. IMPORTANT: Match your tone to your emotional state ({customer_emotion}) - be {emotion_tone}
10. **Your response must be a direct answer to the agent's last message above. Read it carefully and respond accordingly.**

Your response (respond ONLY to the agent's message above, with {emotion_tone} tone):
"""
    try:
        reaction = run_llm(call_prompt)
        reaction_text = reaction.strip()
        
        # ⭐ 추가: 종료 확인 질문에 대한 응답 검증 및 강제 적용
        if is_closing_question:
            if L_local['customer_no_more_inquiries'] in reaction_text:
                return L_local['customer_no_more_inquiries']
            elif L_local['customer_has_additional_inquiries'] in reaction_text:
                return reaction_text  # 추가 문의 내용 포함 가능
            else:
                # LLM이 규칙을 따르지 않으면, 대화가 해결된 것으로 가정하고 종료 응답 반환
                return L_local['customer_no_more_inquiries']
        
        return reaction_text
    except Exception as e:
        return f"❌ 고객 반응 생성 오류: {e}"



def generate_customer_reaction_for_first_greeting(current_lang_key: str, agent_greeting: str, initial_query: str) -> str:
    """전화 시뮬레이터 전용: 첫 인사말에 대한 고객의 맞춤형 반응 생성 (초기 문의 고려)"""
    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}[current_lang_key]
    L_local = LANG[current_lang_key]
    
    # ⭐ 추가: 고객 성별 및 감정 상태 가져오기
    customer_gender = st.session_state.customer_avatar.get("gender", "male")
    customer_emotion = st.session_state.customer_avatar.get("state", "NEUTRAL")
    
    # 감정 상태에 따른 톤 설정
    emotion_tone_map = {
        "HAPPY": "friendly, positive, and satisfied",
        "ASKING": "slightly frustrated, questioning, and seeking clarification",
        "ANGRY": "angry, frustrated, and demanding",
        "SAD": "sad, depressed, and disappointed",
        "NEUTRAL": "neutral, calm, and polite"
    }
    emotion_tone = emotion_tone_map.get(customer_emotion, "neutral, calm, and polite")
    
    website_url = st.session_state.get("call_website_url", "").strip()
    website_context = f"\nWebsite URL: {website_url}" if website_url else ""
    
    agent_greeting_text = agent_greeting.strip() if agent_greeting else "None"
    initial_query_text = initial_query.strip() if initial_query else "None"
    
    call_prompt = f"""
You are a CUSTOMER in a phone call. You are a {customer_gender} customer. Respond naturally in {lang_name}.

Your current emotional state: {customer_emotion}
Your response tone should be: {emotion_tone}

═══════════════════════════════════════════════════════════════════
🎯 YOUR SITUATION:
═══════════════════════════════════════════════════════════════════

You called because: "{initial_query_text}"

The agent just greeted you and said: "{agent_greeting_text}"
{website_context}

═══════════════════════════════════════════════════════════════════
YOUR TASK: Respond to the agent's greeting in a way that:
═══════════════════════════════════════════════════════════════════

1. Acknowledge the agent's greeting naturally
2. Briefly mention your inquiry/concern: "{initial_query_text}"
3. Show that you're ready to discuss your issue
4. Keep it conversational and natural (1-2 sentences max)
5. DO NOT be overly formal - this is a phone call, be natural
6. IMPORTANT: Match your tone to your emotional state ({customer_emotion}) - be {emotion_tone}

Example good responses (adjust tone based on your emotional state):
- If {customer_emotion}: [Respond with {emotion_tone} tone]
- "Hello, thank you. I'm calling because [brief mention of issue]..."
- "Hi, yes. I need help with [your issue]..."
- "Thank you. I have a question about [your issue]..."

Your response (respond naturally to the greeting and briefly mention your inquiry, with {emotion_tone} tone):
"""
    try:
        reaction = run_llm(call_prompt)
        return reaction.strip()
    except Exception as e:
        return f"❌ 고객 반응 생성 오류: {e}"



def summarize_history_for_call(call_logs: List[Dict[str, str]], initial_query: str, current_lang_key: str) -> str:
    """전화 통화 로그와 초기 문의를 바탕으로 요약본을 생성"""
    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}[current_lang_key]

    # 로그 재구성 (phone_exchange 역할만 사용)
    full_log_text = f"--- Initial Customer Query ---\nCustomer: {initial_query}\n"
    for log in call_logs:
        if log["role"] == "phone_exchange":
            full_log_text += f"{log['content']}\n"
        elif log["role"] == "agent" and "content" in log:
            # 최초 에이전트 인사말은 여기에 포함
            full_log_text += f"Agent (Greeting): {log['content']}\n"

    summary_prompt = f"""
You are an AI Supervisor. Analyze the following telephone support conversation log.
Provide a concise, neutral summary of the key issue, the steps taken by the agent, and the final outcome.
The summary MUST be STRICTLY in {lang_name}.

--- Conversation Log ---
{full_log_text}
---

Summary:
"""
    if not st.session_state.is_llm_ready:
        return f"❌ LLM Key is missing. Cannot generate summary. Log length: {len(full_log_text.splitlines())}"

    try:
        summary = run_llm(summary_prompt)
        return summary.strip()
    except Exception as e:
        return f"❌ Summary Generation Error: {e}"



def generate_customer_closing_response(current_lang_key: str) -> str:
    """에이전트의 마지막 확인 질문에 대한 고객의 최종 답변 생성 (채팅용)"""
    history_text = get_chat_history_for_prompt()
    # 언어 키 검증
    if current_lang_key not in ["ko", "en", "ja"]:
        current_lang_key = st.session_state.get("language", "ko")
        if current_lang_key not in ["ko", "en", "ja"]:
            current_lang_key = "ko"
    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}[current_lang_key]
    L_local = LANG.get(current_lang_key, LANG["ko"])  # ⭐ 수정: 함수 내에서 사용할 언어 팩

    # 마지막 메시지가 에이전트의 종료 확인 메시지인지 확인 (프롬프트에 포함)
    closing_msg = L_local['customer_closing_confirm']

    # 첨부 파일 컨텍스트 추가
    attachment_context = st.session_state.sim_attachment_context_for_llm
    if attachment_context:
        attachment_context = f"[INITIAL ATTACHMENT CONTEXT (for customer reference only, do not repeat to agent)]\n{attachment_context}\n\n"
    else:
        attachment_context = ""

    final_prompt = f"""
{attachment_context}
You are now ROLEPLAYING as the CUSTOMER.

The agent's final message was the closing confirmation: "{closing_msg}".
You MUST respond to this confirmation based on the overall conversation.

Conversation history:
{history_text}

RULES:
1. If the conversation seems resolved and you have NO additional questions:
    - You MUST reply with "{L_local['customer_no_more_inquiries']}".
2. If the conversation is NOT fully resolved and you DO have additional questions (or the agent provided a cancellation denial that you want to appeal):
    - You MUST reply with "{L_local['customer_has_additional_inquiries']}" AND MUST FOLLOW UP WITH THE NEW INQUIRY DETAILS. DO NOT just repeat that you have an additional question.
3. Your reply MUST be ONLY one of the two options above, in {lang_name}.
4. Output ONLY the customer's next message (must be one of the two rule options).
"""
    try:
        reaction = run_llm(final_prompt)
        # LLM의 출력이 규칙을 따르지 않을 경우를 대비하여 강제 적용
        reaction_text = reaction.strip()
        # "추가 문의 사항도 있습니다"가 포함되어 있으면 그대로 반환 (상세 내용 포함 가정)
        if L_local['customer_no_more_inquiries'] in reaction_text:
            return L_local['customer_no_more_inquiries']
        elif L_local['customer_has_additional_inquiries'] in reaction_text:
            return reaction_text
        else:
            # LLM이 규칙을 어겼을 경우, "추가 문의 사항이 있다"고 가정하고 에이전트 턴으로 넘김
            return L_local['customer_has_additional_inquiries']
    except Exception as e:
        st.error(f"고객 최종 반응 생성 오류: {e}")
        return L_local['customer_has_additional_inquiries']  # 오류 시 에이전트 턴으로 유도


# ----------------------------------------
# Initial Advice/Draft Generation (이관 후 재사용) (요청 4 반영)
# ----------------------------------------

def generate_agent_first_greeting(lang_key: str, initial_query: str) -> str:
    """전화 통화 시작 시 에이전트의 첫 인사말을 생성 (임시 함수)"""
    # 언어 키 검증
    if lang_key not in ["ko", "en", "ja"]:
        lang_key = st.session_state.get("language", "ko")
        if lang_key not in ["ko", "en", "ja"]:
            lang_key = "ko"
    L_local = LANG.get(lang_key, LANG["ko"])
    # 문의 내용의 첫 10자만 사용 (too long)
    topic = initial_query.strip()[:15].replace('\n', ' ')
    if len(initial_query.strip()) > 15:
        topic += "..."

    if lang_key == 'ko':
        return f"안녕하세요, {topic} 관련 문의 주셨죠? 상담원 000입니다. 무엇을 도와드릴까요?"
    elif lang_key == 'en':
        return f"Hello, thank you for calling. I see you're calling about {topic}. My name is 000. How may I help you today?"
    elif lang_key == 'ja':
        return f"お電話ありがとうございます。{topic}の件ですね。担当の000と申します。どのようなご用件でしょうか?"
    return "Hello, how may I help you?"



def detect_text_language(text: str) -> str:
    """
    텍스트의 언어를 자동 감지합니다.
    Returns: "ko", "en", "ja" 중 하나 (기본값: "ko")
    """
    if not text or not text.strip():
        return "ko"  # 기본값
    
    try:
        # 간단한 휴리스틱: 일본어 문자(히라가나, 가타카나, 한자)가 많이 포함되어 있으면 일본어
        japanese_chars = sum(1 for char in text if '\u3040' <= char <= '\u309F' or '\u30A0' <= char <= '\u30FF' or '\u4E00' <= char <= '\u9FAF')
        if japanese_chars > len(text) * 0.1:  # 10% 이상 일본어 문자
            return "ja"
        
        # 영어 문자 비율이 높으면 영어
        english_chars = sum(1 for char in text if char.isascii() and char.isalpha())
        if english_chars > len(text) * 0.7:  # 70% 이상 영어 문자
            return "en"
        
        # LLM을 사용한 정확한 언어 감지 시도 (오류 발생 시 무시하고 휴리스틱 결과 사용)
        if st.session_state.is_llm_ready:
            try:
                detection_prompt = f"""Detect the language of the following text. Respond with ONLY one word: "ko" (Korean), "en" (English), or "ja" (Japanese).

Text: {text[:200]}

Language:"""
                detected = run_llm(detection_prompt).strip().lower()
                # 오류 메시지가 아닌 경우에만 사용
                if detected and detected not in ["❌", "error", "failed"] and detected in ["ko", "en", "ja"]:
                    return detected
            except Exception as e:
                # LLM 호출 실패 시 휴리스틱 결과 사용
                print(f"Language detection LLM call failed: {e}")
                pass
    except Exception as e:
        # 전체 함수에서 예외 발생 시 기본값 반환
        print(f"Language detection error: {e}")
        return "ko"
    
    # 기본값: 한국어
    return "ko"



def analyze_customer_profile(customer_query: str, current_lang_key: str = None) -> Dict[str, Any]:
    """신규 고객의 문의사항과 말투를 분석하여 고객성향 점수를 실시간으로 계산 (요청 4)"""
    # 입력 텍스트의 언어를 자동 감지 (오류 발생 시 안전하게 처리)
    try:
        detected_lang = detect_text_language(customer_query)
    except Exception as e:
        print(f"Language detection failed in analyze_customer_profile: {e}")
        detected_lang = "ko"  # 기본값 사용
    
    # current_lang_key가 제공되지 않으면 감지된 언어 사용
    lang_key_to_use = current_lang_key if current_lang_key else detected_lang
    # lang_key_to_use가 유효한지 확인
    if lang_key_to_use not in ["ko", "en", "ja"]:
        lang_key_to_use = "ko"  # 기본값으로 폴백
    
    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}[lang_key_to_use]

    analysis_prompt = f"""
You are an AI analyst analyzing a customer's inquiry to determine their profile and sentiment.

Analyze the following customer inquiry and provide a structured analysis in JSON format (ONLY JSON, no markdown).

Analyze:
1. Customer gender (male/female/unknown - analyze based on name, language patterns, or cultural hints)
2. Customer sentiment score (0-100, where 0=very negative/angry, 50=neutral, 100=very positive/happy)
3. Communication style (formal/casual, brief/detailed, polite/direct)
4. Urgency level (low/medium/high)
5. Customer type prediction (normal/difficult/very_dissatisfied)
6. Language and cultural hints (if any)
7. Key concerns or pain points

Output format (JSON only):
{{
  "gender": "male",
  "sentiment_score": 45,
  "communication_style": "brief, direct, slightly frustrated",
  "urgency_level": "high",
  "predicted_customer_type": "difficult",
  "cultural_hints": "unknown",
  "key_concerns": ["issue 1", "issue 2"],
  "tone_analysis": "brief description of tone"
}}

Customer Inquiry:
{customer_query}

JSON Output:
"""

    if not st.session_state.is_llm_ready:
        return {
            "gender": "unknown",
            "sentiment_score": 50,
            "communication_style": "unknown",
            "urgency_level": "medium",
            "predicted_customer_type": "normal",
            "cultural_hints": "unknown",
            "key_concerns": [],
            "tone_analysis": "Unable to analyze"
        }

    try:
        analysis_text = run_llm(analysis_prompt).strip()
        # JSON 추출
        if "```json" in analysis_text:
            analysis_text = analysis_text.split("```json")[1].split("```")[0].strip()
        elif "```" in analysis_text:
            analysis_text = analysis_text.split("```")[1].split("```")[0].strip()

        import json
        analysis_data = json.loads(analysis_text)
        return analysis_data
    except Exception as e:
        return {
            "gender": "unknown",
            "sentiment_score": 50,
            "communication_style": "unknown",
            "urgency_level": "medium",
            "predicted_customer_type": "normal",
            "cultural_hints": "unknown",
            "key_concerns": [],
            "tone_analysis": f"Analysis error: {str(e)}"
        }



def find_similar_cases(customer_query: str, customer_profile: Dict[str, Any], current_lang_key: str,
                       limit: int = 5) -> List[Dict[str, Any]]:
    """저장된 요약 데이터에서 유사한 케이스를 찾아 반환 (요청 4)"""
    histories = load_simulation_histories_local(current_lang_key)

    if not histories:
        return []

    # 요약 데이터가 있는 케이스만 필터링
    cases_with_summary = [
        h for h in histories
        if h.get("summary") and isinstance(h.get("summary"), dict) and h.get("is_chat_ended", False)
           and not h.get("is_call", False)  # 전화 이력 제외
    ]

    if not cases_with_summary:
        return []

    # 유사도 계산 (간단한 키워드 매칭 + 점수 유사도)
    similar_cases = []
    query_lower = customer_query.lower()
    customer_sentiment = customer_profile.get("sentiment_score", 50)
    customer_style = customer_profile.get("communication_style", "")

    for case in cases_with_summary:
        summary = case.get("summary", {})
        main_inquiry = summary.get("main_inquiry", "").lower()
        case_sentiment = summary.get("customer_sentiment_score", 50)
        case_satisfaction = summary.get("customer_satisfaction_score", 50)

        # 유사도 점수 계산
        similarity_score = 0

        # 1. 문의 내용 유사도 (키워드 매칭)
        query_words = set(query_lower.split())
        inquiry_words = set(main_inquiry.split())
        if query_words and inquiry_words:
            word_overlap = len(query_words & inquiry_words) / len(query_words | inquiry_words)
            similarity_score += word_overlap * 40

        # 2. 감정 점수 유사도
        sentiment_diff = abs(customer_sentiment - case_sentiment)
        sentiment_similarity = max(0, 1 - (sentiment_diff / 100)) * 30
        similarity_score += sentiment_similarity

        # 3. 만족도 점수 (높을수록 좋은 케이스)
        satisfaction_bonus = (case_satisfaction / 100) * 30
        similarity_score += satisfaction_bonus

        if similarity_score > 30:  # 최소 유사도 임계값
            similar_cases.append({
                "case": case,
                "similarity_score": similarity_score,
                "summary": summary
            })

    # 유사도 순으로 정렬
    similar_cases.sort(key=lambda x: x["similarity_score"], reverse=True)
    return similar_cases[:limit]


def visualize_customer_profile_scores(customer_profile: Dict[str, Any], current_lang_key: str):
    """고객 프로필 점수를 시각화 (감정 점수, 긴급도)"""
    if not IS_PLOTLY_AVAILABLE:
        return None

    # 언어 키 검증
    if current_lang_key not in ["ko", "en", "ja"]:
        current_lang_key = st.session_state.get("language", "ko")
        if current_lang_key not in ["ko", "en", "ja"]:
            current_lang_key = "ko"
    L = LANG.get(current_lang_key, LANG["ko"])

    sentiment_score = customer_profile.get("sentiment_score", 50)
    urgency_map = {"low": 25, "medium": 50, "high": 75}
    urgency_level = customer_profile.get("urgency_level", "medium")
    urgency_score = urgency_map.get(urgency_level.lower(), 50)

    # 게이지 차트 생성
    fig = make_subplots(
        rows=1, cols=2,
        specs=[[{"type": "indicator"}, {"type": "indicator"}]],
        subplot_titles=(
            L.get("sentiment_score_label", "고객 감정 점수"),
            L.get("urgency_score_label", "긴급도 점수")
        )
    )

    # 감정 점수 게이지
    fig.add_trace(
        go.Indicator(
            mode="gauge+number+delta",
            value=sentiment_score,
            domain={'x': [0, 1], 'y': [0, 1]},
            title={'text': L.get("sentiment_score_label", "감정 점수")},
            delta={'reference': 50},
            gauge={
                'axis': {'range': [None, 100]},
                'bar': {'color': "darkblue"},
                'steps': [
                    {'range': [0, 33], 'color': "lightgray"},
                    {'range': [33, 66], 'color': "gray"},
                    {'range': [66, 100], 'color': "lightgreen"}
                ],
                'threshold': {
                    'line': {'color': "red", 'width': 4},
                    'thickness': 0.75,
                    'value': 90
                }
            }
        ),
        row=1, col=1
    )

    # 긴급도 점수 게이지
    fig.add_trace(
        go.Indicator(
            mode="gauge+number",
            value=urgency_score,
            domain={'x': [0, 1], 'y': [0, 1]},
            title={'text': L.get("urgency_score_label", "긴급도")},
            gauge={
                'axis': {'range': [None, 100]},
                'bar': {'color': "darkred"},
                'steps': [
                    {'range': [0, 50], 'color': "lightgreen"},
                    {'range': [50, 75], 'color': "yellow"},
                    {'range': [75, 100], 'color': "lightcoral"}
                ],
            }
        ),
        row=1, col=2
    )

    fig.update_layout(height=300, margin=dict(l=20, r=20, t=50, b=20))
    return fig


def visualize_similarity_cases(similar_cases: List[Dict[str, Any]], current_lang_key: str):
    """유사 케이스 추천을 시각화"""
    if not IS_PLOTLY_AVAILABLE or not similar_cases:
        return None

    # 언어 키 검증
    if current_lang_key not in ["ko", "en", "ja"]:
        current_lang_key = st.session_state.get("language", "ko")
        if current_lang_key not in ["ko", "en", "ja"]:
            current_lang_key = "ko"
    L = LANG.get(current_lang_key, LANG["ko"])

    case_labels = []
    similarity_scores = []
    sentiment_scores = []
    satisfaction_scores = []

    for idx, similar_case in enumerate(similar_cases, 1):
        summary = similar_case["summary"]
        similarity = similar_case["similarity_score"]
        case_labels.append(f"Case {idx}")
        similarity_scores.append(similarity)
        sentiment_scores.append(summary.get("customer_sentiment_score", 50))
        satisfaction_scores.append(summary.get("customer_satisfaction_score", 50))

    # 유사도 차트
    fig = make_subplots(
        rows=2, cols=1,
        subplot_titles=(
            L.get("similarity_chart_title", "유사 케이스 유사도"),
            L.get("scores_comparison_title",
                  "감정 및 만족도 점수 비교")
        ),
        vertical_spacing=0.15
    )

    # 유사도 바 차트
    fig.add_trace(
        go.Bar(
            x=case_labels,
            y=similarity_scores,
            name=L.get("similarity_score_label", "유사도"),
            marker_color='lightblue',
            text=[f"{s:.1f}%" for s in similarity_scores],
            textposition='outside'
        ),
        row=1, col=1
    )

    # 감정 및 만족도 점수 비교
    fig.add_trace(
        go.Bar(
            x=case_labels,
            y=sentiment_scores,
            name=L.get("sentiment_score_label", "감정 점수"),
            marker_color='lightcoral'
        ),
        row=2, col=1
    )

    fig.add_trace(
        go.Bar(
            x=case_labels,
            y=satisfaction_scores,
            name=L.get("satisfaction_score_label", "만족도"),
            marker_color='lightgreen'
        ),
        row=2, col=1
    )

    fig.update_layout(
        height=600,
        showlegend=True,
        margin=dict(l=20, r=20, t=50, b=20),
        barmode='group'
    )
    fig.update_yaxes(title_text="점수", row=2, col=1)
    fig.update_yaxes(title_text="유사도 (%)", row=1, col=1)

    return fig


def visualize_case_trends(histories: List[Dict[str, Any]], current_lang_key: str):
    """과거 성공 사례 트렌드를 시각화"""
    if not IS_PLOTLY_AVAILABLE or not histories:
        return None

    # 언어 키 검증
    if current_lang_key not in ["ko", "en", "ja"]:
        current_lang_key = st.session_state.get("language", "ko")
        if current_lang_key not in ["ko", "en", "ja"]:
            current_lang_key = "ko"
    L = LANG.get(current_lang_key, LANG["ko"])

    # 요약 데이터가 있는 케이스만 필터링
    cases_with_summary = [
        h for h in histories
        if h.get("summary") and isinstance(h.get("summary"), dict) and h.get("is_chat_ended", False)
    ]

    if not cases_with_summary:
        return None

    # 날짜별로 정렬
    cases_with_summary.sort(key=lambda x: x.get("timestamp", ""))

    dates = []
    sentiment_scores = []
    satisfaction_scores = []

    for case in cases_with_summary:
        summary = case.get("summary", {})
        timestamp = case.get("timestamp", "")
        try:
            dt = datetime.fromisoformat(timestamp)
            dates.append(dt)
            sentiment_scores.append(summary.get("customer_sentiment_score", 50))
            satisfaction_scores.append(summary.get("customer_satisfaction_score", 50))
        except Exception:
            continue

    if not dates:
        return None

    # 트렌드 라인 차트
    fig = go.Figure()

    fig.add_trace(go.Scatter(
        x=dates,
        y=sentiment_scores,
        mode='lines+markers',
        name=L.get("sentiment_trend_label", "감정 점수 추이"),
        line=dict(color='lightcoral', width=2),
        marker=dict(size=6)
    ))

    fig.add_trace(go.Scatter(
        x=dates,
        y=satisfaction_scores,
        mode='lines+markers',
        name=L.get("satisfaction_trend_label", "만족도 점수 추이"),
        line=dict(color='lightgreen', width=2),
        marker=dict(size=6)
    ))

    fig.update_layout(
        title=L.get("case_trends_title", "과거 케이스 점수 추이"),
        xaxis_title=L.get("date_label", "날짜"),
        yaxis_title=L.get("score_label", "점수 (0-100)"),
        height=400,
        hovermode='x unified',
        legend=dict(yanchor="top", y=0.99, xanchor="left", x=0.01)
    )

    return fig


def visualize_customer_characteristics(summary: Dict[str, Any], current_lang_key: str):
    """고객 특성을 시각화 (언어, 문화권, 지역 등)"""
    if not IS_PLOTLY_AVAILABLE or not summary:
        return None

    # 언어 키 검증
    if current_lang_key not in ["ko", "en", "ja"]:
        current_lang_key = st.session_state.get("language", "ko")
        if current_lang_key not in ["ko", "en", "ja"]:
            current_lang_key = "ko"
    L = LANG.get(current_lang_key, LANG["ko"])

    characteristics = summary.get("customer_characteristics", {})
    privacy_info = summary.get("privacy_info", {})

    # 특성 데이터 준비
    labels = []
    values = []

    # 언어 정보
    language = characteristics.get("language", "unknown")
    if language != "unknown":
        labels.append(L.get("language_label", "언어"))
        lang_map = {"ko": "한국어", "en": "English", "ja": "日本語"}
        values.append(lang_map.get(language, language))

    # 개인정보 제공 여부
    if privacy_info.get("has_email"):
        labels.append(L.get("email_provided_label", "이메일 제공"))
        values.append("Yes")
    if privacy_info.get("has_phone"):
        labels.append(L.get("phone_provided_label", "전화번호 제공"))
        values.append("Yes")

    # 지역 정보
    region = privacy_info.get("region_hint", characteristics.get("region", "unknown"))
    if region != "unknown":
        labels.append(L.get("region_label", "지역"))
        values.append(region)

    if not labels:
        return None

    # 파이 차트
    fig = go.Figure(data=[go.Pie(
        labels=labels,
        values=[1] * len(labels),
        hole=0.4,
        marker_colors=px.colors.qualitative.Set3[:len(labels)]
    )])

    fig.update_layout(
        title=L.get("customer_characteristics_title",
                    "고객 특성 분포"),
        height=300,
        showlegend=True,
        margin=dict(l=20, r=20, t=50, b=20)
    )

    return fig



def generate_guideline_from_past_cases(customer_query: str, customer_profile: Dict[str, Any],
                                       similar_cases: List[Dict[str, Any]], current_lang_key: str) -> str:
    """과거 유사 케이스의 성공적인 해결 방법을 바탕으로 가이드라인 생성"""
    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}[current_lang_key]

    if not similar_cases:
        return ""

    # 유사 케이스 요약
    past_cases_text = ""
    for idx, similar_case in enumerate(similar_cases, 1):
        case = similar_case["case"]
        summary = similar_case["summary"]
        similarity = similar_case["similarity_score"]

        past_cases_text += f"""
[Case {idx}] (Similarity: {similarity:.1f}%)
- Inquiry: {summary.get('main_inquiry', 'N/A')}
- Customer Sentiment: {summary.get('customer_sentiment_score', 50)}/100
- Customer Satisfaction: {summary.get('customer_satisfaction_score', 50)}/100
- Key Responses: {', '.join(summary.get('key_responses', [])[:3])}
- Summary: {summary.get('summary', 'N/A')[:200]}
"""

    guideline_prompt = f"""
You are an AI Customer Support Supervisor analyzing past successful cases to provide guidance.

Based on the following similar past cases and their successful resolution strategies, provide actionable guidelines for handling the current customer inquiry.

Current Customer Inquiry:
{customer_query}

Current Customer Profile:
- Gender: {customer_profile.get('gender', 'unknown')}
- Sentiment Score: {customer_profile.get('sentiment_score', 50)}/100
- Communication Style: {customer_profile.get('communication_style', 'unknown')}
- Urgency: {customer_profile.get('urgency_level', 'medium')}
- Predicted Type: {customer_profile.get('predicted_customer_type', 'normal')}

Similar Past Cases (Successful Resolutions):
{past_cases_text}

Provide a concise guideline in {lang_name} that:
1. Identifies what worked well in similar past cases
2. Suggests specific approaches based on successful patterns
3. Warns about potential pitfalls based on past experiences
4. Recommends response strategies that led to high customer satisfaction

Guideline (in {lang_name}):
"""

    if not st.session_state.is_llm_ready:
        return ""

    try:
        guideline = run_llm(guideline_prompt).strip()
        return guideline
    except Exception as e:
        return f"가이드라인 생성 오류: {str(e)}"



def _generate_initial_advice(customer_query, customer_type_display, customer_email, customer_phone, current_lang_key,
                             customer_attachment_file):
    """Supervisor 가이드라인과 초안을 생성하는 함수 (저장된 데이터 활용)"""
    # 입력 텍스트의 언어를 자동 감지 (오류 발생 시 안전하게 처리)
    try:
        detected_lang = detect_text_language(customer_query)
    except Exception as e:
        print(f"Language detection failed in _generate_initial_advice: {e}")
        detected_lang = current_lang_key if current_lang_key else "ko"
    
    # 감지된 언어를 우선 사용하되, current_lang_key가 명시적으로 제공되면 그것을 사용
    lang_key_to_use = detected_lang if detected_lang else current_lang_key
    # lang_key_to_use가 유효한지 확인
    if lang_key_to_use not in ["ko", "en", "ja"]:
        lang_key_to_use = current_lang_key if current_lang_key else "ko"
    
    # 언어 키 검증
    if lang_key_to_use not in ["ko", "en", "ja"]:
        lang_key_to_use = st.session_state.get("language", "ko")
        if lang_key_to_use not in ["ko", "en", "ja"]:
            lang_key_to_use = "ko"
    L = LANG.get(lang_key_to_use, LANG["ko"])
    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}[lang_key_to_use]

    contact_info_block = ""
    if customer_email or customer_phone:
        contact_info_block = (
            f"\n\n[Customer contact info for reference (DO NOT use these in your reply draft!)]"
            f"\n- Email: {customer_email or 'N/A'}"
            f"\n- Phone: {customer_phone or 'N/A'}"
        )

    attachment_block = ""
    if customer_attachment_file:
        file_name = customer_attachment_file.name
        attachment_block = f"\n\n[ATTACHMENT NOTE]: {L['attachment_info_llm'].format(filename=file_name)}"

    # 고객 프로필 분석 (감지된 언어 사용)
    customer_profile = analyze_customer_profile(customer_query, lang_key_to_use)

    # 유사 케이스 찾기 (감지된 언어 사용)
    similar_cases = find_similar_cases(customer_query, customer_profile, lang_key_to_use, limit=5)

    # 과거 케이스 기반 가이드라인 생성 (감지된 언어 사용)
    past_cases_guideline = ""
    if similar_cases:
        past_cases_guideline = generate_guideline_from_past_cases(
            customer_query, customer_profile, similar_cases, lang_key_to_use
        )

    # 고객 프로필 정보
    gender_display = customer_profile.get('gender', 'unknown')
    profile_block = f"""
[Customer Profile Analysis]
- Gender: {gender_display}
- Sentiment Score: {customer_profile.get('sentiment_score', 50)}/100
- Communication Style: {customer_profile.get('communication_style', 'unknown')}
- Urgency Level: {customer_profile.get('urgency_level', 'medium')}
- Predicted Type: {customer_profile.get('predicted_customer_type', 'normal')}
- Key Concerns: {', '.join(customer_profile.get('key_concerns', []))}
- Tone: {customer_profile.get('tone_analysis', 'unknown')}
"""

    # 과거 케이스 기반 가이드라인 블록
    past_cases_block = ""
    if past_cases_guideline:
        past_cases_block = f"""
[Guidelines Based on {len(similar_cases)} Similar Past Cases]
{past_cases_guideline}
"""
    elif similar_cases:
        past_cases_block = f"""
[Note: Found {len(similar_cases)} similar past cases, but unable to generate detailed guidelines.
Consider reviewing past cases manually for patterns.]
"""

    # Output ALL text (guidelines and draft) STRICTLY in {lang_name}. <--- 강력한 언어 강제 지시
    initial_prompt = f"""
Output ALL text (guidelines and draft) STRICTLY in {lang_name}.

You are an AI Customer Support Supervisor. Your role is to analyze the following customer inquiry
from a **{st.session_state.customer_type_sim_select}** and provide:

1) A detailed **response guideline for the human agent** (step-by-step).
2) A **ready-to-send draft reply** in {lang_name}.

[FORMAT]
- Use the exact markdown headers:
  - "### {L['simulation_advice_header']}"
  - "### {L['simulation_draft_header']}"

[CRITICAL GUIDELINE RULES]
1. **Initial Information Collection (Req 3):** The first step in the guideline MUST be to request the necessary initial diagnostic information (e.g., device compatibility, local status/location, order number) BEFORE attempting to troubleshoot or solve the problem.
2. **Empathy for Difficult Customers (Req 5):** If the customer type is 'Difficult Customer' or 'Highly Dissatisfied Customer', the guideline MUST emphasize extreme politeness, empathy, and apologies, even if the policy (e.g., no refund) must be enforced.
3. **24-48 Hour Follow-up (Req 6):** If the issue cannot be solved immediately or requires confirmation from a local partner/supervisor, the guideline MUST state the procedure:
   - Acknowledge the issue.
   - Inform the customer they will receive a definite answer within 24 or 48 hours.
   - Request the customer's email or phone number for follow-up contact. (Use provided contact info if available)
4. **Past Cases Learning:** If past cases guidelines are provided, incorporate successful strategies from those cases into your recommendations.

Customer Inquiry:
{customer_query}
{contact_info_block}
{attachment_block}
{profile_block}
{past_cases_block}
"""
    if not st.session_state.is_llm_ready:
        mock_text = (
            f"### {L['simulation_advice_header']}\n\n"
            f"- (Mock) {st.session_state.customer_type_sim_select} 유형 고객 응대 가이드입니다. (요청 3, 5, 6 반영)\n\n"
            f"### {L['simulation_draft_header']}\n\n"
            f"(Mock) 에이전트 응대 초안이 여기에 들어갑니다。\n\n"
        )
        return mock_text
    else:
        with st.spinner(L["response_generating"]):
            try:
                return run_llm(initial_prompt)
            except Exception as e:
                st.error(f"AI 조언 생성 중 오류 발생: {e}")
                return f"❌ AI Advice Generation Error: {e}"


# ========================================
# 고객 검증 관련 함수
# ========================================

def mask_email(email: str, show_chars: int = 2) -> str:
    """
    이메일 주소를 마스킹합니다.
    앞/뒤 show_chars 자리만 표시하고 나머지는 * 처리합니다.
    
    Args:
        email: 마스킹할 이메일 주소
        show_chars: 앞/뒤에 표시할 문자 수 (기본값: 2)
    
    Returns:
        마스킹된 이메일 주소 (예: "ab***@ex***.com")
    """
    if not email or "@" not in email:
        return email
    
    local_part, domain = email.split("@", 1)
    domain_parts = domain.split(".", 1)
    
    # 로컬 부분 마스킹 (앞 show_chars 자리만 표시)
    if len(local_part) <= show_chars:
        masked_local = local_part
    else:
        masked_local = local_part[:show_chars] + "*" * (len(local_part) - show_chars)
    
    # 도메인 부분 마스킹 (앞 show_chars 자리만 표시)
    if len(domain_parts[0]) <= show_chars:
        masked_domain = domain_parts[0]
    else:
        masked_domain = domain_parts[0][:show_chars] + "*" * (len(domain_parts[0]) - show_chars)
    
    if len(domain_parts) > 1:
        return f"{masked_local}@{masked_domain}.{domain_parts[1]}"
    else:
        return f"{masked_local}@{masked_domain}"


def verify_customer_info(
    provided_info: Dict[str, str],
    stored_info: Dict[str, str]
) -> Tuple[bool, Dict[str, bool]]:
    """
    고객이 제공한 정보와 저장된 정보를 비교하여 검증합니다.
    시스템 내부에서만 실행되며, confidential 정보를 보호합니다.
    
    Args:
        provided_info: 고객이 제공한 검증 정보
            - receipt_number: 영수증/예약 번호
            - card_last4: 카드 뒷자리 4자리 (카드 결제인 경우)
            - account_number: 계좌번호 (온라인뱅킹인 경우)
            - payment_method: 결제 수단
            - customer_name: 고객 성함
            - customer_email: 고객 이메일
            - customer_phone: 고객 연락처
            - file_uploaded: 파일 업로드 여부
        stored_info: 시스템에 저장된 검증 정보 (동일한 구조)
    
    Returns:
        (전체 검증 성공 여부, 각 필드별 검증 결과)
    """
    verification_results = {
        "receipt_number": False,
        "payment_info": False,  # 카드 뒷자리 또는 계좌번호
        "customer_name": False,
        "customer_email": False,
        "customer_phone": False,
        "file_uploaded": False
    }
    
    # 파일 업로드 확인 (파일이 있으면 추가 점수)
    if provided_info.get("file_uploaded"):
        verification_results["file_uploaded"] = True
    
    # 영수증/예약 번호 검증 (대소문자 무시, 공백 제거)
    if provided_info.get("receipt_number") and stored_info.get("receipt_number"):
        provided_receipt = provided_info["receipt_number"].strip().upper().replace(" ", "")
        stored_receipt = stored_info["receipt_number"].strip().upper().replace(" ", "")
        verification_results["receipt_number"] = provided_receipt == stored_receipt
    
    # 결제 정보 검증 (결제 수단에 따라 다르게 처리)
    payment_method = provided_info.get("payment_method", "")
    
    # 카드 결제인 경우
    if "카드" in payment_method or "card" in payment_method.lower():
        if provided_info.get("card_last4") and stored_info.get("card_last4"):
            provided_card = "".join(filter(str.isdigit, provided_info["card_last4"]))[-4:]
            stored_card = "".join(filter(str.isdigit, stored_info["card_last4"]))[-4:]
            verification_results["payment_info"] = provided_card == stored_card and len(provided_card) == 4
    
    # 온라인뱅킹인 경우
    elif "온라인뱅킹" in payment_method or "online banking" in payment_method.lower():
        if provided_info.get("account_number") and stored_info.get("account_number"):
            provided_account = "".join(filter(str.isdigit, provided_info["account_number"]))
            stored_account = "".join(filter(str.isdigit, stored_info["account_number"]))
            # 계좌번호는 마지막 4-6자리 비교
            verification_results["payment_info"] = provided_account[-6:] == stored_account[-6:] or provided_account[-4:] == stored_account[-4:]
    
    # 카카오페이, 네이버페이, GrabPay, Touch N Go 등은 결제 수단 정보만으로 확인 가능
    elif payment_method and stored_info.get("payment_method"):
        provided_payment = payment_method.strip().lower()
        stored_payment = stored_info["payment_method"].strip().lower()
        verification_results["payment_info"] = provided_payment == stored_payment
    
    # 고객 성함 검증 (대소문자 무시, 공백 정규화)
    if provided_info.get("customer_name") and stored_info.get("customer_name"):
        provided_name = " ".join(provided_info["customer_name"].strip().split()).upper()
        stored_name = " ".join(stored_info["customer_name"].strip().split()).upper()
        verification_results["customer_name"] = provided_name == stored_name
    
    # 이메일 검증 (대소문자 무시)
    if provided_info.get("customer_email") and stored_info.get("customer_email"):
        provided_email = provided_info["customer_email"].strip().lower()
        stored_email = stored_info["customer_email"].strip().lower()
        verification_results["customer_email"] = provided_email == stored_email
    
    # 연락처 검증 (숫자만 추출하여 비교)
    if provided_info.get("customer_phone") and stored_info.get("customer_phone"):
        provided_phone = "".join(filter(str.isdigit, provided_info["customer_phone"]))
        stored_phone = "".join(filter(str.isdigit, stored_info["customer_phone"]))
        # 마지막 4-10자리 비교 (국가코드 제외)
        verification_results["customer_phone"] = provided_phone[-10:] == stored_phone[-10:] or provided_phone[-4:] == stored_phone[-4:]
    
    # 전체 검증 성공 여부: 최소 3개 이상 일치해야 함 (파일 업로드 시 추가 점수)
    matched_count = sum(verification_results.values())
    # 파일이 업로드된 경우 2점 추가
    if verification_results["file_uploaded"]:
        matched_count += 1
    is_verified = matched_count >= 3
    
    return is_verified, verification_results


def check_if_customer_provided_verification_info(messages: List[Dict[str, Any]]) -> bool:
    """
    고객이 검증 정보를 제공했는지 확인합니다.
    고객 메시지에서 영수증, 예약번호, 카드, 결제수단, 성함, 이메일, 연락처 등의 키워드가 있는지 확인합니다.
    
    Args:
        messages: 대화 메시지 리스트
    
    Returns:
        검증 정보 제공 여부
    """
    if not messages:
        return False
    
    # 최근 고객 메시지 확인 (최근 10개까지 확인)
    # 모든 메시지를 확인하여 고객 관련 메시지 추출
    recent_customer_messages = []
    for msg in messages[-10:]:
        role = msg.get("role", "")
        content = msg.get("content", "")
        # 고객 역할 확인 (더 포괄적으로)
        if role in ["customer", "customer_rebuttal", "initial_query"] or "customer" in role.lower():
            recent_customer_messages.append(content)
    
    # 메시지가 없으면 False
    if not recent_customer_messages:
        # 디버깅: 메시지가 없는 경우 로그
        print(f"[DEBUG] No customer messages found. Total messages: {len(messages)}")
        if messages:
            print(f"[DEBUG] Available roles: {[msg.get('role') for msg in messages[-5:]]}")
        return False
    
    # 모든 고객 메시지를 하나의 텍스트로 결합 (원본 대소문자 유지)
    combined_text_original = " ".join(recent_customer_messages)
    combined_text = combined_text_original.lower()
    
    import re
    
    # 숫자 패턴 확인 (예약 번호, 전화번호 등)
    has_numbers = bool(re.search(r'\d{4,}', combined_text))  # 4자리 이상 숫자
    
    # 이메일 패턴 확인
    has_email = bool(re.search(r'[\w\.-]+@[\w\.-]+\.\w+', combined_text_original))
    
    # 전화번호 패턴 확인 (010-1234-5678, 010 1234 5678, 01012345678 등)
    has_phone = bool(re.search(r'010[-.\s]?\d{3,4}[-.\s]?\d{4}', combined_text_original) or 
                     re.search(r'010\d{8}', combined_text_original))
    
    # 최소 2개 이상의 검증 정보가 있어야 함
    info_count = 0
    
    # 1. 예약/영수증 번호 확인 (더 포괄적으로)
    if (re.search(r'예약\s*번호', combined_text) or 
        re.search(r'영수증\s*번호', combined_text) or
        re.search(r'예약.*[:：]\s*\d{4,}', combined_text_original) or
        re.search(r'영수증.*[:：]\s*\d{4,}', combined_text_original) or
        re.search(r'예약번호.*[:：]\s*\d{4,}', combined_text_original) or
        re.search(r'예약.*\d{4,}', combined_text) or
        re.search(r'영수증.*\d{4,}', combined_text) or
        re.search(r'booking.*number', combined_text) or
        re.search(r'receipt.*number', combined_text) or
        re.search(r'\d{5,}', combined_text_original)):  # 5자리 이상 숫자도 예약번호로 간주
        info_count += 1
    
    # 2. 결제 수단 확인 (카드, 카카오페이, 네이버페이, VISA, Master 등)
    payment_keywords = [
        "카드", "card", "visa", "master", "amex", "american express",
        "신용카드", "체크카드", "credit card", "debit card",
        "카카오페이", "kakao", "kakaopay",
        "네이버페이", "naver", "naverpay",
        "온라인뱅킹", "online banking", "online",
        "grabpay", "grab pay", "grab",
        "touch n go", "touch n' go", "tng",
        "결제 수단", "payment method", "payment", "결제하", "결제 내역"
    ]
    # 결제 수단 키워드가 있거나, "결제 수단 :" 같은 패턴이 있는 경우
    if (any(kw in combined_text for kw in payment_keywords) or
        re.search(r'결제\s*수단\s*[:：]', combined_text_original)):
        info_count += 1
    
    # 3. 성함 확인 (이름, 성함 키워드 + 한글/영문 이름 패턴)
    name_keywords = ["성함", "이름", "name", "제 이름", "고객님의 성함", "my name", "고객님의 이름"]
    # 한글 이름 패턴 (2-4자 한글, 성함/이름 키워드 뒤에 오는 경우)
    korean_name_pattern = (
        re.search(r'성함\s*[:：]\s*[가-힣]{2,4}', combined_text_original) or
        re.search(r'이름\s*[:：]\s*[가-힣]{2,4}', combined_text_original) or
        re.search(r'고객님의\s*성함\s*[:：]\s*[가-힣]{2,4}', combined_text_original) or
        re.search(r'[가-힣]{2,4}', combined_text_original)  # 단순 한글 이름 패턴도 확인
    )
    # 영문 이름 패턴
    english_name_pattern = re.search(r'\b[A-Z][a-z]+\s+[A-Z][a-z]+\b', combined_text_original)
    
    if (any(kw in combined_text for kw in name_keywords) or 
        korean_name_pattern or 
        english_name_pattern):
        info_count += 1
    
    # 4. 이메일 확인
    if has_email:
        info_count += 1
    
    # 5. 연락처 확인
    if (has_phone or 
        re.search(r'연락처', combined_text) or 
        re.search(r'전화번호', combined_text) or
        re.search(r'phone', combined_text)):
        info_count += 1
    
    # 최소 2개 이상의 정보가 제공되었거나, 이메일/전화번호 중 하나와 다른 정보가 함께 있는 경우
    # 또는 숫자(예약번호)와 결제수단/성함이 함께 있는 경우
    result = info_count >= 2 or (has_email and info_count >= 1) or (has_phone and info_count >= 1)
    
    # 디버깅: 감지 결과 로그
    print(f"[DEBUG] Verification info detection:")
    print(f"  - Combined text (first 200 chars): {combined_text_original[:200]}")
    print(f"  - Info count: {info_count}")
    print(f"  - Has email: {has_email}")
    print(f"  - Has phone: {has_phone}")
    print(f"  - Has numbers: {has_numbers}")
    print(f"  - Result: {result}")
    
    return result


def check_if_login_related_inquiry(customer_query: str) -> bool:
    """
    고객 문의가 로그인/계정 관련인지 확인합니다.
    
    Args:
        customer_query: 고객 문의 내용
    
    Returns:
        로그인/계정 관련 문의인지 여부
    """
    if not customer_query or not customer_query.strip():
        return False
    
    login_keywords = [
        "로그인", "login", "ログイン",
        "계정", "account", "アカウント",
        "비밀번호", "password", "パスワード",
        "아이디", "id", "ID", "ユーザーID",
        "접속", "access", "アクセス",
        "인증", "authentication", "認証",
        "로그인 안됨", "cannot login", "ログインできない",
        "로그인 오류", "login error", "ログインエラー",
        "로그인 안", "로그인 실패", "로그인 문제",
        "계정 문제", "계정 잠금", "계정 오류",
        "account problem", "account error", "account locked",
        "password reset", "비밀번호 재설정", "パスワードリセット",
        "forgot password", "비밀번호 분실", "パスワード忘れ"
    ]
    
    query_lower = customer_query.lower()
    # 각 키워드를 개별적으로 확인 (부분 문자열 매칭)
    for keyword in login_keywords:
        if keyword.lower() in query_lower:
            return True
    
    return False


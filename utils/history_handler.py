# Copyright (c) Streamlit Inc. (2018-2022) Snowflake Inc. (2022-2025)
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""
시뮬레이션 이력 관리 모듈
로컬 저장, 요약 생성, 가이드 생성 등의 이력 관련 기능을 제공합니다.
"""

import os
import json
import uuid
import hashlib
from datetime import datetime
from typing import List, Dict, Any, Optional
import streamlit as st

from config import SIM_META_FILE, DATA_DIR
from utils import _load_json, _save_json
from llm_client import get_api_key, run_llm
from lang_pack import LANG

# Word, PPTX, PDF 내보내기 라이브러리
try:
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    IS_DOCX_AVAILABLE = True
except ImportError:
    IS_DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    IS_PPTX_AVAILABLE = True
except ImportError:
    IS_PPTX_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import black
    from reportlab.lib.units import inch
    IS_REPORTLAB_AVAILABLE = True
except ImportError:
    IS_REPORTLAB_AVAILABLE = False


def load_simulation_histories_local(lang_key: str) -> List[Dict[str, Any]]:
    """로컬에 저장된 시뮬레이션 이력을 로드합니다."""
    histories = _load_json(SIM_META_FILE, [])
    return [
        h for h in histories
        if h.get("language_key") == lang_key and (isinstance(h.get("messages"), list) or h.get("summary"))
    ]


def generate_call_summary(messages: List[Dict[str, Any]], initial_query: str, customer_type: str,
                          current_lang_key: str) -> Dict[str, Any]:
    """전화 통화 이력 요약 생성 (문의 내용 + 솔루션 요점만, 다국어 지원)"""
    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}[current_lang_key]
    
    # 고객 문의 내용 추출
    customer_inquiries = []
    agent_solutions = []
    
    for msg in messages:
        role = msg.get("role", "")
        content = msg.get("content", "")
        if role in ["customer", "customer_rebuttal"]:
            customer_inquiries.append(content)
        elif role in ["agent", "agent_response"]:
            agent_solutions.append(content)
    
    # 요약 프롬프트 (문의 내용 + 솔루션 요점만)
    summary_prompt = f"""
You are an AI analyst summarizing a phone call customer support conversation. Generate a concise summary focusing ONLY on:
1. Customer inquiry/inquiries (what the customer asked about)
2. Key solutions provided by the agent (main points only, not full scripts)

The summary MUST be in {lang_name} language.

Conversation:
Initial Query: {initial_query}

Customer Messages:
{chr(10).join([f"- {inq}" for inq in customer_inquiries[:5]])}

Agent Responses:
{chr(10).join([f"- {sol}" for sol in agent_solutions[:5]])}

Generate a JSON summary with the following structure (ONLY JSON, no markdown):
{{
    "customer_inquiry": "Brief summary of what the customer asked about",
    "key_solutions": ["Solution point 1", "Solution point 2", "Solution point 3"],
    "summary": "Brief overall summary in {lang_name}"
}}

IMPORTANT:
- Keep it concise (not full scripts)
- Focus on main inquiry and solution points
- Use {lang_name} language
- Maximum 3 solution points
"""
    
    try:
        if get_api_key("gemini") or get_api_key("openai"):
            summary_json = run_llm(summary_prompt)
            # JSON 추출
            if summary_json.strip().startswith("```"):
                summary_json = summary_json.strip().split("```")[1]
                if summary_json.startswith("json"):
                    summary_json = summary_json[4:]
            summary_json = summary_json.strip()
            
            import json
            summary_data = json.loads(summary_json)
            
            # 다국어 번역 추가
            summary_data["summary_ko"] = summary_data.get("summary", "")
            summary_data["summary_en"] = ""
            summary_data["summary_ja"] = ""
            
            # 영어 번역
            if current_lang_key != "en":
                try:
                    from utils.translation import translate_text_with_llm
                    summary_data["summary_en"] = translate_text_with_llm(
                        summary_data.get("summary", ""), "en", current_lang_key
                    ) or ""
                except:
                    pass
            
            # 일본어 번역
            if current_lang_key != "ja":
                try:
                    from utils.translation import translate_text_with_llm
                    summary_data["summary_ja"] = translate_text_with_llm(
                        summary_data.get("summary", ""), "ja", current_lang_key
                    ) or ""
                except:
                    pass
            
            return summary_data
        else:
            return {
                "customer_inquiry": initial_query,
                "key_solutions": [],
                "summary": f"Phone call conversation about {initial_query}",
                "summary_ko": "",
                "summary_en": "",
                "summary_ja": ""
            }
    except Exception as e:
        return {
            "customer_inquiry": initial_query,
            "key_solutions": [],
            "summary": f"Error generating summary: {str(e)}",
            "summary_ko": "",
            "summary_en": "",
            "summary_ja": ""
        }


def generate_chat_summary(messages: List[Dict[str, Any]], initial_query: str, customer_type: str,
                          current_lang_key: str) -> Dict[str, Any]:
    """채팅 내용을 AI로 요약하여 주요 정보와 점수를 추출"""
    lang_name = {"ko": "Korean", "en": "English", "ja": "Japanese"}[current_lang_key]

    conversation_text = f"Initial Query: {initial_query}\n\n"
    for msg in messages:
        role = msg.get("role", "")
        content = msg.get("content", "")
        if role in ["customer", "customer_rebuttal", "phone_exchange"]:
            conversation_text += f"Customer: {content}\n"
        elif role == "agent_response" or role == "agent":
            conversation_text += f"Agent: {content}\n"

    summary_prompt = f"""
You are an AI analyst summarizing a customer support conversation. Your task is to extract comprehensive customer profile data and score various aspects numerically.

Analyze the conversation and provide a structured summary in JSON format (ONLY JSON, no markdown).

Extract and score:
1. Main inquiry topic (what the customer asked about)
2. Key responses provided by the agent (list of max 3 core actions/solutions)
3. Customer sentiment score (0-100, where 0=very negative, 50=neutral, 100=very positive)
4. Customer satisfaction score (0-100, based on final response)
5. Customer characteristics with detailed scoring:
   - Language preference (detected language code: ko/en/ja)
   - Cultural background hints (score 0-100, where higher = more cultural context detected)
   - Location/region (general region only, anonymize specific addresses)
   - Communication style (formal/casual, brief/detailed) with scores:
     * Formality score (0-100, 0=casual, 100=very formal)
     * Detail level score (0-100, 0=brief, 100=very detailed)
   - Customer personality traits (score each 0-100):
     * Patience level (0-100)
     * Assertiveness (0-100)
     * Politeness level (0-100)
     * Technical proficiency (0-100, if technical inquiry)
6. Privacy-sensitive information (anonymize: names, emails, phone numbers, specific addresses)
   - Extract patterns only (e.g., "email provided", "phone number provided", "resides in Asia region")
7. Customer behavior patterns:
   - Response time pattern (fast/moderate/slow based on message frequency)
   - Question complexity (simple/moderate/complex)
   - Escalation tendency (0-100, likelihood to escalate)

Output format (JSON only):
{{
  "main_inquiry": "brief description of main issue",
  "key_responses": ["response 1", "response 2"],
  "customer_sentiment_score": 75,
  "customer_satisfaction_score": 80,
  "customer_characteristics": {{
    "language": "ko/en/ja or unknown",
    "cultural_hints": "brief description or unknown",
    "cultural_score": 60,
    "region": "general region or unknown",
    "communication_style": "formal/casual/brief/detailed",
    "formality_score": 70,
    "detail_level_score": 65,
    "personality_traits": {{
      "patience_level": 60,
      "assertiveness": 70,
      "politeness_level": 80,
      "technical_proficiency": 50
    }}
  }},
  "privacy_info": {{
    "has_email": true/false,
    "has_phone": true/false,
    "has_address": true/false,
    "region_hint": "general region or unknown"
  }},
  "behavior_patterns": {{
    "response_time": "fast/moderate/slow",
    "question_complexity": "simple/moderate/complex",
    "escalation_tendency": 30
  }},
  "summary": "overall conversation summary in {lang_name}"
}}

Conversation:
{conversation_text}

JSON Output:
"""

    if not st.session_state.is_llm_ready:
        return {
            "main_inquiry": initial_query[:100],
            "key_responses": [],
            "customer_sentiment_score": 50,
            "customer_satisfaction_score": 50,
            "customer_characteristics": {
                "language": current_lang_key,
                "cultural_hints": "unknown",
                "cultural_score": 0,
                "region": "unknown",
                "communication_style": "unknown",
                "formality_score": 50,
                "detail_level_score": 50,
                "personality_traits": {
                    "patience_level": 50,
                    "assertiveness": 50,
                    "politeness_level": 50,
                    "technical_proficiency": 50
                }
            },
            "privacy_info": {
                "has_email": False,
                "has_phone": False,
                "has_address": False,
                "region_hint": "unknown"
            },
            "behavior_patterns": {
                "response_time": "moderate",
                "question_complexity": "moderate",
                "escalation_tendency": 50
            },
            "summary": f"Customer inquiry about: {initial_query[:100]}"
        }

    try:
        summary_text = run_llm(summary_prompt).strip()
        if "```json" in summary_text:
            summary_text = summary_text.split("```json")[1].split("```")[0].strip()
        elif "```" in summary_text:
            summary_text = summary_text.split("```")[1].split("```")[0].strip()

        summary_data = json.loads(summary_text)
        return summary_data
    except Exception as e:
        st.warning(f"요약 생성 중 오류 발생: {e}")
        return {
            "main_inquiry": initial_query[:100],
            "key_responses": [],
            "customer_sentiment_score": 50,
            "customer_satisfaction_score": 50,
            "customer_characteristics": {
                "language": current_lang_key,
                "cultural_hints": "unknown",
                "cultural_score": 0,
                "region": "unknown",
                "communication_style": "unknown",
                "formality_score": 50,
                "detail_level_score": 50,
                "personality_traits": {
                    "patience_level": 50,
                    "assertiveness": 50,
                    "politeness_level": 50,
                    "technical_proficiency": 50
                }
            },
            "privacy_info": {
                "has_email": False,
                "has_phone": False,
                "has_address": False,
                "region_hint": "unknown"
            },
            "behavior_patterns": {
                "response_time": "moderate",
                "question_complexity": "moderate",
                "escalation_tendency": 50
            },
            "summary": f"Error generating summary: {str(e)}"
        }


def save_simulation_history_local(initial_query: str, customer_type: str, messages: List[Dict[str, Any]],
                                  is_chat_ended: bool, attachment_context: str, is_call: bool = False,
                                  customer_name: str = "", customer_phone: str = "", customer_email: str = "",
                                  customer_id: str = ""):
    """AI 요약 데이터를 중심으로 이력을 저장 (고객 정보 포함)"""
    histories = _load_json(SIM_META_FILE, [])
    doc_id = str(uuid.uuid4())
    ts = datetime.utcnow().isoformat()

    summary_data = None
    if is_chat_ended or len(messages) > 4 or is_call:
        # 전화 통화의 경우 요약 생성 (문의 내용 + 솔루션 요점만)
        if is_call:
            summary_data = generate_call_summary(messages, initial_query, customer_type, st.session_state.language)
        else:
            # 채팅의 경우 기존 요약 함수 사용
            summary_data = generate_chat_summary(messages, initial_query, customer_type, st.session_state.language)

    if summary_data:
        data = {
            "id": doc_id,
            "initial_query": initial_query,
            "customer_type": customer_type,
            "messages": [],
            "summary": summary_data,
            "language_key": st.session_state.language,
            "timestamp": ts,
            "is_chat_ended": is_chat_ended,
            "attachment_context": attachment_context if attachment_context else "",
            "is_call": is_call,
        }
    else:
        data = {
            "id": doc_id,
            "initial_query": initial_query,
            "customer_type": customer_type,
            "messages": messages[:10] if len(messages) > 10 else messages,
            "summary": None,
            "language_key": st.session_state.language,
            "timestamp": ts,
            "is_chat_ended": is_chat_ended,
            "attachment_context": attachment_context if attachment_context else "",
            "is_call": is_call,
        }

    histories.insert(0, data)
    
    today_str = datetime.utcnow().strftime("%Y-%m-%d")
    today_histories = [h for h in histories if h.get("timestamp", "").startswith(today_str)]
    
    if len(today_histories) > 20:
        today_histories_sorted = sorted(today_histories, key=lambda x: x.get("timestamp", ""))
        excess_count = len(today_histories) - 20
        excess_ids = {h.get("id") for h in today_histories_sorted[:excess_count]}
        histories = [h for h in histories if h.get("id") not in excess_ids]
    
    _save_json(SIM_META_FILE, histories[:100])
    
    if summary_data and is_chat_ended:
        try:
            all_histories = _load_json(SIM_META_FILE, [])
            today_str = datetime.now().strftime("%y%m%d")
            guide_filename = f"{today_str}_고객가이드.TXT"
            guide_filepath = os.path.join(DATA_DIR, guide_filename)
            
            guide_content = generate_daily_customer_guide(all_histories, st.session_state.language)
            
            if guide_content:
                saved_path = save_daily_customer_guide(guide_content, st.session_state.language)
                if saved_path:
                    print(f"✅ 고객 가이드가 자동 생성/업데이트되었습니다: {saved_path}")
        except Exception as e:
            print(f"고객 가이드 자동 생성 중 오류 발생 (무시됨): {e}")
    
    return doc_id


def delete_all_history_local():
    """모든 이력을 삭제합니다."""
    _save_json(SIM_META_FILE, [])


def get_daily_data_statistics(language: str = "ko") -> Dict[str, Any]:
    """일일 데이터 수집 통계를 반환합니다."""
    histories = _load_json(SIM_META_FILE, [])
    today = datetime.now().date()
    
    today_histories = []
    for h in histories:
        try:
            ts = datetime.fromisoformat(h.get("timestamp", "")).date()
            if ts == today and h.get("summary") and isinstance(h.get("summary"), dict):
                today_histories.append(h)
        except:
            continue
    
    unique_customers = set()
    for h in today_histories:
        customer_id = h.get("id", "")
        initial_query = h.get("initial_query", "")
        customer_hash = hashlib.md5(f"{customer_id}_{initial_query[:50]}".encode()).hexdigest()
        unique_customers.add(customer_hash)
    
    return {
        "date": today.isoformat(),
        "total_cases": len(today_histories),
        "unique_customers": len(unique_customers),
        "target_met": len(unique_customers) >= 5,
        "cases_with_summary": len([h for h in today_histories if h.get("summary")])
    }


def recommend_guideline_for_customer(new_customer_summary: Dict[str, Any], histories: List[Dict[str, Any]], language: str = "ko") -> Optional[str]:
    """신규 고객의 문의사항과 말투 등을 종합하여 고객 성향 점수를 수치화하고, 저장된 데이터를 바탕으로 최적의 가이드라인을 추천합니다."""
    if not histories or not get_api_key("gemini") and not get_api_key("openai"):
        return None
    
    try:
        similar_customers = []
        new_scores = {
            "sentiment": new_customer_summary.get("customer_sentiment_score", 50),
            "satisfaction": new_customer_summary.get("customer_satisfaction_score", 50),
            "formality": new_customer_summary.get("customer_characteristics", {}).get("formality_score", 50),
            "patience": new_customer_summary.get("customer_characteristics", {}).get("personality_traits", {}).get("patience_level", 50),
            "assertiveness": new_customer_summary.get("customer_characteristics", {}).get("personality_traits", {}).get("assertiveness", 50),
        }
        
        for h in histories:
            if not h.get("summary") or not isinstance(h.get("summary"), dict):
                continue
            
            summary = h["summary"]
            old_scores = {
                "sentiment": summary.get("customer_sentiment_score", 50),
                "satisfaction": summary.get("customer_satisfaction_score", 50),
                "formality": summary.get("customer_characteristics", {}).get("formality_score", 50),
                "patience": summary.get("customer_characteristics", {}).get("personality_traits", {}).get("patience_level", 50),
                "assertiveness": summary.get("customer_characteristics", {}).get("personality_traits", {}).get("assertiveness", 50),
            }
            
            similarity = sum(abs(new_scores[k] - old_scores[k]) for k in new_scores.keys())
            
            if similarity < 100:
                similar_customers.append({
                    "history": h,
                    "similarity": similarity,
                    "scores": old_scores
                })
        
        similar_customers.sort(key=lambda x: x["similarity"])
        
        if similar_customers:
            lang_name = {"ko": "한국어", "en": "English", "ja": "日本語"}.get(language, "한국어")
            
            similar_cases_text = json.dumps([
                {
                    "initial_query": c["history"].get("initial_query", ""),
                    "key_responses": c["history"].get("summary", {}).get("key_responses", []),
                    "scores": c["scores"],
                    "satisfaction": c["history"].get("summary", {}).get("customer_satisfaction_score", 50)
                }
                for c in similar_customers[:5]
            ], ensure_ascii=False, indent=2)
            
            recommendation_prompt = (
                f"당신은 CS 센터 전문가입니다. 신규 고객의 성향 점수를 분석하고, 유사한 과거 고객들의 성공 사례를 바탕으로 최적의 응대 가이드라인을 추천하세요.\n\n"
                f"신규 고객 프로필:\n{json.dumps(new_customer_summary, ensure_ascii=False, indent=2)}\n\n"
                f"유사한 과거 고객 사례 (상위 5개):\n{similar_cases_text}\n\n"
                f"다음 내용을 포함하여 {lang_name}로 가이드라인을 작성하세요:\n"
                f"1. 고객 성향 분석 (점수 기반)\n"
                f"2. 예상되는 고객 반응 패턴\n"
                f"3. 효과적인 응대 전략 (유사 사례 기반)\n"
                f"4. 주의해야 할 사항\n"
                f"5. 권장 응대 톤 및 스타일\n\n"
                f"실용적이고 구체적인 가이드라인을 제공하세요."
            )
            
            recommendation = run_llm(recommendation_prompt)
            return recommendation if recommendation and not recommendation.startswith("❌") else None
        
        return None
        
    except Exception as e:
        print(f"가이드라인 추천 중 오류 발생: {e}")
        return None


def generate_daily_customer_guide(histories: List[Dict[str, Any]], language: str = "ko") -> Optional[str]:
    """일일 고객 가이드 생성 함수"""
    if not histories or not get_api_key("gemini") and not get_api_key("openai"):
        return None
    
    try:
        histories_with_summary = [h for h in histories if h.get("summary") and isinstance(h.get("summary"), dict)]
        
        if not histories_with_summary:
            return None
        
        recent_histories = histories_with_summary[:50]
        
        customer_data_map = {}
        for h in recent_histories:
            customer_id = h.get("id", "")
            customer_type = h.get("customer_type", "")
            summary = h.get("summary", {})
            
            if customer_id not in customer_data_map:
                customer_data_map[customer_id] = {
                    "customer_type": customer_type,
                    "histories": [],
                    "total_interactions": 0
                }
            
            customer_data_map[customer_id]["histories"].append({
                "initial_query": h.get("initial_query", ""),
                "summary": summary,
                "timestamp": h.get("timestamp", ""),
                "language": h.get("language_key", language)
            })
            customer_data_map[customer_id]["total_interactions"] += 1
        
        lang_name = {"ko": "한국어", "en": "English", "ja": "日本語"}.get(language, "한국어")
        
        guide_prompt = (
            f"당신은 CS 센터 교육 전문가입니다. 다음 고객 응대 이력 데이터를 분석하여 종합적인 고객 응대 가이드라인을 작성하세요.\n\n"
            f"분석할 이력 데이터 (고객별 누적 데이터 포함):\n{json.dumps(list(customer_data_map.values())[:20], ensure_ascii=False, indent=2)}\n\n"
            f"다음 내용을 포함하여 가이드라인을 {lang_name}로 작성하세요:\n"
            f"1. 고객 유형별 응대 전략 (일반/까다로운/매우 불만족)\n"
            f"2. 문화권별 응대 가이드 (언어, 문화적 배경 고려)\n"
            f"3. 주요 문의 유형별 해결 방법\n"
            f"4. 고객 감정 점수에 따른 응대 전략\n"
            f"5. 개인정보 처리 가이드\n"
            f"6. 효과적인 소통 스타일 권장사항\n"
            f"7. 동일 고객의 반복 문의에 대한 대응 전략\n"
            f"8. 강성 고객 가이드라인 (까다로운 고객, 매우 불만족 고객)\n\n"
            f"가이드라인을 {lang_name}로 작성하세요. 실제 사례를 바탕으로 구체적이고 실용적인 내용으로 작성해주세요."
        )
        
        guide_content = run_llm(guide_prompt)
        
        if not guide_content or guide_content.startswith("❌"):
            return None
        
        today_str = datetime.now().strftime("%y%m%d")
        formatted_guide = (
            f"고객 응대 가이드라인\n"
            f"생성일: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
            f"분석 이력 수: {len(recent_histories)}\n"
            f"고객 수: {len(customer_data_map)}\n"
            f"=" * 80 + "\n\n"
            f"{guide_content}\n\n"
            f"=" * 80 + "\n"
            f"이 가이드는 AI 고객 응대 시뮬레이터 데이터를 기반으로 자동 생성되었습니다.\n"
            f"고객 데이터가 추가될 때마다 업데이트됩니다."
        )
        
        return formatted_guide
        
    except Exception as e:
        print(f"고객 가이드 생성 중 오류 발생: {e}")
        return None


def save_daily_customer_guide(guide_content: str, language: str = "ko") -> Optional[str]:
    """일일 고객 가이드를 파일로 저장합니다."""
    try:
        today_str = datetime.now().strftime("%y%m%d")
        guide_filename = f"{today_str}_고객가이드.TXT"
        guide_filepath = os.path.join(DATA_DIR, guide_filename)
        
        if os.path.exists(guide_filepath):
            with open(guide_filepath, "r", encoding="utf-8") as f:
                existing_content = f.read()
            
            updated_content = (
                f"{existing_content}\n\n"
                f"{'=' * 80}\n"
                f"업데이트: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
                f"{'=' * 80}\n\n"
                f"{guide_content}"
            )
            
            with open(guide_filepath, "w", encoding="utf-8") as f:
                f.write(updated_content)
        else:
            with open(guide_filepath, "w", encoding="utf-8") as f:
                f.write(guide_content)
        
        return guide_filepath
        
    except Exception as e:
        print(f"고객 가이드 저장 중 오류 발생: {e}")
        return None


# ========================================
# 이력 내보내기 (Word/PPTX/PDF)
# ========================================

def export_history_to_word(histories: List[Dict[str, Any]], filename: str = None, lang: str = "ko") -> str:
    """이력을 Word 파일로 저장"""
    if not IS_DOCX_AVAILABLE:
        raise ImportError("Word 저장을 위해 python-docx가 필요합니다: pip install python-docx")
    
    # 언어 설정 확인 및 기본값 설정
    if lang not in ["ko", "en", "ja"]:
        lang = "ko"
    L = LANG.get(lang, LANG["ko"])
    
    if filename is None:
        filename = f"customer_history_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
    filepath = os.path.join(DATA_DIR, filename)
    
    doc = Document()
    
    # 제목 추가
    title = doc.add_heading(L.get("download_history_title", "고객 응대 이력 요약"), 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # 각 이력 추가
    for i, hist in enumerate(histories, 1):
        # 이력 제목
        doc.add_heading(f'{L.get("download_history_number", "이력 #")}{i}', level=1)
        
        # 기본 정보
        doc.add_paragraph(f'ID: {hist.get("id", "N/A")}')
        doc.add_paragraph(f'{L.get("date_label", "날짜")}: {hist.get("timestamp", "N/A")}')
        doc.add_paragraph(f'{L.get("download_initial_inquiry", "초기 문의")}: {hist.get("initial_query", "N/A")}')
        doc.add_paragraph(f'{L.get("customer_type_label", "고객 유형")}: {hist.get("customer_type", "N/A")}')
        doc.add_paragraph(f'{L.get("language_label", "언어")}: {hist.get("language_key", "N/A")}')
        
        summary = hist.get('summary', {})
        if summary:
            # 요약 섹션
            doc.add_heading(L.get("download_summary", "요약"), level=2)
            doc.add_paragraph(f'{L.get("download_main_inquiry", "주요 문의")}: {summary.get("main_inquiry", "N/A")}')
            doc.add_paragraph(f'{L.get("download_key_response", "핵심 응답")}: {", ".join(summary.get("key_responses", []))}')
            doc.add_paragraph(f'{L.get("sentiment_score_label", "고객 감정 점수")}: {summary.get("customer_sentiment_score", "N/A")}/100')
            doc.add_paragraph(f'{L.get("customer_satisfaction_score_label", "고객 만족도 점수")}: {summary.get("customer_satisfaction_score", "N/A")}/100')
            
            # 고객 특성
            characteristics = summary.get('customer_characteristics', {})
            doc.add_heading(L.get("download_customer_characteristics", "고객 특성"), level=2)
            doc.add_paragraph(f'{L.get("language_label", "언어")}: {characteristics.get("language", "N/A")}')
            doc.add_paragraph(f'{L.get("download_cultural_background", "문화적 배경")}: {characteristics.get("cultural_hints", "N/A")}')
            doc.add_paragraph(f'{L.get("region_label", "지역")}: {characteristics.get("region", "N/A")}')
            doc.add_paragraph(f'{L.get("download_communication_style", "소통 스타일")}: {characteristics.get("communication_style", "N/A")}')
            
            # 개인정보 요약
            privacy = summary.get('privacy_info', {})
            doc.add_heading(L.get("download_privacy_summary", "개인정보 요약"), level=2)
            doc.add_paragraph(f'{L.get("email_provided_label", "이메일 제공")}: {L.get("download_yes", "예") if privacy.get("has_email") else L.get("download_no", "아니오")}')
            doc.add_paragraph(f'{L.get("phone_provided_label", "전화번호 제공")}: {L.get("download_yes", "예") if privacy.get("has_phone") else L.get("download_no", "아니오")}')
            doc.add_paragraph(f'{L.get("download_address_provided", "주소 제공")}: {L.get("download_yes", "예") if privacy.get("has_address") else L.get("download_no", "아니오")}')
            doc.add_paragraph(f'{L.get("download_region_hint", "지역 힌트")}: {privacy.get("region_hint", "N/A")}')
            
            # 전체 요약
            doc.add_paragraph(f'{L.get("download_overall_summary", "전체 요약")}: {summary.get("summary", "N/A")}')
        
        # 구분선
        if i < len(histories):
            doc.add_paragraph('-' * 80)
    
    doc.save(filepath)
    return filepath


def export_history_to_pptx(histories: List[Dict[str, Any]], filename: str = None, lang: str = "ko") -> str:
    """이력을 PPTX 파일로 저장"""
    if not IS_PPTX_AVAILABLE:
        raise ImportError("PPTX 저장을 위해 python-pptx가 필요합니다: pip install python-pptx")
    
    # 언어 설정 확인 및 기본값 설정
    if lang not in ["ko", "en", "ja"]:
        lang = "ko"
    L = LANG.get(lang, LANG["ko"])
    
    if filename is None:
        filename = f"customer_history_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    filepath = os.path.join(DATA_DIR, filename)
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 제목 슬라이드
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = L.get("download_history_title", "고객 응대 이력 요약")
    subtitle.text = f"{L.get('download_created_date', '생성일')}: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    # 각 이력에 대해 슬라이드 생성
    for i, hist in enumerate(histories, 1):
        # 제목 및 내용 레이아웃
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = f"{L.get('download_history_number', '이력 #')}{i}"
        
        tf = body_shape.text_frame
        tf.text = f"ID: {hist.get('id', 'N/A')}"
        
        p = tf.add_paragraph()
        p.text = f"{L.get('date_label', '날짜')}: {hist.get('timestamp', 'N/A')}"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = f"{L.get('download_initial_inquiry', '초기 문의')}: {hist.get('initial_query', 'N/A')}"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = f"{L.get('customer_type_label', '고객 유형')}: {hist.get('customer_type', 'N/A')}"
        p.level = 0
        
        summary = hist.get('summary', {})
        if summary:
            p = tf.add_paragraph()
            p.text = f"{L.get('download_main_inquiry', '주요 문의')}: {summary.get('main_inquiry', 'N/A')}"
            p.level = 0
            
            p = tf.add_paragraph()
            p.text = f"{L.get('sentiment_score_label', '고객 감정 점수')}: {summary.get('customer_sentiment_score', 'N/A')}/100"
            p.level = 0
            
            p = tf.add_paragraph()
            p.text = f"{L.get('customer_satisfaction_score_label', '고객 만족도 점수')}: {summary.get('customer_satisfaction_score', 'N/A')}/100"
            p.level = 0
    
    prs.save(filepath)
    return filepath


def export_history_to_pdf(histories: List[Dict[str, Any]], filename: str = None, lang: str = "ko") -> str:
    """이력을 PDF 파일로 저장 (한글/일본어 인코딩 지원 강화)"""
    if not IS_REPORTLAB_AVAILABLE:
        raise ImportError("PDF 저장을 위해 reportlab이 필요합니다: pip install reportlab")
    
    # 언어 설정 확인 및 기본값 설정
    if lang not in ["ko", "en", "ja"]:
        lang = "ko"
    L = LANG.get(lang, LANG["ko"])
    
    if filename is None:
        filename = f"customer_history_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
    filepath = os.path.join(DATA_DIR, filename)
    
    # 한글/일본어 폰트 지원
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    
    korean_font_registered = False
    japanese_font_registered = False
    korean_font_name = 'KoreanFont'
    japanese_font_name = 'JapaneseFont'
    
    def register_font(font_name: str, font_path: str) -> bool:
        """폰트를 등록하는 헬퍼 함수"""
        try:
            if font_path.endswith('.ttf'):
                font = TTFont(font_name, font_path)
                pdfmetrics.registerFont(font)
                return font_name in pdfmetrics.getRegisteredFontNames()
            elif font_path.endswith('.ttc'):
                for subfont_idx in range(8):
                    try:
                        font = TTFont(font_name, font_path, subfontIndex=subfont_idx)
                        pdfmetrics.registerFont(font)
                        if font_name in pdfmetrics.getRegisteredFontNames():
                            return True
                    except Exception:
                        continue
            elif font_path.endswith('.otf'):
                try:
                    font = TTFont(font_name, font_path)
                    pdfmetrics.registerFont(font)
                    return font_name in pdfmetrics.getRegisteredFontNames()
                except Exception:
                    pass
            return False
        except Exception:
            return False
    
    def download_font_if_needed(font_dir: str = None) -> str:
        """폰트가 없을 경우 Noto Sans CJK 폰트를 자동 다운로드"""
        if font_dir is None:
            font_dir = os.path.join(DATA_DIR, "fonts")
        os.makedirs(font_dir, exist_ok=True)
        
        font_files = [
            ("NotoSansCJK-Regular.ttf", "https://github.com/googlefonts/noto-cjk/raw/main/Sans/Subset/TTF/NotoSansCJK-Regular.ttf"),
            ("NotoSansCJK-Regular.otf", "https://github.com/googlefonts/noto-cjk/raw/main/Sans/Subset/OTF/NotoSansCJK-Regular.otf"),
        ]
        
        for font_filename, font_url in font_files:
            font_path = os.path.join(font_dir, font_filename)
            if os.path.exists(font_path) and os.path.getsize(font_path) > 1000:
                return font_path
        
        try:
            import requests
            for font_filename, font_url in font_files:
                font_path = os.path.join(font_dir, font_filename)
                try:
                    response = requests.get(font_url, timeout=30, stream=True)
                    response.raise_for_status()
                    with open(font_path, 'wb') as f:
                        for chunk in response.iter_content(chunk_size=8192):
                            f.write(chunk)
                    if os.path.exists(font_path) and os.path.getsize(font_path) > 1000:
                        return font_path
                except Exception:
                    continue
        except ImportError:
            pass
        return None
    
    try:
        import platform
        system = platform.system()
        
        if system == 'Windows':
            korean_font_paths = [
                "C:/Windows/Fonts/malgun.ttf", "C:/Windows/Fonts/malgunsl.ttf",
                "C:/Windows/Fonts/NanumGothic.ttf", "C:/Windows/Fonts/gulim.ttc",
                "C:/Windows/Fonts/batang.ttc", "C:/Windows/Fonts/malgun.ttc",
            ]
            japanese_font_paths = [
                "C:/Windows/Fonts/msgothic.ttc", "C:/Windows/Fonts/msmincho.ttc",
                "C:/Windows/Fonts/meiryo.ttc", "C:/Windows/Fonts/yuanti.ttc",
            ]
        elif system == 'Darwin':
            korean_font_paths = [
                "/System/Library/Fonts/Supplemental/AppleGothic.ttf",
                "/Library/Fonts/AppleGothic.ttf",
            ]
            japanese_font_paths = [
                "/System/Library/Fonts/Supplemental/AppleGothic.ttf",
                "/System/Library/Fonts/Hiragino Sans GB.ttc",
            ]
        else:
            korean_font_paths = [
                "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
                "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
            ]
            japanese_font_paths = [
                "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
                "/usr/share/fonts/truetype/takao/TakaoGothic.ttf",
            ]
        
        for font_path in korean_font_paths:
            if os.path.exists(font_path) and register_font(korean_font_name, font_path):
                korean_font_registered = True
                break
        
        for font_path in japanese_font_paths:
            if os.path.exists(font_path) and register_font(japanese_font_name, font_path):
                japanese_font_registered = True
                break
        
        if not korean_font_registered and not japanese_font_registered and system == 'Linux':
            downloaded_font = download_font_if_needed()
            if downloaded_font and os.path.exists(downloaded_font):
                if register_font(korean_font_name, downloaded_font):
                    korean_font_registered = True
                    japanese_font_registered = True
    except Exception:
        pass
    
    doc = SimpleDocTemplate(filepath, pagesize=letter)
    story = []
    styles = getSampleStyleSheet()
    
    def get_multilingual_style(base_style_name, default_font=None, **kwargs):
        """다국어 지원 스타일 생성"""
        base_style = styles[base_style_name]
        style_kwargs = {'parent': base_style, **kwargs}
        registered_fonts = pdfmetrics.getRegisteredFontNames()
        if default_font and default_font in registered_fonts:
            style_kwargs['fontName'] = default_font
        elif korean_font_registered and korean_font_name in registered_fonts:
            style_kwargs['fontName'] = korean_font_name
        elif japanese_font_registered and japanese_font_name in registered_fonts:
            style_kwargs['fontName'] = japanese_font_name
        return ParagraphStyle(f'Multilingual{base_style_name}', **style_kwargs)
    
    title_style = get_multilingual_style('Heading1', fontSize=24, textColor=black, spaceAfter=30, alignment=1,
                                         default_font=korean_font_name if korean_font_registered else japanese_font_name)
    normal_style = get_multilingual_style('Normal')
    heading1_style = get_multilingual_style('Heading1')
    heading2_style = get_multilingual_style('Heading2')
    
    def safe_text(text, detect_font=True):
        """텍스트를 안전하게 처리"""
        if text is None:
            return ("N/A", None)
        if isinstance(text, bytes):
            try:
                text_str = text.decode('utf-8', errors='replace')
            except:
                text_str = str(text)
        else:
            text_str = str(text)
        if text_str is None:
            return ("N/A", None)
        try:
            import unicodedata
            text_str = unicodedata.normalize('NFC', text_str)
        except:
            pass
        text_str = text_str.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        text_str = text_str.replace('"', '&quot;').replace("'", '&#39;')
        
        recommended_font = None
        if detect_font:
            try:
                has_korean = any('\uAC00' <= char <= '\uD7AF' or '\u1100' <= char <= '\u11FF' for char in text_str)
                has_japanese = any('\u3040' <= char <= '\u309F' or '\u30A0' <= char <= '\u30FF' or '\u4E00' <= char <= '\u9FFF' for char in text_str)
                registered_fonts = pdfmetrics.getRegisteredFontNames()
                if has_korean and korean_font_registered and korean_font_name in registered_fonts:
                    recommended_font = korean_font_name
                elif has_japanese and japanese_font_registered and japanese_font_name in registered_fonts:
                    recommended_font = japanese_font_name
            except Exception:
                pass
        return (text_str, recommended_font)
    
    def create_paragraph(text, style, auto_font=True):
        """Paragraph 생성"""
        text_str, recommended_font = safe_text(text, detect_font=auto_font)
        if recommended_font and auto_font:
            style_with_font = ParagraphStyle(name=f'{style.name}_with_font', parent=style, fontName=recommended_font)
            return Paragraph(text_str, style_with_font)
        return Paragraph(text_str, style)
    
    title_text, _ = safe_text(L.get("download_history_title", "고객 응대 이력 요약"))
    story.append(Paragraph(title_text, title_style))
    story.append(Spacer(1, 0.2*inch))
    
    for i, hist in enumerate(histories, 1):
        heading_text, _ = safe_text(f'{L.get("download_history_number", "이력 #")}{i}')
        story.append(Paragraph(heading_text, heading1_style))
        story.append(Spacer(1, 0.1*inch))
        
        id_text, _ = safe_text(f'ID: {hist.get("id", "N/A")}')
        story.append(create_paragraph(id_text, normal_style))
        
        timestamp_text, _ = safe_text(f'{L.get("date_label", "날짜")}: {hist.get("timestamp", "N/A")}')
        story.append(create_paragraph(timestamp_text, normal_style))
        
        query_text, _ = safe_text(f'{L.get("download_initial_inquiry", "초기 문의")}: {hist.get("initial_query", "N/A")}')
        story.append(create_paragraph(query_text, normal_style))
        
        customer_type_text, _ = safe_text(f'{L.get("customer_type_label", "고객 유형")}: {hist.get("customer_type", "N/A")}')
        story.append(create_paragraph(customer_type_text, normal_style))
        
        language_text, _ = safe_text(f'{L.get("language_label", "언어")}: {hist.get("language_key", "N/A")}')
        story.append(create_paragraph(language_text, normal_style))
        
        summary = hist.get('summary', {})
        if summary:
            story.append(Spacer(1, 0.1*inch))
            summary_title, _ = safe_text(L.get("download_summary", "요약"))
            story.append(Paragraph(summary_title, heading2_style))
            
            main_inquiry_text, _ = safe_text(f'{L.get("download_main_inquiry", "주요 문의")}: {summary.get("main_inquiry", "N/A")}')
            story.append(create_paragraph(main_inquiry_text, normal_style))
            
            key_responses = summary.get("key_responses", [])
            if isinstance(key_responses, list):
                responses_list = [safe_text(r)[0] for r in key_responses]
                responses_text = ", ".join(responses_list)
            else:
                responses_text, _ = safe_text(key_responses)
            responses_para_text, _ = safe_text(f'{L.get("download_key_response", "핵심 응답")}: {responses_text}')
            story.append(create_paragraph(responses_para_text, normal_style))
            
            sentiment_text, _ = safe_text(f'{L.get("sentiment_score_label", "고객 감정 점수")}: {summary.get("customer_sentiment_score", "N/A")}/100')
            story.append(create_paragraph(sentiment_text, normal_style))
            
            satisfaction_text, _ = safe_text(f'{L.get("customer_satisfaction_score_label", "고객 만족도 점수")}: {summary.get("customer_satisfaction_score", "N/A")}/100')
            story.append(create_paragraph(satisfaction_text, normal_style))
            
            characteristics = summary.get('customer_characteristics', {})
            story.append(Spacer(1, 0.1*inch))
            char_title, _ = safe_text(L.get("download_customer_characteristics", "고객 특성"))
            story.append(Paragraph(char_title, heading2_style))
            
            lang_char_text, _ = safe_text(f'{L.get("language_label", "언어")}: {characteristics.get("language", "N/A")}')
            story.append(create_paragraph(lang_char_text, normal_style))
            
            cultural_text, _ = safe_text(f'{L.get("download_cultural_background", "문화적 배경")}: {characteristics.get("cultural_hints", "N/A")}')
            story.append(create_paragraph(cultural_text, normal_style))
            
            region_text, _ = safe_text(f'{L.get("region_label", "지역")}: {characteristics.get("region", "N/A")}')
            story.append(create_paragraph(region_text, normal_style))
            
            comm_style_text, _ = safe_text(f'{L.get("download_communication_style", "소통 스타일")}: {characteristics.get("communication_style", "N/A")}')
            story.append(create_paragraph(comm_style_text, normal_style))
            
            privacy = summary.get('privacy_info', {})
            story.append(Spacer(1, 0.1*inch))
            privacy_title, _ = safe_text(L.get("download_privacy_summary", "개인정보 요약"))
            story.append(Paragraph(privacy_title, heading2_style))
            
            email_text, _ = safe_text(f'{L.get("email_provided_label", "이메일 제공")}: {L.get("download_yes", "예") if privacy.get("has_email") else L.get("download_no", "아니오")}')
            story.append(create_paragraph(email_text, normal_style))
            
            phone_text, _ = safe_text(f'{L.get("phone_provided_label", "전화번호 제공")}: {L.get("download_yes", "예") if privacy.get("has_phone") else L.get("download_no", "아니오")}')
            story.append(create_paragraph(phone_text, normal_style))
            
            address_text, _ = safe_text(f'{L.get("download_address_provided", "주소 제공")}: {L.get("download_yes", "예") if privacy.get("has_address") else L.get("download_no", "아니오")}')
            story.append(create_paragraph(address_text, normal_style))
            
            region_hint_text, _ = safe_text(f'{L.get("download_region_hint", "지역 힌트")}: {privacy.get("region_hint", "N/A")}')
            story.append(create_paragraph(region_hint_text, normal_style))
            
            full_summary_text, _ = safe_text(f'{L.get("download_overall_summary", "전체 요약")}: {summary.get("summary", "N/A")}')
            story.append(create_paragraph(full_summary_text, normal_style))
        
        if i < len(histories):
            story.append(Spacer(1, 0.2*inch))
            divider_text, _ = safe_text('-' * 80)
            story.append(Paragraph(divider_text, normal_style))
            story.append(Spacer(1, 0.2*inch))
    
    try:
        doc.build(story)
    except Exception as e:
        raise Exception(f"PDF 생성 실패: {str(e)}")
    
    return filepath


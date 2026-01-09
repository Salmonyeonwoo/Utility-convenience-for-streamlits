# ========================================
# utils/history_handler_exports.py
# 이력 관리 - 내보내기 함수들 (Word, PPTX, PDF)
# ========================================

import os
from datetime import datetime
from typing import List, Dict, Any
from lang_pack import LANG
from config import DATA_DIR

# Word, PPTX, PDF 내보내기 라이브러리
try:
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    IS_DOCX_AVAILABLE = True
except ImportError:
    IS_DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    IS_PPTX_AVAILABLE = True
except ImportError:
    IS_PPTX_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import black
    from reportlab.lib.units import inch
    IS_REPORTLAB_AVAILABLE = True
except ImportError:
    IS_REPORTLAB_AVAILABLE = False

def export_history_to_word(histories: List[Dict[str, Any]], filename: str = None, lang: str = "ko") -> str:
    """이력을 Word 파일로 저장"""
    if not IS_DOCX_AVAILABLE:
        raise ImportError("Word 저장을 위해 python-docx가 필요합니다: pip install python-docx")
    
    if lang not in ["ko", "en", "ja"]:
        lang = "ko"
    L = LANG.get(lang, LANG["ko"])
    
    if filename is None:
        filename = f"customer_history_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
    filepath = os.path.join(DATA_DIR, filename)
    
    doc = Document()
    
    # 제목 추가
    title = doc.add_heading(L.get("download_history_title", "고객 응대 이력 요약"), 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # 각 이력 추가
    for i, hist in enumerate(histories, 1):
        # 이력 제목
        doc.add_heading(f'{L.get("download_history_number", "이력 #")}{i}', level=1)
        
        # 기본 정보
        doc.add_paragraph(f'ID: {hist.get("id", "N/A")}')
        doc.add_paragraph(f'{L.get("date_label", "날짜")}: {hist.get("timestamp", "N/A")}')
        doc.add_paragraph(f'{L.get("download_initial_inquiry", "초기 문의")}: {hist.get("initial_query", "N/A")}')
        doc.add_paragraph(f'{L.get("customer_type_label", "고객 유형")}: {hist.get("customer_type", "N/A")}')
        doc.add_paragraph(f'{L.get("language_label", "언어")}: {hist.get("language_key", "N/A")}')
        
        summary = hist.get('summary', {})
        if summary:
            # 요약 섹션
            doc.add_heading(L.get("download_summary", "요약"), level=2)
            doc.add_paragraph(f'{L.get("download_main_inquiry", "주요 문의")}: {summary.get("main_inquiry", "N/A")}')
            doc.add_paragraph(f'{L.get("download_key_response", "핵심 응답")}: {", ".join(summary.get("key_responses", []))}')
            doc.add_paragraph(f'{L.get("sentiment_score_label", "고객 감정 점수")}: {summary.get("customer_sentiment_score", "N/A")}/100')
            doc.add_paragraph(f'{L.get("customer_satisfaction_score_label", "고객 만족도 점수")}: {summary.get("customer_satisfaction_score", "N/A")}/100')
            
            # 고객 특성
            characteristics = summary.get('customer_characteristics', {})
            doc.add_heading(L.get("download_customer_characteristics", "고객 특성"), level=2)
            doc.add_paragraph(f'{L.get("language_label", "언어")}: {characteristics.get("language", "N/A")}')
            doc.add_paragraph(f'{L.get("download_cultural_background", "문화적 배경")}: {characteristics.get("cultural_hints", "N/A")}')
            doc.add_paragraph(f'{L.get("region_label", "지역")}: {characteristics.get("region", "N/A")}')
            doc.add_paragraph(f'{L.get("download_communication_style", "소통 스타일")}: {characteristics.get("communication_style", "N/A")}')
            
            # 개인정보 요약
            privacy = summary.get('privacy_info', {})
            doc.add_heading(L.get("download_privacy_summary", "개인정보 요약"), level=2)
            doc.add_paragraph(f'{L.get("email_provided_label", "이메일 제공")}: {L.get("download_yes", "예") if privacy.get("has_email") else L.get("download_no", "아니오")}')
            doc.add_paragraph(f'{L.get("phone_provided_label", "전화번호 제공")}: {L.get("download_yes", "예") if privacy.get("has_phone") else L.get("download_no", "아니오")}')
            doc.add_paragraph(f'{L.get("download_address_provided", "주소 제공")}: {L.get("download_yes", "예") if privacy.get("has_address") else L.get("download_no", "아니오")}')
            doc.add_paragraph(f'{L.get("download_region_hint", "지역 힌트")}: {privacy.get("region_hint", "N/A")}')
            
            # 전체 요약
            doc.add_paragraph(f'{L.get("download_overall_summary", "전체 요약")}: {summary.get("summary", "N/A")}')
        
        # 구분선
        if i < len(histories):
            doc.add_paragraph('-' * 80)
    
    doc.save(filepath)
    return filepath

def export_history_to_pptx(histories: List[Dict[str, Any]], filename: str = None, lang: str = "ko") -> str:
    """이력을 PPTX 파일로 저장"""
    if not IS_PPTX_AVAILABLE:
        raise ImportError("PPTX 저장을 위해 python-pptx가 필요합니다: pip install python-pptx")
    
    if lang not in ["ko", "en", "ja"]:
        lang = "ko"
    L = LANG.get(lang, LANG["ko"])
    
    if filename is None:
        filename = f"customer_history_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    filepath = os.path.join(DATA_DIR, filename)
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 제목 슬라이드
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = L.get("download_history_title", "고객 응대 이력 요약")
    subtitle.text = f"{L.get('download_created_date', '생성일')}: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    # 각 이력에 대해 슬라이드 생성
    for i, hist in enumerate(histories, 1):
        # 제목 및 내용 레이아웃
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = f"{L.get('download_history_number', '이력 #')}{i}"
        
        tf = body_shape.text_frame
        tf.text = f"ID: {hist.get('id', 'N/A')}"
        
        p = tf.add_paragraph()
        p.text = f"{L.get('date_label', '날짜')}: {hist.get('timestamp', 'N/A')}"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = f"{L.get('download_initial_inquiry', '초기 문의')}: {hist.get('initial_query', 'N/A')}"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = f"{L.get('customer_type_label', '고객 유형')}: {hist.get('customer_type', 'N/A')}"
        p.level = 0
        
        summary = hist.get('summary', {})
        if summary:
            p = tf.add_paragraph()
            p.text = f"{L.get('download_main_inquiry', '주요 문의')}: {summary.get('main_inquiry', 'N/A')}"
            p.level = 0
            
            p = tf.add_paragraph()
            p.text = f"{L.get('sentiment_score_label', '고객 감정 점수')}: {summary.get('customer_sentiment_score', 'N/A')}/100"
            p.level = 0
            
            p = tf.add_paragraph()
            p.text = f"{L.get('customer_satisfaction_score_label', '고객 만족도 점수')}: {summary.get('customer_satisfaction_score', 'N/A')}/100"
            p.level = 0
    
    prs.save(filepath)
    return filepath

# PDF export는 별도 파일로 분리 (너무 큼)
from utils.history_handler_exports_pdf import export_history_to_pdf


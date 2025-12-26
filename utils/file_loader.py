# ========================================
# utils/file_loader.py
# 파일 자동 인식 및 로드 모듈
# ========================================

import os
import json
import csv
import io
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path

# Word, PPTX, PDF 파싱 라이브러리
try:
    from docx import Document
    IS_DOCX_AVAILABLE = True
except ImportError:
    IS_DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    IS_PPTX_AVAILABLE = True
except ImportError:
    IS_PPTX_AVAILABLE = False

try:
    import PyPDF2
    IS_PDF_AVAILABLE = True
except ImportError:
    try:
        import pdfplumber
        IS_PDF_AVAILABLE = True
        USE_PDFPLUMBER = True
    except ImportError:
        IS_PDF_AVAILABLE = False
        USE_PDFPLUMBER = False


def scan_data_directory(data_dir: str = None) -> List[Dict[str, Any]]:
    """
    데이터 디렉토리에서 지원되는 파일들을 스캔하여 메타데이터 반환
    
    Args:
        data_dir: 스캔할 디렉토리 경로 (None이면 기본 경로 사용)
    
    Returns:
        파일 메타데이터 리스트
    """
    if data_dir is None:
        # 기본 경로: 프로젝트 루트의 data 폴더
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        data_dir = os.path.join(base_dir, "data")
    
    files_metadata = []
    
    # 로컬 data 디렉토리 스캔
    if os.path.exists(data_dir):
        for root, dirs, files in os.walk(data_dir):
            # .git 폴더는 제외
            if '.git' in root:
                continue
                
            for file in files:
                file_path = os.path.join(root, file)
                file_ext = os.path.splitext(file)[1].lower()
                
                if file_ext in ['.json', '.docx', '.pdf', '.pptx', '.csv']:
                    try:
                        file_stat = os.stat(file_path)
                        # 상대 경로 계산
                        rel_path = os.path.relpath(file_path, data_dir)
                        
                        metadata = {
                            "file_path": file_path,
                            "file_name": file,
                            "file_type": file_ext[1:],  # .json -> json
                            "file_size": file_stat.st_size,
                            "modified_time": datetime.fromtimestamp(file_stat.st_mtime).isoformat(),
                            "source": "local",
                            "relative_path": rel_path
                        }
                        files_metadata.append(metadata)
                    except Exception as e:
                        print(f"파일 메타데이터 읽기 오류: {file_path}, {e}")
    
    # GitHub customers 디렉토리도 스캔 (로컬 클론된 경우)
    github_data_dir = os.path.join(data_dir, "customers")
    if os.path.exists(github_data_dir):
        for root, dirs, files in os.walk(github_data_dir):
            for file in files:
                file_path = os.path.join(root, file)
                file_ext = os.path.splitext(file)[1].lower()
                
                if file_ext in ['.json', '.docx', '.pdf', '.pptx', '.csv']:
                    try:
                        file_stat = os.stat(file_path)
                        rel_path = os.path.relpath(file_path, data_dir)
                        
                        metadata = {
                            "file_path": file_path,
                            "file_name": file,
                            "file_type": file_ext[1:],
                            "file_size": file_stat.st_size,
                            "modified_time": datetime.fromtimestamp(file_stat.st_mtime).isoformat(),
                            "source": "github",
                            "relative_path": rel_path
                        }
                        files_metadata.append(metadata)
                    except Exception as e:
                        print(f"파일 메타데이터 읽기 오류: {file_path}, {e}")
    
    # 수정 시간 기준으로 정렬 (최신순)
    files_metadata.sort(key=lambda x: x.get("modified_time", ""), reverse=True)
    
    return files_metadata


def scan_github_repository(repo_path: str = "Salmonyeonwoo/Utility-convenience-for-streamlits", 
                           data_path: str = "data/customers",
                           github_token: str = None) -> List[Dict[str, Any]]:
    """
    GitHub 저장소에서 파일 목록 가져오기 (선택적 기능)
    
    Args:
        repo_path: GitHub 저장소 경로 (owner/repo 형식)
        data_path: 저장소 내 데이터 경로
        github_token: GitHub Personal Access Token (선택적, rate limit 증가용)
    
    Returns:
        파일 메타데이터 리스트 (GitHub API 사용 시)
    """
    files_metadata = []
    
    try:
        import requests
        
        # GitHub API를 사용하여 파일 목록 가져오기
        api_url = f"https://api.github.com/repos/{repo_path}/contents/{data_path}"
        
        # 헤더 설정 (토큰이 있으면 사용)
        headers = {}
        if github_token:
            headers["Authorization"] = f"token {github_token}"
        
        response = requests.get(api_url, headers=headers, timeout=10)
        
        if response.status_code == 200:
            contents = response.json()
            for item in contents:
                if item.get("type") == "file":
                    file_name = item.get("name", "")
                    file_ext = os.path.splitext(file_name)[1].lower()
                    
                    if file_ext in ['.json', '.docx', '.pdf', '.pptx', '.csv']:
                        metadata = {
                            "file_path": item.get("download_url", ""),
                            "file_name": file_name,
                            "file_type": file_ext[1:],
                            "file_size": item.get("size", 0),
                            "modified_time": item.get("updated_at", datetime.now().isoformat()),
                            "source": "github_api",
                            "download_url": item.get("download_url", ""),
                            "sha": item.get("sha", ""),
                            "github_token": github_token  # 다운로드 시 사용
                        }
                        files_metadata.append(metadata)
        elif response.status_code == 404:
            # 경로가 없으면 무시 (선택적 기능)
            pass
        elif response.status_code == 401:
            print("GitHub API 인증 실패: 토큰이 유효하지 않거나 권한이 없습니다.")
        elif response.status_code == 403:
            print("GitHub API rate limit 초과: 토큰을 사용하면 rate limit이 증가합니다.")
    except ImportError:
        # requests 라이브러리가 없으면 GitHub API 기능 비활성화
        pass
    except Exception as e:
        # GitHub API 오류는 무시 (선택적 기능)
        print(f"GitHub API 오류: {e}")
    
    return files_metadata


def load_json_file(file_path: str) -> Optional[Dict[str, Any]]:
    """JSON 파일 로드"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
            return data
    except Exception as e:
        print(f"JSON 파일 로드 오류: {file_path}, {e}")
        return None


def load_word_file(file_path: str) -> Optional[Dict[str, Any]]:
    """Word 파일에서 텍스트 추출 및 파싱"""
    if not IS_DOCX_AVAILABLE:
        return None
    
    try:
        doc = Document(file_path)
        text_content = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)
        
        # 테이블에서도 텍스트 추출
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text)
                if row_text:
                    text_content.append(" | ".join(row_text))
        
        full_text = "\n".join(text_content)
        
        # JSON 형식으로 파싱 시도
        try:
            # JSON이 포함되어 있는지 확인
            if "{" in full_text and "}" in full_text:
                json_start = full_text.find("{")
                json_end = full_text.rfind("}") + 1
                json_str = full_text[json_start:json_end]
                data = json.loads(json_str)
                return data
        except:
            pass
        
        # JSON 파싱 실패 시 텍스트로 반환
        return {
            "content": full_text,
            "file_type": "word",
            "parsed": False
        }
    except Exception as e:
        print(f"Word 파일 로드 오류: {file_path}, {e}")
        return None


def load_pdf_file(file_path: str) -> Optional[Dict[str, Any]]:
    """PDF 파일에서 텍스트 추출 및 파싱"""
    if not IS_PDF_AVAILABLE:
        return None
    
    try:
        text_content = []
        
        if USE_PDFPLUMBER:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        text_content.append(text)
        else:
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text:
                        text_content.append(text)
        
        full_text = "\n".join(text_content)
        
        # JSON 형식으로 파싱 시도
        try:
            if "{" in full_text and "}" in full_text:
                json_start = full_text.find("{")
                json_end = full_text.rfind("}") + 1
                json_str = full_text[json_start:json_end]
                data = json.loads(json_str)
                return data
        except:
            pass
        
        # JSON 파싱 실패 시 텍스트로 반환
        return {
            "content": full_text,
            "file_type": "pdf",
            "parsed": False
        }
    except Exception as e:
        print(f"PDF 파일 로드 오류: {file_path}, {e}")
        return None


def load_pptx_file(file_path: str) -> Optional[Dict[str, Any]]:
    """PPTX 파일에서 텍스트 추출 및 파싱"""
    if not IS_PPTX_AVAILABLE:
        return None
    
    try:
        prs = Presentation(file_path)
        text_content = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)
        
        full_text = "\n".join(text_content)
        
        # JSON 형식으로 파싱 시도
        try:
            if "{" in full_text and "}" in full_text:
                json_start = full_text.find("{")
                json_end = full_text.rfind("}") + 1
                json_str = full_text[json_start:json_end]
                data = json.loads(json_str)
                return data
        except:
            pass
        
        # JSON 파싱 실패 시 텍스트로 반환
        return {
            "content": full_text,
            "file_type": "pptx",
            "parsed": False
        }
    except Exception as e:
        print(f"PPTX 파일 로드 오류: {file_path}, {e}")
        return None


def load_csv_file(file_path: str) -> Optional[Dict[str, Any]]:
    """CSV 파일 로드 및 JSON 형식으로 변환"""
    try:
        data = []
        with open(file_path, 'r', encoding='utf-8-sig') as f:
            reader = csv.DictReader(f)
            for row in reader:
                data.append(row)
        
        return {
            "data": data,
            "file_type": "csv",
            "parsed": True
        }
    except Exception as e:
        print(f"CSV 파일 로드 오류: {file_path}, {e}")
        return None


def download_file_from_url(url: str, save_path: str = None, github_token: str = None) -> Optional[str]:
    """
    URL에서 파일 다운로드 (GitHub API용)
    
    Args:
        url: 다운로드할 파일 URL
        save_path: 저장할 경로 (None이면 임시 파일 생성)
        github_token: GitHub Personal Access Token (선택적)
    
    Returns:
        저장된 파일 경로 또는 None
    """
    try:
        import requests
        
        # 헤더 설정 (토큰이 있으면 사용)
        headers = {}
        if github_token:
            headers["Authorization"] = f"token {github_token}"
        
        response = requests.get(url, headers=headers, timeout=10)
        if response.status_code == 200:
            if save_path is None:
                import tempfile
                _, ext = os.path.splitext(url.split('?')[0])  # 쿼리 파라미터 제거
                fd, save_path = tempfile.mkstemp(suffix=ext)
                os.close(fd)
            
            with open(save_path, 'wb') as f:
                f.write(response.content)
            
            return save_path
        else:
            print(f"파일 다운로드 실패: HTTP {response.status_code}")
    except ImportError:
        print("requests 라이브러리가 필요합니다: pip install requests")
    except Exception as e:
        print(f"파일 다운로드 오류: {url}, {e}")
    
    return None


def load_file_by_type(file_path: str, file_type: str, github_token: str = None) -> Optional[Dict[str, Any]]:
    """
    파일 타입에 따라 적절한 로더 함수 호출
    
    Args:
        file_path: 파일 경로 또는 URL
        file_type: 파일 확장자 (json, docx, pdf, pptx, csv)
        github_token: GitHub Personal Access Token (URL 다운로드 시 사용)
    
    Returns:
        파싱된 데이터 또는 None
    """
    file_type = file_type.lower()
    
    # URL인 경우 다운로드
    if file_path.startswith("http://") or file_path.startswith("https://"):
        downloaded_path = download_file_from_url(file_path, github_token=github_token)
        if downloaded_path:
            file_path = downloaded_path
        else:
            return None
    
    if file_type == "json":
        return load_json_file(file_path)
    elif file_type == "docx":
        return load_word_file(file_path)
    elif file_type == "pdf":
        return load_pdf_file(file_path)
    elif file_type == "pptx":
        return load_pptx_file(file_path)
    elif file_type == "csv":
        return load_csv_file(file_path)
    else:
        return None


def parse_history_from_file_data(file_data: Dict[str, Any], file_name: str) -> Optional[Dict[str, Any]]:
    """
    파일 데이터를 시뮬레이션 이력 형식으로 변환
    
    Args:
        file_data: 파일에서 로드한 데이터
        file_name: 원본 파일명
    
    Returns:
        시뮬레이션 이력 형식의 딕셔너리 또는 None
    """
    try:
        # 이미 시뮬레이션 이력 형식인 경우
        if isinstance(file_data, list) and len(file_data) > 0:
            # 리스트의 첫 번째 항목 반환
            history = file_data[0]
            if "messages" in history or "initial_query" in history:
                return history
        
        # 딕셔너리인 경우
        if isinstance(file_data, dict):
            # messages 필드가 있으면 그대로 사용
            if "messages" in file_data and isinstance(file_data.get("messages"), list):
                return file_data
            
            # initial_query가 있으면 시뮬레이션 이력 형식으로 변환
            if "initial_query" in file_data:
                return file_data
            
            # basic_info가 있는 경우 (고객 데이터 형식)
            if "basic_info" in file_data:
                basic_info = file_data.get("basic_info", {})
                customer_data = file_data.get("data", {})
                
                # 상담 이력에서 최근 문의 찾기
                consultation_history = customer_data.get("consultation_history", [])
                initial_query = ""
                customer_type = "일반"
                
                if consultation_history and len(consultation_history) > 0:
                    # 최근 상담 이력에서 문의 내용 추출
                    latest_consultation = consultation_history[-1]
                    initial_query = (
                        latest_consultation.get("inquiry", "") or 
                        latest_consultation.get("content", "") or
                        latest_consultation.get("message", "") or
                        f"{basic_info.get('customer_name', '고객')}님의 상담 이력"
                    )
                    customer_type = latest_consultation.get("customer_type", "일반")
                else:
                    # 상담 이력이 없으면 기본 정보로 초기 문의 생성
                    customer_name = basic_info.get("customer_name", "고객")
                    initial_query = f"{customer_name}님의 고객 정보를 불러왔습니다."
                
                # CRM 프로필에서 고객 유형 가져오기
                crm_profile = customer_data.get("crm_profile", {})
                if crm_profile and crm_profile.get("personality"):
                    customer_type = crm_profile.get("personality", customer_type)
                
                # 메시지 생성
                messages = []
                if initial_query:
                    messages.append({
                        "role": "customer",
                        "content": initial_query
                    })
                
                return {
                    "initial_query": initial_query,
                    "customer_type": customer_type,
                    "messages": messages,
                    "summary": None,
                    "timestamp": datetime.now().isoformat(),
                    "source_file": file_name,
                    "customer_id": basic_info.get("customer_id", ""),
                    "customer_name": basic_info.get("customer_name", "")
                }
            
            # data 필드가 있는 경우 (customer_data 형식 - 다른 구조)
            if "data" in file_data:
                customer_info = file_data.get("data", {})
                
                # 여러 경로에서 initial_query 찾기
                initial_query = (
                    customer_info.get("last_inquiry", "") or 
                    customer_info.get("initial_query", "") or
                    customer_info.get("recent_inquiry", "") or
                    customer_info.get("inquiry", "")
                )
                
                # 상담 이력에서 찾기
                if not initial_query:
                    consultation_history = customer_info.get("consultation_history", [])
                    if consultation_history and len(consultation_history) > 0:
                        latest = consultation_history[-1]
                        initial_query = (
                            latest.get("inquiry", "") or 
                            latest.get("content", "") or
                            latest.get("message", "")
                        )
                
                # 기본 문의 생성
                if not initial_query:
                    customer_name = customer_info.get("name", "") or file_name.replace(".json", "").replace("CUST-", "")
                    initial_query = f"{customer_name}님의 고객 정보를 불러왔습니다."
                
                customer_type = "일반"
                crm_profile = customer_info.get("crm_profile", {})
                if crm_profile:
                    customer_type = crm_profile.get("personality", customer_type)
                
                return {
                    "initial_query": initial_query,
                    "customer_type": customer_type,
                    "messages": [{"role": "customer", "content": initial_query}] if initial_query else [],
                    "summary": None,
                    "timestamp": datetime.now().isoformat(),
                    "source_file": file_name
                }
            
            # 직접 필드가 있는 경우 (다양한 형식 지원)
            if "customer_name" in file_data or "customer_id" in file_data:
                customer_name = file_data.get("customer_name", file_data.get("name", "고객"))
                initial_query = (
                    file_data.get("inquiry", "") or 
                    file_data.get("last_inquiry", "") or
                    f"{customer_name}님의 고객 정보를 불러왔습니다."
                )
                
                return {
                    "initial_query": initial_query,
                    "customer_type": file_data.get("customer_type", file_data.get("personality", "일반")),
                    "messages": [{"role": "customer", "content": initial_query}],
                    "summary": None,
                    "timestamp": datetime.now().isoformat(),
                    "source_file": file_name
                }
        
        # 텍스트 콘텐츠만 있는 경우
        if "content" in file_data:
            content = file_data["content"]
            # JSON 형식으로 파싱 시도
            try:
                if isinstance(content, str) and "{" in content:
                    json_start = content.find("{")
                    json_end = content.rfind("}") + 1
                    json_str = content[json_start:json_end]
                    parsed_data = json.loads(json_str)
                    return parse_history_from_file_data(parsed_data, file_name)
            except:
                pass
            
            # 파싱 실패 시 텍스트를 initial_query로 사용
            return {
                "initial_query": content[:500] if isinstance(content, str) else str(content)[:500],
                "customer_type": "일반",
                "messages": [{"role": "customer", "content": content[:500] if isinstance(content, str) else str(content)[:500]}],
                "summary": None,
                "timestamp": datetime.now().isoformat(),
                "source_file": file_name
            }
        
        # 모든 파싱 실패 시 기본 이력 생성 (파일명 기반)
        if file_name:
            customer_id = file_name.replace(".json", "").replace("CUST-", "").replace(".docx", "").replace(".pdf", "").replace(".pptx", "").replace(".csv", "")
            return {
                "initial_query": f"고객 정보 파일을 불러왔습니다. (파일: {file_name})",
                "customer_type": "일반",
                "messages": [{"role": "customer", "content": f"고객 정보 파일을 불러왔습니다. (파일: {file_name})"}],
                "summary": None,
                "timestamp": datetime.now().isoformat(),
                "source_file": file_name,
                "raw_data": file_data  # 원본 데이터 보관
            }
        
        return None
    except Exception as e:
        print(f"이력 파싱 오류: {file_name}, {e}")
        import traceback
        traceback.print_exc()
        # 오류 발생 시에도 기본 이력 생성
        return {
            "initial_query": f"파일을 불러왔습니다: {file_name}",
            "customer_type": "일반",
            "messages": [{"role": "customer", "content": f"파일을 불러왔습니다: {file_name}"}],
            "summary": None,
            "timestamp": datetime.now().isoformat(),
            "source_file": file_name,
            "parse_error": str(e)
        }

